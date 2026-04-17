# 顧客提案資料向けに名称・文言を改善した版
from __future__ import annotations

# =========================================================
# 0. Imports / Page Config
# =========================================================
from dataclasses import dataclass, field
from datetime import date, timedelta
from decimal import Decimal, ROUND_HALF_UP
from io import BytesIO
from numbers import Number
import math
import os
import re
import unicodedata
from typing import Any, Optional


@dataclass
class CalculationContext:
    """
    計算プロセス全体で共有されるコンテキスト情報を保持するデータクラス。
    一次相続・二次相続の入力データおよび計算結果を統合管理します。
    """
    primary_inputs: dict = field(default_factory=dict)
    secondary_inputs: Optional[dict] = field(default_factory=dict)
    primary_result: Optional[dict] = field(default_factory=dict)
    common_config: dict = field(default_factory=dict)

    # 拡張性を考慮し、動的な属性追加を許可
    def __post_init__(self):
        if self.primary_inputs is None:
            self.primary_inputs = {}
        if self.secondary_inputs is None:
            self.secondary_inputs = {}


import pandas as pd
import plotly.graph_objects as go
import streamlit as st
from openpyxl import load_workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.page import PageMargins
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.pdfmetrics import registerFont
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.shapes import Drawing, Line, String
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

st.set_page_config(page_title="SUMMIT v31.16 PRO", layout="wide")

# =========================================================
# 1. Constants
# =========================================================
APP_TITLE = "山根会計 相続シミュレーション"
APP_LOGIN_USER_LABEL = "ログイン: 川東"
APP_PASSWORD_ENV_KEY = "SUMMIT_APP_PASSWORD"
EXCEL_FILE_NAME = "相続シミュレーション_ご提案資料.xlsx"
EXCEL_TITLE = "山根会計 相続税シミュレーション ご提案資料"
PDF_FILE_NAME = "相続シミュレーション_ご提案資料.pdf"
PPT_FILE_NAME = "相続シミュレーション_ご提案資料.pptx"
GLOBAL_RISK_NOTICE = "本資料は現時点でご提供資料・条件に基づくシミュレーションです。正式申告・実行時には追加確認により変動する場合があります。"
SECONDARY_RISK_NOTICE = "二次相続は概算ロジックを含む参考表示です。相次相続控除・相続人構成・個別事情の確認前に断定利用しないでください。"
SMALL_SCALE_RISK_NOTICE = "小規模宅地等は概算判定です。適用可・要確認の表示にかかわらず、実務利用前に要件を別途確認してください。"
INSURANCE_GIFT_RISK_NOTICE = "生命保険・贈与加算・精算課税は入力内容に依存する概算整理です。証憑確認前の断定利用は禁止です。"
OUTPUT_RISK_NOTICE = "本資料は現時点でご提示いただいた資料・条件に基づく試算です。正式申告時には評価資料、分割内容、各種特例の適用可否に応じて金額が変動する場合があります。"

COLOR_NAVY = "1f2c4d"
COLOR_GOLD = "c5a059"
COLOR_RED = "a61d24"
COLOR_BLUE = "1E5AA8"
COLOR_SKY = "DCEAFB"
COLOR_LIGHT_BLUE = "EDF5FF"
COLOR_LIGHT_RED = "FDECEC"
COLOR_GRAY = "666666"
COLOR_LIGHT_GRAY = "E9EDF3"
COLOR_DARK = "22324A"
COLOR_GREEN = "3D7A57"
COLOR_TEXT = "333333"

SMALL_SCALE_HOME_LIMIT = Decimal("330")
SMALL_SCALE_BUSINESS_LIMIT = Decimal("400")
SMALL_SCALE_RENT_LIMIT = Decimal("200")
SMALL_SCALE_HOME_RATE = Decimal("0.8")
SMALL_SCALE_BUSINESS_RATE = Decimal("0.8")
SMALL_SCALE_RENT_RATE = Decimal("0.5")

SPOUSE_TAX_REDUCTION_REFERENCE = Decimal("160000000")
LIFE_INSURANCE_EXEMPT_PER_HEIR = Decimal("5000000")
BASIC_DEDUCTION_BASE = Decimal("30000000")
BASIC_DEDUCTION_PER_HEIR = Decimal("6000000")
PERCENT_DENOMINATOR = Decimal("100")
TWO_TENTHS_SURTAX_RATE = Decimal("0.2")
MAX_INSURANCE_RECIPIENT_ROWS = 5
MAX_GIFT_RECORD_ROWS = 10
ANNUAL_GIFT_LOOKBACK_YEARS = 7  # 概算。制度経過措置の精密判定は今後拡張
SEISAN_ANNUAL_BASIC_EXEMPTION = Decimal("1100000")
GIFT_TYPE_ANNUAL = "暦年課税"
GIFT_TYPE_SEISAN = "相続時精算課税"
GIFT_TYPE_OPTIONS = [GIFT_TYPE_ANNUAL, GIFT_TYPE_SEISAN]

SMALL_SCALE_STATUS_APPLICABLE = "適用可"
SMALL_SCALE_STATUS_REVIEW = "要確認"
SMALL_SCALE_STATUS_NOT_APPLICABLE = "適用不可"

LAND_CATEGORY_HOME = "特定居住用"
LAND_CATEGORY_BUSINESS = "特定事業用"
LAND_CATEGORY_RENTAL = "貸付事業用"

HEIR_TYPE_CHILD = "子"
HEIR_TYPE_GRANDCHILD = "孫（養子含む）"
HEIR_TYPE_PARENT = "親"
HEIR_TYPE_FULL_SIBLING = "兄弟姉妹（全血）"
HEIR_TYPE_HALF_SIBLING = "兄弟姉妹（半血）"

HEIR_TYPE_OPTIONS = [
    HEIR_TYPE_CHILD,
    HEIR_TYPE_GRANDCHILD,
    HEIR_TYPE_PARENT,
    HEIR_TYPE_FULL_SIBLING,
    HEIR_TYPE_HALF_SIBLING,
]

TAB_LABELS = [
    " 👥  1.基本構成",
    " 💰  2.一次財産詳細",
    " 📑  3.一次相続明細（概算）",
    " 📑  4.二次相続明細（概算参考）",
    " ⏳  5.二次推移予測（参考）",
    " 📊  6.分析結果（内部確認用）",
]

# =========================================================
# 2. Data Models
# =========================================================


@dataclass
class PrimaryInputs:
    heir_count: int
    has_spouse: bool
    heirs_info: list[dict[str, str]]
    date_of_death: date
    v_home: int
    a_home: int
    v_biz: int
    a_biz: int
    v_rent: int
    a_rent: int
    small_scale_inputs: dict[str, SmallScaleInput]
    v_build: int
    v_stock: int
    v_cash: int
    v_ins: int
    insurance_entries: list[InsuranceRecipientInput]
    v_others: int
    v_debt: int
    v_funeral: int
    gift_records: list[GiftRecord]


@dataclass
class SecondaryInputs:
    spouse_acquisition_pct: int
    s_own: int
    annual_spend: int
    interval_years: int
    use_individual_allocations: bool
    actual_acquisition_inputs: list[int]


@dataclass
class InsuranceRecipientInput:
    recipient_name: str
    recipient_type: str
    amount: int
    is_statutory_heir: bool
    is_two_tenths_target: bool


@dataclass
class GiftRecord:
    gift_date: date
    recipient_name: str
    recipient_type: str
    amount: int
    tax_type: str


@dataclass
class GiftComputationRecord:
    gift_date: date
    recipient_name: str
    recipient_type: str
    tax_type: str
    amount: Decimal
    calendar_year: int
    is_addback_target: bool
    addback_amount: Decimal
    reason: str


@dataclass
class SmallScaleInput:
    category: str
    acquirer_name: str
    apply_special_rule: bool
    is_spouse_acquirer: bool = False
    is_cohabiting_heir: bool = False
    is_no_house_heir: bool = False
    will_hold_until_filing: bool = False
    will_reside_until_filing: bool = False
    is_business_successor: bool = False
    will_continue_business: bool = False
    will_continue_rental: bool = False


@dataclass
class SmallScaleLandRecord:
    category: str
    land_name: str
    acquirer_name: str
    status: str
    reason: str
    original_value: Decimal
    area_sqm: Decimal
    eligible_area_sqm: Decimal
    reduction_rate: Decimal
    reduction_amount: Decimal


@dataclass
class HeirTaxRecord:
    name: str
    heir_type: str
    legal_share: Decimal
    actual_share: Decimal
    input_acquisition_amount: Decimal
    normalized_acquisition_amount: Decimal
    insurance_gross: Decimal
    insurance_nontaxable: Decimal
    insurance_taxable: Decimal
    annual_gift_addback: Decimal
    seisan_gift_addback: Decimal
    base_taxable_price: Decimal
    taxable_price: Decimal
    preliminary_tax: Decimal
    two_tenths_surtax: Decimal
    spouse_credit: Decimal
    final_tax: Decimal
    is_two_tenths_target: bool


@dataclass
class PrimaryResult:
    st_count: int
    land_eval: Decimal
    total_red: Decimal
    small_scale_records: list[SmallScaleLandRecord]
    ins_ded: Decimal
    pure_as: Decimal
    tax_p: Decimal
    basic_1: Decimal
    taxable_1: Decimal
    total_tax_1: Decimal
    spouse_legal_share: Decimal
    heir_legal_shares: list[Decimal]
    spouse_actual_share: Decimal
    spouse_actual_taxable_price: Decimal
    spouse_tax_limit: Decimal
    total_final_tax: Decimal
    total_insurance_gross: Decimal
    total_insurance_nontaxable: Decimal
    total_insurance_taxable: Decimal
    total_annual_gift_addback: Decimal
    total_seisan_gift_addback: Decimal
    gift_detail_records: list[GiftComputationRecord]
    heir_tax_records: list[HeirTaxRecord]


@dataclass
class HeirCarryForwardSnapshot:
    heir_id: str
    heir_name: str
    relation_type: str
    birth_date: Optional[date]
    disability_flag: bool
    acquired_total_amount: Decimal
    special_disability_flag: bool = False
    acquired_cash_amount: Decimal = Decimal("0")
    acquired_securities_amount: Decimal = Decimal("0")
    acquired_real_estate_amount: Decimal = Decimal("0")
    acquired_insurance_amount: Decimal = Decimal("0")
    acquired_other_amount: Decimal = Decimal("0")
    paid_inheritance_tax_amount: Decimal = Decimal("0")
    net_assets_after_first_tax: Decimal = Decimal("0")
    real_estate_usage_type: str = ""
    cohabitation_flag: bool = False
    business_use_flag: bool = False
    notes: list[str] = field(default_factory=list)


@dataclass
class PrimaryToSecondarySnapshot:
    first_inheritance_date: date
    inheritance_case_id: str
    division_status: str
    assumption_notes: list[str] = field(default_factory=list)
    first_total_estate_amount: Decimal = Decimal("0")
    first_total_taxable_base: Decimal = Decimal("0")
    first_total_tax: Decimal = Decimal("0")
    first_total_net_assets_after_tax: Decimal = Decimal("0")
    spouse_heir_id: str = ""
    spouse_acquired_total_amount: Decimal = Decimal("0")
    spouse_net_assets_after_first_tax: Decimal = Decimal("0")
    heir_snapshots: list[HeirCarryForwardSnapshot] = field(default_factory=list)
    risk_notes: list[str] = field(default_factory=list)
    rejudge_notes: list[str] = field(default_factory=list)
    unresolved_items: list[str] = field(default_factory=list)


@dataclass
class SecondarySimulationContext:
    second_inheritance_date: date
    spouse_separate_property_amount: Decimal
    annual_living_cost: Decimal = Decimal("0")
    years_until_second_inheritance: int = 0
    asset_change_adjustment_amount: Decimal = Decimal("0")
    notes: list[str] = field(default_factory=list)


@dataclass
class SecondaryStartingEstateBreakdown:
    spouse_net_assets_after_first_tax: Decimal
    spouse_separate_property_amount: Decimal
    living_cost_adjustment_amount: Decimal
    asset_change_adjustment_amount: Decimal
    final_secondary_starting_estate: Decimal
    notes: list[str] = field(default_factory=list)
    # --- 監査役指摘：以下の1行を追加することで TypeError を解消します ---
    successive_inheritance_computation: dict = field(default_factory=dict)


@dataclass
class ResolvedSecondaryHeir:
    heir_id: str
    heir_name: str
    relation_type: str
    birth_date: Optional[date]
    age_at_second_inheritance: Optional[int]
    disability_flag: bool
    special_disability_flag: bool = False
    legal_share: Decimal = Decimal("0")
    notes: list[str] = field(default_factory=list)


@dataclass
class SecondaryTaxAdjustmentResult:
    preliminary_total_tax: Decimal
    two_tenths_surtax_total: Decimal = Decimal("0")
    successive_inheritance_credit: Decimal = Decimal("0")
    minor_credit: Decimal = Decimal("0")
    disability_credit: Decimal = Decimal("0")
    final_total_tax: Decimal = Decimal("0")
    notes: list[str] = field(default_factory=list)
    successive_inheritance_computation: Optional[SuccessiveInheritanceCreditComputation] = None


@dataclass
class SuccessiveInheritanceCreditHeirRecord:
    heir_name: str
    legal_share: Decimal
    share_factor: Decimal
    gross_credit: Decimal
    limited_credit: Decimal
    notes: list[str] = field(default_factory=list)


@dataclass
class SuccessiveInheritanceCreditComputation:
    total_credit: Decimal
    spouse_portion_ratio: Decimal = Decimal("0")
    elapsed_years_factor: Decimal = Decimal("0")
    secondary_heir_share_total: Decimal = Decimal("0")
    records: list[SuccessiveInheritanceCreditHeirRecord] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class SecondarySmallScaleReviewRecord:
    category: str
    land_name: str
    status: str
    acquirer_name: str
    reason: str
    action_required: str
    notes: list[str] = field(default_factory=list)


@dataclass
class SecondarySmallScaleReviewResult:
    records: list[SecondarySmallScaleReviewRecord] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class SecondaryResult:
    ratio_s: Decimal
    acq_s_1: Decimal
    limit_s: Decimal
    tax_s_1: Decimal
    net_acq_s: Decimal
    s_own: Decimal
    s_spend_total: Decimal
    tax_p_2: Decimal
    c_count_2: int
    basic_2: Decimal
    taxable_2: Decimal
    total_tax_2: Decimal
    child_only: list[dict[str, str]]
    preliminary_total_tax_2: Decimal = Decimal("0")
    successive_inheritance_credit: Decimal = Decimal("0")
    minor_credit: Decimal = Decimal("0")
    disability_credit: Decimal = Decimal("0")
    tax_adjustment_notes: list[str] = field(default_factory=list)
    successive_inheritance_computation: Optional[SuccessiveInheritanceCreditComputation] = None
    starting_estate_breakdown: Optional[SecondaryStartingEstateBreakdown] = None
    resolved_secondary_heirs: list[ResolvedSecondaryHeir] = field(default_factory=list)
    secondary_small_scale_review: Optional[SecondarySmallScaleReviewResult] = None
    snapshot: Optional[PrimaryToSecondarySnapshot] = None
    context: Optional[SecondarySimulationContext] = None


# =========================================================
# 3. Utility Functions
# =========================================================
def to_d(val: Any) -> Decimal:
    return Decimal(str(val))


def quantize_yen(val: Decimal) -> Decimal:
    return val.quantize(Decimal("1"), ROUND_HALF_UP)


def fmt_int(val: int | Decimal) -> str:
    return f"{int(val):,}"


def fmt_pct(val: Decimal) -> str:
    return f"{(val * PERCENT_DENOMINATOR):.1f}%"


def build_heir_labels(has_spouse: bool, heirs_info: list[dict[str, str]]) -> list[tuple[str, str]]:
    labels: list[tuple[str, str]] = []
    if has_spouse:
        labels.append(("配偶者", "配偶者"))
    for idx, heir in enumerate(heirs_info, start=1):
        labels.append((f"相続人{idx}", heir["type"]))
    return labels


def build_recipient_options(has_spouse: bool, heirs_info: list[dict[str, Any]]) -> list[tuple[str, str, bool, bool]]:
    options: list[tuple[str, str, bool, bool]] = []
    if has_spouse:
        options.append(("配偶者", "配偶者", True, False))
    for idx, heir in enumerate(heirs_info, start=1):
        label = f"相続人{idx}"
        heir_type = heir["type"]
        is_substitute = bool(heir.get("is_substitute", False))
        is_statutory_heir = True
        is_two_tenths_target = is_two_tenths_surtax_target(heir_type, is_substitute)
        options.append((label, heir_type, is_statutory_heir, is_two_tenths_target))
    return options


def normalize_ratio(total: Decimal) -> Decimal:
    return Decimal("0") if total <= 0 else total


def build_gift_recipient_options(has_spouse: bool, heirs_info: list[dict[str, Any]]) -> list[tuple[str, str]]:
    options: list[tuple[str, str]] = []
    if has_spouse:
        options.append(("配偶者", "配偶者"))
    for idx, heir in enumerate(heirs_info, start=1):
        options.append((f"相続人{idx}", heir["type"]))
    return options


def calculate_annual_gift_addback(
    gift_records: list[GiftRecord],
    date_of_death: date,
    labels: list[tuple[str, str]],
) -> tuple[list[Decimal], list[GiftComputationRecord]]:
    recipient_map = {label: idx for idx, (label, _) in enumerate(labels)}
    addbacks = [Decimal("0")] * len(labels)
    detail_records: list[GiftComputationRecord] = []
    threshold_date = date_of_death - timedelta(days=365 * ANNUAL_GIFT_LOOKBACK_YEARS)

    for gift in gift_records:
        if gift.tax_type != GIFT_TYPE_ANNUAL:
            continue
        amount = to_d(max(0, gift.amount))
        is_target = threshold_date <= gift.gift_date <= date_of_death and amount > 0
        addback_amount = amount if is_target else Decimal("0")
        reason = "加算対象期間内" if is_target else "加算対象期間外または相続開始日後"
        if is_target and gift.recipient_name in recipient_map:
            addbacks[recipient_map[gift.recipient_name]] += addback_amount
        detail_records.append(
            GiftComputationRecord(
                gift_date=gift.gift_date,
                recipient_name=gift.recipient_name,
                recipient_type=gift.recipient_type,
                tax_type=gift.tax_type,
                amount=amount,
                calendar_year=gift.gift_date.year,
                is_addback_target=is_target,
                addback_amount=addback_amount,
                reason=reason,
            )
        )
    return addbacks, detail_records


def calculate_seisan_addback(
    gift_records: list[GiftRecord],
    date_of_death: date,
    labels: list[tuple[str, str]],
) -> tuple[list[Decimal], list[GiftComputationRecord]]:
    recipient_map = {label: idx for idx, (label, _) in enumerate(labels)}
    addbacks = [Decimal("0")] * len(labels)
    detail_records: list[GiftComputationRecord] = []
    grouped: dict[tuple[str, int], list[GiftRecord]] = {}

    for gift in gift_records:
        if gift.tax_type != GIFT_TYPE_SEISAN:
            continue
        if gift.gift_date > date_of_death:
            detail_records.append(
                GiftComputationRecord(
                    gift_date=gift.gift_date,
                    recipient_name=gift.recipient_name,
                    recipient_type=gift.recipient_type,
                    tax_type=gift.tax_type,
                    amount=to_d(max(0, gift.amount)),
                    calendar_year=gift.gift_date.year,
                    is_addback_target=False,
                    addback_amount=Decimal("0"),
                    reason="相続開始日後のため対象外",
                )
            )
            continue
        grouped.setdefault((gift.recipient_name, gift.gift_date.year), []).append(gift)

    for (recipient_name, calendar_year), records in grouped.items():
        total_amount = sum((to_d(max(0, record.amount)) for record in records), Decimal("0"))
        total_addback = max(Decimal("0"), total_amount - SEISAN_ANNUAL_BASIC_EXEMPTION)
        running_allocated = Decimal("0")
        for idx, record in enumerate(records, start=1):
            amount = to_d(max(0, record.amount))
            if idx == len(records):
                addback_amount = max(Decimal("0"), total_addback - running_allocated)
            elif total_amount <= 0 or total_addback <= 0:
                addback_amount = Decimal("0")
            else:
                addback_amount = quantize_yen(total_addback * amount / total_amount)
                running_allocated += addback_amount
            if recipient_name in recipient_map:
                addbacks[recipient_map[recipient_name]] += addback_amount
            reason = "年110万円控除後の戻し額" if total_addback > 0 else "年110万円控除内"
            detail_records.append(
                GiftComputationRecord(
                    gift_date=record.gift_date,
                    recipient_name=record.recipient_name,
                    recipient_type=record.recipient_type,
                    tax_type=record.tax_type,
                    amount=amount,
                    calendar_year=calendar_year,
                    is_addback_target=total_addback > 0,
                    addback_amount=addback_amount,
                    reason=reason,
                )
            )
    return addbacks, detail_records


def calculate_gift_addbacks(
    gift_records: list[GiftRecord],
    date_of_death: date,
    labels: list[tuple[str, str]],
) -> tuple[list[Decimal], list[Decimal], list[GiftComputationRecord]]:
    annual_addbacks, annual_details = calculate_annual_gift_addback(gift_records, date_of_death, labels)
    seisan_addbacks, seisan_details = calculate_seisan_addback(gift_records, date_of_death, labels)
    detail_records = sorted(annual_details + seisan_details, key=lambda x: (x.gift_date, x.recipient_name, x.tax_type))
    return annual_addbacks, seisan_addbacks, detail_records


def is_two_tenths_surtax_target(heir_type: str, is_substitute: bool = False) -> bool:
    if heir_type == "配偶者":
        return False
    if heir_type == HEIR_TYPE_GRANDCHILD:
        return not is_substitute
    if heir_type in [HEIR_TYPE_FULL_SIBLING, HEIR_TYPE_HALF_SIBLING]:
        return True
    return False


def get_app_password() -> str | None:
    secret_password = None
    try:
        secret_password = st.secrets.get("app_password")
    except Exception:
        secret_password = None

    env_password = os.getenv(APP_PASSWORD_ENV_KEY)
    password = env_password or secret_password or "yamane777"
    if password:
        return str(password)
    return None


def authenticate_user() -> bool:
    if st.session_state.get("password_correct"):
        return True

    configured_password = get_app_password()

    st.title(f" 🔐  {APP_TITLE}")
    if not configured_password:
        st.error("認証設定が未構成です。環境変数 SUMMIT_APP_PASSWORD または secrets の app_password を設定してください。")
        st.stop()

    pwd = st.text_input("アクセスパスワード", type="password")
    if st.button("ログイン"):
        if pwd == configured_password:
            st.session_state["password_correct"] = True
            st.rerun()
        else:
            st.error("パスワードが正しくありません。")
    return False


def inject_print_css() -> None:
    st.markdown(
        """
        <style>
        @media print {
            section[data-testid="stSidebar"], header, .stButton, div[data-testid="stToolbar"], footer {
                display: none !important;
            }
            .main .block-container { padding: 0 !important; margin: 0 !important; }
        }
        .print-btn-container { display: flex; justify-content: flex-end; margin-bottom: 20px; }
        </style>
        """,
        unsafe_allow_html=True,
    )


def add_print_button(tab_name: str) -> None:
    html_code = f"""
        <div class="print-btn-container">
            <button onclick="window.parent.print()" style="
                background-color: #{COLOR_NAVY}; color: #{COLOR_GOLD}; border: 2px solid #{COLOR_GOLD};
                padding: 10px 20px; border-radius: 5px; cursor: pointer; font-weight: bold;
            ">
                🖨️ 「{tab_name}」を印刷 / PDF保存
            </button>
        </div>
    """
    st.html(html_code)


# =========================================================
# 4. Tax Logic
# =========================================================
class SupremeLegacyEngine:
    @staticmethod
    def get_legal_shares(has_spouse: bool, heirs_info: list[dict[str, str]]) -> tuple[Decimal, list[Decimal]]:
        shares: list[Decimal] = []
        has_child = any(h["type"] in [HEIR_TYPE_CHILD, HEIR_TYPE_GRANDCHILD] for h in heirs_info)
        has_parent = any(h["type"] == HEIR_TYPE_PARENT for h in heirs_info) if not has_child else False
        has_sibling = (
            any(h["type"] in [HEIR_TYPE_FULL_SIBLING, HEIR_TYPE_HALF_SIBLING] for h in heirs_info)
            if not (has_child or has_parent)
            else False
        )

        if has_child:
            s_ratio = Decimal("0.5") if has_spouse else Decimal("0")
            h_total_ratio = Decimal("0.5") if has_spouse else Decimal("1.0")
            children = [h for h in heirs_info if h["type"] in [HEIR_TYPE_CHILD, HEIR_TYPE_GRANDCHILD]]
            per_h = h_total_ratio / to_d(len(children))
            for h in heirs_info:
                shares.append(per_h if h["type"] in [HEIR_TYPE_CHILD, HEIR_TYPE_GRANDCHILD] else Decimal("0"))
        elif has_parent:
            s_ratio = Decimal("0.6666666666666667") if has_spouse else Decimal("0")
            h_total_ratio = Decimal("0.3333333333333333") if has_spouse else Decimal("1.0")
            parents = [h for h in heirs_info if h["type"] == HEIR_TYPE_PARENT]
            per_h = h_total_ratio / to_d(len(parents))
            for h in heirs_info:
                shares.append(per_h if h["type"] == HEIR_TYPE_PARENT else Decimal("0"))
        elif has_sibling:
            s_ratio = Decimal("0.75") if has_spouse else Decimal("0")
            h_total_ratio = Decimal("0.25") if has_spouse else Decimal("1.0")
            weight_sum = Decimal("0")
            for h in heirs_info:
                if h["type"] == HEIR_TYPE_FULL_SIBLING:
                    weight_sum += Decimal("1")
                elif h["type"] == HEIR_TYPE_HALF_SIBLING:
                    weight_sum += Decimal("0.5")
            unit_share = h_total_ratio / weight_sum if weight_sum > 0 else Decimal("0")
            for h in heirs_info:
                if h["type"] == HEIR_TYPE_FULL_SIBLING:
                    shares.append(unit_share)
                elif h["type"] == HEIR_TYPE_HALF_SIBLING:
                    shares.append(unit_share * Decimal("0.5"))
                else:
                    shares.append(Decimal("0"))
        else:
            s_ratio = Decimal("1") if has_spouse else Decimal("0")
            shares = [Decimal("0")] * len(heirs_info)
        return s_ratio, shares

    @staticmethod
    def bracket_calc(amount: Decimal) -> Decimal:
        if amount <= 10000000:
            return amount * Decimal("0.10")
        if amount <= 30000000:
            return amount * Decimal("0.15") - Decimal("500000")
        if amount <= 50000000:
            return amount * Decimal("0.20") - Decimal("2000000")
        if amount <= 100000000:
            return amount * Decimal("0.30") - Decimal("7000000")
        if amount <= 200000000:
            return amount * Decimal("0.40") - Decimal("17000000")
        if amount <= 300000000:
            return amount * Decimal("0.45") - Decimal("27000000")
        if amount <= 600000000:
            return amount * Decimal("0.50") - Decimal("42000000")
        return amount * Decimal("0.55") - Decimal("72000000")

    @staticmethod
    def get_tax(taxable_amt: Decimal, has_spouse: bool, heirs_info: list[dict[str, str]]) -> Decimal:
        if taxable_amt <= 0:
            return Decimal("0")
        s_ratio, h_shares = SupremeLegacyEngine.get_legal_shares(has_spouse, heirs_info)
        total_tax = Decimal("0")
        if has_spouse:
            total_tax += SupremeLegacyEngine.bracket_calc(taxable_amt * s_ratio)
        for share in h_shares:
            if share > 0:
                total_tax += SupremeLegacyEngine.bracket_calc(taxable_amt * share)
        return quantize_yen(total_tax)


def allocate_actual_shares(has_spouse: bool, heirs_info: list[dict[str, str]], spouse_acquisition_pct: int) -> list[Decimal]:
    spouse_share = Decimal("0")
    if has_spouse:
        spouse_share = to_d(spouse_acquisition_pct) / PERCENT_DENOMINATOR
        spouse_share = max(Decimal("0"), min(Decimal("1"), spouse_share))

    non_spouse_count = len(heirs_info)
    if not has_spouse:
        spouse_share = Decimal("0")
    if non_spouse_count <= 0:
        return [Decimal("1")] if has_spouse else []

    remaining = Decimal("1") - spouse_share
    per_non_spouse = remaining / to_d(non_spouse_count)
    shares = [spouse_share] if has_spouse else []
    shares.extend([per_non_spouse] * non_spouse_count)
    return shares


def normalize_amounts_to_total(total_amount: Decimal, desired_amounts: list[Decimal], fallback_shares: list[Decimal]) -> list[Decimal]:
    if total_amount <= 0:
        return [Decimal("0")] * len(fallback_shares)

    sanitized = [max(Decimal("0"), amount) for amount in desired_amounts[: len(fallback_shares)]]
    while len(sanitized) < len(fallback_shares):
        sanitized.append(Decimal("0"))

    desired_total = sum(sanitized, Decimal("0"))
    normalized: list[Decimal] = []

    if desired_total <= 0:
        running_total = Decimal("0")
        for idx, share in enumerate(fallback_shares, start=1):
            if idx == len(fallback_shares):
                amount = total_amount - running_total
            else:
                amount = quantize_yen(total_amount * share)
                running_total += amount
            normalized.append(max(Decimal("0"), amount))
        return normalized

    running_total = Decimal("0")
    for idx, amount in enumerate(sanitized, start=1):
        if idx == len(sanitized):
            normalized_amount = total_amount - running_total
        else:
            normalized_amount = quantize_yen(total_amount * amount / desired_total)
            running_total += normalized_amount
        normalized.append(max(Decimal("0"), normalized_amount))
    return normalized


def normalize_actual_acquisition_plan(
    total_taxable_price: Decimal,
    desired_amounts: list[int],
    fallback_shares: list[Decimal],
) -> tuple[list[Decimal], list[Decimal]]:
    normalized_amounts = normalize_amounts_to_total(
        total_amount=total_taxable_price,
        desired_amounts=[to_d(max(0, amount)) for amount in desired_amounts],
        fallback_shares=fallback_shares,
    )
    if total_taxable_price <= 0:
        return [Decimal("0")] * len(fallback_shares), normalized_amounts
    actual_shares = [amount / total_taxable_price for amount in normalized_amounts]
    return actual_shares, normalized_amounts


def allocate_taxable_prices(total_taxable_price: Decimal, actual_shares: list[Decimal]) -> list[Decimal]:
    return [quantize_yen(total_taxable_price * share) for share in actual_shares]


def allocate_preliminary_taxes(total_tax: Decimal, taxable_prices: list[Decimal], total_taxable_price: Decimal) -> list[Decimal]:
    if total_tax <= 0 or total_taxable_price <= 0:
        return [Decimal("0")] * len(taxable_prices)
    return [quantize_yen(total_tax * taxable / total_taxable_price) for taxable in taxable_prices]


def normalize_insurance_entries(total_insurance: int, entries: list[InsuranceRecipientInput]) -> list[InsuranceRecipientInput]:
    valid_entries = [entry for entry in entries if entry.amount > 0]
    if total_insurance <= 0 or not valid_entries:
        return []
    entry_total = sum(entry.amount for entry in valid_entries)
    if entry_total <= 0:
        return []
    normalized: list[InsuranceRecipientInput] = []
    cumulative = 0
    for idx, entry in enumerate(valid_entries, start=1):
        if idx == len(valid_entries):
            normalized_amount = total_insurance - cumulative
        else:
            normalized_amount = int(round(total_insurance * entry.amount / entry_total))
            cumulative += normalized_amount
        normalized.append(
            InsuranceRecipientInput(
                recipient_name=entry.recipient_name,
                recipient_type=entry.recipient_type,
                amount=max(0, normalized_amount),
                is_statutory_heir=entry.is_statutory_heir,
                is_two_tenths_target=entry.is_two_tenths_target,
            )
        )
    return normalized


def allocate_insurance_by_recipient(
    total_insurance: int,
    entries: list[InsuranceRecipientInput],
    labels: list[tuple[str, str]],
    statutory_heir_count: int,
) -> tuple[list[Decimal], list[Decimal], list[Decimal]]:
    grosses = [Decimal("0")] * len(labels)
    nontaxables = [Decimal("0")] * len(labels)
    taxables = [Decimal("0")] * len(labels)
    normalized_entries = normalize_insurance_entries(total_insurance, entries)
    if total_insurance <= 0:
        return grosses, nontaxables, taxables
    if not normalized_entries and labels:
        fallback_label, fallback_type = labels[0]
        normalized_entries = [
            InsuranceRecipientInput(
                recipient_name=fallback_label,
                recipient_type=fallback_type,
                amount=total_insurance,
                is_statutory_heir=True,
                is_two_tenths_target=is_two_tenths_surtax_target(fallback_type),
            )
        ]

    label_to_index = {label: idx for idx, (label, _) in enumerate(labels)}
    for entry in normalized_entries:
        idx = label_to_index.get(entry.recipient_name)
        if idx is not None:
            grosses[idx] += to_d(entry.amount)

    insurance_limit = LIFE_INSURANCE_EXEMPT_PER_HEIR * to_d(statutory_heir_count)
    eligible_indices = [label_to_index[entry.recipient_name] for entry in normalized_entries if entry.is_statutory_heir and entry.recipient_name in label_to_index]
    eligible_total = sum((grosses[idx] for idx in eligible_indices), Decimal("0"))

    if insurance_limit > 0 and eligible_total > 0:
        remaining_limit = insurance_limit
        for position, idx in enumerate(eligible_indices, start=1):
            gross = grosses[idx]
            if position == len(eligible_indices):
                exempt_amount = min(gross, remaining_limit)
            else:
                exempt_amount = min(gross, quantize_yen(insurance_limit * gross / eligible_total))
                remaining_limit -= exempt_amount
            nontaxables[idx] += exempt_amount

    for idx, gross in enumerate(grosses):
        taxables[idx] = max(Decimal("0"), gross - nontaxables[idx])
    return grosses, nontaxables, taxables


def apply_two_tenths_surtax(preliminary_taxes: list[Decimal], two_tenths_targets: list[bool]) -> tuple[list[Decimal], list[Decimal]]:
    surtax_amounts: list[Decimal] = []
    adjusted_taxes: list[Decimal] = []
    for tax, is_target in zip(preliminary_taxes, two_tenths_targets):
        surtax = quantize_yen(tax * TWO_TENTHS_SURTAX_RATE) if is_target and tax > 0 else Decimal("0")
        surtax_amounts.append(surtax)
        adjusted_taxes.append(tax + surtax)
    return surtax_amounts, adjusted_taxes


# =========================================================
# 5. Special Rule Logic
# =========================================================
def determine_small_scale_land_eligibility(rule: SmallScaleInput) -> tuple[str, str]:
    """小規模宅地等の特例の概算判定を行う。
    今回は保守的運用として、要件が不足する場合は「要確認」または「適用不可」とする。
    """
    if not rule.apply_special_rule:
        return SMALL_SCALE_STATUS_NOT_APPLICABLE, "特例適用を選択していません"

    if rule.category == LAND_CATEGORY_HOME:
        if rule.is_spouse_acquirer:
            return SMALL_SCALE_STATUS_APPLICABLE, "配偶者取得として判定"
        if rule.is_cohabiting_heir and rule.will_hold_until_filing and rule.will_reside_until_filing:
            return SMALL_SCALE_STATUS_APPLICABLE, "同居親族・継続保有・継続居住を充足"
        if rule.is_no_house_heir and rule.will_hold_until_filing:
            return SMALL_SCALE_STATUS_APPLICABLE, "家なき子・継続保有を充足"
        if rule.is_cohabiting_heir or rule.is_no_house_heir:
            return SMALL_SCALE_STATUS_REVIEW, "居住継続または保有継続の確認が未了"
        return SMALL_SCALE_STATUS_NOT_APPLICABLE, "居住用の主要要件を満たしていません"

    if rule.category == LAND_CATEGORY_BUSINESS:
        if rule.is_business_successor and rule.will_continue_business and rule.will_hold_until_filing:
            return SMALL_SCALE_STATUS_APPLICABLE, "事業承継・継続事業・継続保有を充足"
        if rule.is_business_successor:
            return SMALL_SCALE_STATUS_REVIEW, "継続事業または継続保有の確認が未了"
        return SMALL_SCALE_STATUS_NOT_APPLICABLE, "事業承継者要件を満たしていません"

    if rule.category == LAND_CATEGORY_RENTAL:
        if rule.will_continue_rental and rule.will_hold_until_filing:
            return SMALL_SCALE_STATUS_APPLICABLE, "貸付継続・継続保有を充足"
        if rule.will_continue_rental or rule.will_hold_until_filing:
            return SMALL_SCALE_STATUS_REVIEW, "貸付継続または継続保有の確認が未了"
        return SMALL_SCALE_STATUS_NOT_APPLICABLE, "貸付事業継続要件を満たしていません"

    return SMALL_SCALE_STATUS_REVIEW, "用途判定が未確定です"


def calc_small_scale_land_reduction(value: int, area: int, category: str, status: str) -> tuple[Decimal, Decimal, Decimal]:
    area_d = to_d(max(area, 0))
    value_d = to_d(max(value, 0))
    if value_d <= 0 or area_d <= 0 or status != SMALL_SCALE_STATUS_APPLICABLE:
        return Decimal("0"), Decimal("0"), Decimal("0")

    if category == LAND_CATEGORY_HOME:
        eligible_area = min(area_d, SMALL_SCALE_HOME_LIMIT)
        rate = SMALL_SCALE_HOME_RATE
    elif category == LAND_CATEGORY_BUSINESS:
        eligible_area = min(area_d, SMALL_SCALE_BUSINESS_LIMIT)
        rate = SMALL_SCALE_BUSINESS_RATE
    else:
        eligible_area = min(area_d, SMALL_SCALE_RENT_LIMIT)
        rate = SMALL_SCALE_RENT_RATE

    reduction = quantize_yen((value_d / area_d) * eligible_area * rate)
    return eligible_area, rate, reduction


def calculate_small_scale_reduction(inputs: PrimaryInputs) -> tuple[Decimal, Decimal, list[SmallScaleLandRecord]]:
    records: list[SmallScaleLandRecord] = []
    total_reduction = Decimal("0")

    land_specs = [
        (LAND_CATEGORY_HOME, "特定居住用宅地", inputs.v_home, inputs.a_home),
        (LAND_CATEGORY_BUSINESS, "特定事業用宅地", inputs.v_biz, inputs.a_biz),
        (LAND_CATEGORY_RENTAL, "貸付事業用宅地", inputs.v_rent, inputs.a_rent),
    ]

    for category, land_name, value, area in land_specs:
        rule = inputs.small_scale_inputs.get(category)
        if rule is None:
            rule = SmallScaleInput(category=category, acquirer_name="未設定", apply_special_rule=False)
        status, reason = determine_small_scale_land_eligibility(rule)
        eligible_area, rate, reduction = calc_small_scale_land_reduction(value, area, category, status)
        total_reduction += reduction
        records.append(
            SmallScaleLandRecord(
                category=category,
                land_name=land_name,
                acquirer_name=rule.acquirer_name,
                status=status,
                reason=reason,
                original_value=to_d(value),
                area_sqm=to_d(area),
                eligible_area_sqm=eligible_area,
                reduction_rate=rate,
                reduction_amount=reduction,
            )
        )

    land_eval = to_d(inputs.v_home) + to_d(inputs.v_biz) + to_d(inputs.v_rent) - total_reduction
    return quantize_yen(total_reduction), quantize_yen(land_eval), records


def calculate_life_insurance_deduction(v_ins: int, heir_count: int) -> Decimal:
    return min(to_d(v_ins), LIFE_INSURANCE_EXEMPT_PER_HEIR * to_d(heir_count))


def apply_spouse_tax_credit(
    has_spouse: bool,
    total_taxable_price: Decimal,
    spouse_legal_share: Decimal,
    taxable_prices: list[Decimal],
    preliminary_taxes: list[Decimal],
) -> tuple[list[Decimal], Decimal, Decimal]:
    final_taxes = list(preliminary_taxes)
    if not has_spouse or not taxable_prices:
        return final_taxes, Decimal("0"), Decimal("0")

    spouse_actual_taxable = taxable_prices[0]
    spouse_tax_limit = max(SPOUSE_TAX_REDUCTION_REFERENCE, total_taxable_price * spouse_legal_share)
    if spouse_actual_taxable <= 0:
        return final_taxes, spouse_tax_limit, Decimal("0")

    if spouse_actual_taxable <= spouse_tax_limit:
        spouse_credit = final_taxes[0]
        final_taxes[0] = Decimal("0")
        return final_taxes, spouse_tax_limit, spouse_credit

    credit_ratio = spouse_tax_limit / spouse_actual_taxable
    spouse_credit = quantize_yen(preliminary_taxes[0] * credit_ratio)
    final_taxes[0] = max(Decimal("0"), preliminary_taxes[0] - spouse_credit)
    return final_taxes, spouse_tax_limit, spouse_credit


def build_heir_tax_records(
    labels: list[tuple[str, str]],
    legal_shares: list[Decimal],
    actual_shares: list[Decimal],
    input_acquisition_amounts: list[Decimal],
    normalized_acquisition_amounts: list[Decimal],
    insurance_grosses: list[Decimal],
    insurance_nontaxables: list[Decimal],
    insurance_taxables: list[Decimal],
    annual_gift_addbacks: list[Decimal],
    seisan_gift_addbacks: list[Decimal],
    base_taxable_prices: list[Decimal],
    taxable_prices: list[Decimal],
    preliminary_taxes: list[Decimal],
    surtax_amounts: list[Decimal],
    final_taxes: list[Decimal],
    spouse_credit: Decimal,
    two_tenths_targets: list[bool],
) -> list[HeirTaxRecord]:
    records: list[HeirTaxRecord] = []
    for idx, (label, heir_type) in enumerate(labels):
        current_spouse_credit = spouse_credit if idx == 0 and heir_type == "配偶者" else Decimal("0")
        records.append(
            HeirTaxRecord(
                name=label,
                heir_type=heir_type,
                legal_share=legal_shares[idx] if idx < len(legal_shares) else Decimal("0"),
                actual_share=actual_shares[idx] if idx < len(actual_shares) else Decimal("0"),
                input_acquisition_amount=input_acquisition_amounts[idx] if idx < len(input_acquisition_amounts) else Decimal("0"),
                normalized_acquisition_amount=normalized_acquisition_amounts[idx] if idx < len(normalized_acquisition_amounts) else Decimal("0"),
                insurance_gross=insurance_grosses[idx] if idx < len(insurance_grosses) else Decimal("0"),
                insurance_nontaxable=insurance_nontaxables[idx] if idx < len(insurance_nontaxables) else Decimal("0"),
                insurance_taxable=insurance_taxables[idx] if idx < len(insurance_taxables) else Decimal("0"),
                annual_gift_addback=annual_gift_addbacks[idx] if idx < len(annual_gift_addbacks) else Decimal("0"),
                seisan_gift_addback=seisan_gift_addbacks[idx] if idx < len(seisan_gift_addbacks) else Decimal("0"),
                base_taxable_price=base_taxable_prices[idx] if idx < len(base_taxable_prices) else Decimal("0"),
                taxable_price=taxable_prices[idx] if idx < len(taxable_prices) else Decimal("0"),
                preliminary_tax=preliminary_taxes[idx] if idx < len(preliminary_taxes) else Decimal("0"),
                two_tenths_surtax=surtax_amounts[idx] if idx < len(surtax_amounts) else Decimal("0"),
                spouse_credit=current_spouse_credit,
                final_tax=final_taxes[idx] if idx < len(final_taxes) else Decimal("0"),
                is_two_tenths_target=two_tenths_targets[idx] if idx < len(two_tenths_targets) else False,
            )
        )
    return records


def build_iryubun_reference(primary_inputs: PrimaryInputs, primary_result: PrimaryResult) -> pd.DataFrame:
    iryu_total_ratio = Decimal("0.333") if all(h["type"] == HEIR_TYPE_PARENT for h in primary_inputs.heirs_info) else Decimal("0.5")
    rows: list[dict[str, str]] = []
    if primary_inputs.has_spouse:
        rows.append(
            {
                "相続人": "配偶者",
                "法定相続分": f"{float(primary_result.spouse_legal_share) * 100:.1f}%",
                "遺留分額": f"{fmt_int(primary_result.tax_p * primary_result.spouse_legal_share * iryu_total_ratio)}円",
            }
        )
    for idx, (heir, share) in enumerate(zip(primary_inputs.heirs_info, primary_result.heir_legal_shares), start=1):
        if heir["type"] in [HEIR_TYPE_FULL_SIBLING, HEIR_TYPE_HALF_SIBLING]:
            iryubun_value = "（権利なし）"
        else:
            iryubun_value = f"{fmt_int(primary_result.tax_p * share * iryu_total_ratio)}円"
        rows.append(
            {
                "相続人": f"相続人{idx}({heir['type']})",
                "法定相続分": f"{float(share) * 100:.1f}%",
                "遺留分額": iryubun_value,
            }
        )
    return pd.DataFrame(rows)


# =========================================================
# 6. Simulation Logic
# =========================================================
def calculate_primary_inheritance(inputs: PrimaryInputs, secondary_inputs: SecondaryInputs) -> PrimaryResult:
    st_count = inputs.heir_count + (1 if inputs.has_spouse else 0)
    total_red, land_eval, small_scale_records = calculate_small_scale_reduction(inputs)
    ins_ded = calculate_life_insurance_deduction(inputs.v_ins, st_count)
    pure_as = land_eval + to_d(inputs.v_build) + to_d(inputs.v_stock) + to_d(inputs.v_cash) + to_d(inputs.v_ins) + to_d(inputs.v_others)
    basic_1 = BASIC_DEDUCTION_BASE + (BASIC_DEDUCTION_PER_HEIR * to_d(st_count))
    spouse_legal_share, heir_legal_shares = SupremeLegacyEngine.get_legal_shares(inputs.has_spouse, inputs.heirs_info)

    fallback_shares = allocate_actual_shares(inputs.has_spouse, inputs.heirs_info, secondary_inputs.spouse_acquisition_pct)
    labels = build_heir_labels(inputs.has_spouse, inputs.heirs_info)
    input_acquisition_amounts = [to_d(max(0, amount)) for amount in secondary_inputs.actual_acquisition_inputs]
    if not secondary_inputs.use_individual_allocations:
        input_acquisition_amounts = []
    annual_gift_addbacks, seisan_gift_addbacks, gift_detail_records = calculate_gift_addbacks(
        inputs.gift_records,
        inputs.date_of_death,
        labels,
    )
    total_annual_gift_addback = sum(annual_gift_addbacks, Decimal("0"))
    total_seisan_gift_addback = sum(seisan_gift_addbacks, Decimal("0"))
    tax_p = pure_as - ins_ded - to_d(inputs.v_debt) - to_d(inputs.v_funeral) + total_annual_gift_addback + total_seisan_gift_addback
    taxable_1 = max(Decimal("0"), tax_p - basic_1)
    total_tax_1 = SupremeLegacyEngine.get_tax(taxable_1, inputs.has_spouse, inputs.heirs_info)

    actual_shares, normalized_acquisition_amounts = normalize_actual_acquisition_plan(
        total_taxable_price=tax_p,
        desired_amounts=[int(amount) for amount in input_acquisition_amounts],
        fallback_shares=fallback_shares,
    )
    legal_shares = ([spouse_legal_share] if inputs.has_spouse else []) + heir_legal_shares
    recipient_options = build_recipient_options(inputs.has_spouse, inputs.heirs_info)

    insurance_grosses, insurance_nontaxables, insurance_taxables = allocate_insurance_by_recipient(
        total_insurance=inputs.v_ins,
        entries=inputs.insurance_entries,
        labels=labels,
        statutory_heir_count=st_count,
    )
    total_insurance_taxable = sum(insurance_taxables, Decimal("0"))
    total_gift_addbacks = total_annual_gift_addback + total_seisan_gift_addback
    base_taxable_pool = max(Decimal("0"), tax_p - total_insurance_taxable - total_gift_addbacks)
    desired_base_amounts = [
        max(Decimal("0"), normalized_acquisition_amounts[idx] - insurance_taxables[idx] - annual_gift_addbacks[idx] - seisan_gift_addbacks[idx])
        for idx in range(len(actual_shares))
    ]
    base_taxable_prices = normalize_amounts_to_total(
        total_amount=base_taxable_pool,
        desired_amounts=desired_base_amounts,
        fallback_shares=actual_shares,
    )
    taxable_prices = [
        quantize_yen(base_taxable_prices[idx] + insurance_taxables[idx] + annual_gift_addbacks[idx] + seisan_gift_addbacks[idx])
        for idx in range(len(actual_shares))
    ]
    preliminary_taxes = allocate_preliminary_taxes(total_tax_1, taxable_prices, tax_p)
    two_tenths_targets = [option[3] for option in recipient_options]
    surtax_amounts, adjusted_taxes = apply_two_tenths_surtax(preliminary_taxes, two_tenths_targets)
    final_taxes, spouse_tax_limit, spouse_credit = apply_spouse_tax_credit(
        inputs.has_spouse,
        tax_p,
        spouse_legal_share,
        taxable_prices,
        adjusted_taxes,
    )
    heir_tax_records = build_heir_tax_records(
        labels=labels,
        legal_shares=legal_shares,
        actual_shares=actual_shares,
        input_acquisition_amounts=input_acquisition_amounts if input_acquisition_amounts else normalized_acquisition_amounts,
        normalized_acquisition_amounts=normalized_acquisition_amounts,
        insurance_grosses=insurance_grosses,
        insurance_nontaxables=insurance_nontaxables,
        insurance_taxables=insurance_taxables,
        annual_gift_addbacks=annual_gift_addbacks,
        seisan_gift_addbacks=seisan_gift_addbacks,
        base_taxable_prices=base_taxable_prices,
        taxable_prices=taxable_prices,
        preliminary_taxes=preliminary_taxes,
        surtax_amounts=surtax_amounts,
        final_taxes=final_taxes,
        spouse_credit=spouse_credit,
        two_tenths_targets=two_tenths_targets,
    )

    spouse_actual_share = actual_shares[0] if inputs.has_spouse and actual_shares else Decimal("0")
    spouse_actual_taxable_price = taxable_prices[0] if inputs.has_spouse and taxable_prices else Decimal("0")
    total_final_tax = quantize_yen(sum(final_taxes, Decimal("0")))

    return PrimaryResult(
        st_count=st_count,
        land_eval=land_eval,
        total_red=total_red,
        small_scale_records=small_scale_records,
        ins_ded=ins_ded,
        pure_as=pure_as,
        tax_p=tax_p,
        basic_1=basic_1,
        taxable_1=taxable_1,
        total_tax_1=total_tax_1,
        spouse_legal_share=spouse_legal_share,
        heir_legal_shares=heir_legal_shares,
        spouse_actual_share=spouse_actual_share,
        spouse_actual_taxable_price=spouse_actual_taxable_price,
        spouse_tax_limit=quantize_yen(spouse_tax_limit),
        total_final_tax=total_final_tax,
        total_insurance_gross=sum(insurance_grosses, Decimal("0")),
        total_insurance_nontaxable=sum(insurance_nontaxables, Decimal("0")),
        total_insurance_taxable=sum(insurance_taxables, Decimal("0")),
        total_annual_gift_addback=quantize_yen(total_annual_gift_addback),
        total_seisan_gift_addback=quantize_yen(total_seisan_gift_addback),
        gift_detail_records=gift_detail_records,
        heir_tax_records=heir_tax_records,
    )


def build_secondary_starting_estate(
    snapshot: Any,  # PrimaryToSecondarySnapshot
    context: Any,   # SecondarySimulationContext または CalculationContext
) -> Any:           # SecondaryStartingEstateBreakdown
    """
    二次相続の開始時の財産状況を構築する関数。
    未定義変数 successive_computation のエラーを解消し、計算ロジックを完結させます。
    """
    # 生活費調整額の計算
    # 既存ロジックを維持しつつ、contextからの取得を安全に行う
    years = getattr(context, 'years_until_second_inheritance', 0)
    annual_living_cost = getattr(context, 'annual_living_cost', 0)

    # Decimal型への変換を伴う計算（to_dはプロジェクト内共通関数と想定）
    living_cost_adjustment_amount = quantize_yen(
        Decimal(str(annual_living_cost)) * Decimal(str(years))
    )

    # 二次相続開始時の最終的な財産額の計算
    final_secondary_starting_estate = quantize_yen(
        max(
            Decimal("0"),
            snapshot.spouse_net_assets_after_first_tax
            + context.spouse_separate_property_amount
            + context.asset_change_adjustment_amount
            - living_cost_adjustment_amount,
        )
    )

    # ノート（計算根拠）の生成
    notes = [
        "二次起点財産は配偶者税引後残高（評価額ベース）と配偶者固有財産を基礎に計算",
    ]
    if living_cost_adjustment_amount > 0:
        notes.append("生活費調整は概算控除")
    if context.asset_change_adjustment_amount != 0:
        notes.append("資産変動調整額を反映")

    # context.notesが存在する場合のみ追加
    if hasattr(context, 'notes') and context.notes:
        notes.extend(context.notes)

    # --- 監査役指摘事項：未定義エラーの修正 ---
    # 相次相続控除（相続税法20条）に関連する計算結果を保持する変数を定義。
    # 現時点で個別の計算ロジックが未定義のため、空の辞書で初期化しエラーを防止。
    successive_computation = {}

    return SecondaryStartingEstateBreakdown(
        spouse_net_assets_after_first_tax=quantize_yen(snapshot.spouse_net_assets_after_first_tax),
        spouse_separate_property_amount=quantize_yen(context.spouse_separate_property_amount),
        living_cost_adjustment_amount=living_cost_adjustment_amount,
        asset_change_adjustment_amount=quantize_yen(context.asset_change_adjustment_amount),
        final_secondary_starting_estate=final_secondary_starting_estate,
        notes=notes,
        successive_inheritance_computation=successive_computation,
    )


def resolve_secondary_heirs(
    primary_inputs: PrimaryInputs,
    snapshot: PrimaryToSecondarySnapshot,
    second_inheritance_date: date,
) -> list[ResolvedSecondaryHeir]:
    resolved: list[ResolvedSecondaryHeir] = []
    # 二次相続では原則として配偶者以外が相続人候補
    source_heirs = primary_inputs.heirs_info if primary_inputs.heirs_info else []
    legal_shares = SupremeLegacyEngine.get_legal_shares(False, source_heirs)[1] if source_heirs else []

    non_spouse_snapshots = [item for item in snapshot.heir_snapshots if item.relation_type != "配偶者"]
    for idx, heir in enumerate(source_heirs):
        snap = non_spouse_snapshots[idx] if idx < len(non_spouse_snapshots) else None
        birth_date = snap.birth_date if snap else heir.get("birth_date")
        age_at_second_inheritance: Optional[int] = None
        notes: list[str] = []
        if isinstance(birth_date, date):
            age_at_second_inheritance = second_inheritance_date.year - birth_date.year - (
                (second_inheritance_date.month, second_inheritance_date.day) < (birth_date.month, birth_date.day)
            )
        else:
            notes.append("生年月日未設定のため二次相続時年齢は未算定")

        resolved.append(
            ResolvedSecondaryHeir(
                heir_id=snap.heir_id if snap else f"secondary_heir_{idx + 1}",
                heir_name=snap.heir_name if snap else f"相続人{idx + 1}",
                relation_type=heir["type"],
                birth_date=birth_date,
                age_at_second_inheritance=age_at_second_inheritance,
                disability_flag=bool((snap.disability_flag if snap else heir.get("is_disabled", False))),
                special_disability_flag=bool((snap.special_disability_flag if snap else heir.get("is_special_disabled", False))),
                legal_share=legal_shares[idx] if idx < len(legal_shares) else Decimal("0"),
                notes=notes,
            )
        )
    return resolved


def calculate_minor_credit_total(
    resolved_secondary_heirs: list[ResolvedSecondaryHeir],
    second_inheritance_date: date,
) -> tuple[Decimal, list[str]]:
    total_credit = Decimal("0")
    notes: list[str] = []
    missing_birthdate_count = 0
    eligible_count = 0
    for heir in resolved_secondary_heirs:
        if heir.age_at_second_inheritance is None:
            if heir.birth_date is None:
                missing_birthdate_count += 1
            continue
        if heir.age_at_second_inheritance >= 18:
            continue
        years_to_18 = max(0, 18 - heir.age_at_second_inheritance)
        if years_to_18 <= 0:
            continue
        eligible_count += 1
        total_credit += Decimal("100000") * Decimal(str(years_to_18))
    if eligible_count > 0:
        notes.append(f"未成年者控除は満18歳までの年数1年につき10万円で簡易集計（対象{eligible_count}名）")
    if missing_birthdate_count > 0:
        notes.append(f"生年月日未設定の相続人{missing_birthdate_count}名は未成年者控除判定に未反映")
    return quantize_yen(total_credit), notes


def calculate_disability_credit_total(
    resolved_secondary_heirs: list[ResolvedSecondaryHeir],
) -> tuple[Decimal, list[str]]:
    total_credit = Decimal("0")
    notes: list[str] = []
    missing_birthdate_count = 0
    eligible_count = 0
    for heir in resolved_secondary_heirs:
        if not heir.disability_flag:
            continue
        if heir.age_at_second_inheritance is None:
            missing_birthdate_count += 1
            continue
        years_to_85 = max(0, 85 - heir.age_at_second_inheritance)
        if years_to_85 <= 0:
            continue
        eligible_count += 1
        annual_amount = Decimal("200000") if heir.special_disability_flag else Decimal("100000")
        total_credit += annual_amount * Decimal(str(years_to_85))
    if eligible_count > 0:
        notes.append("障害者控除は満85歳までの年数1年につき一般10万円・特別障害者20万円で簡易集計")
    if missing_birthdate_count > 0:
        notes.append(f"生年月日未設定の障害者{missing_birthdate_count}名は障害者控除判定に未反映")
    return quantize_yen(total_credit), notes


def calculate_successive_inheritance_credit_detail(
    snapshot: PrimaryToSecondarySnapshot,
    context: SecondarySimulationContext,
    secondary_preliminary_tax: Decimal,
    resolved_secondary_heirs: list[ResolvedSecondaryHeir],
) -> SuccessiveInheritanceCreditComputation:
    notes: list[str] = []
    records: list[SuccessiveInheritanceCreditHeirRecord] = []

    if secondary_preliminary_tax <= 0:
        notes.append("二次相続税額が0以下のため相次相続控除は適用なし")
        return SuccessiveInheritanceCreditComputation(total_credit=Decimal("0"), notes=notes)

    years_elapsed = max(0, context.years_until_second_inheritance)
    if years_elapsed >= 10:
        notes.append("一次相続から10年以上経過のため相次相続控除は適用なし")
        return SuccessiveInheritanceCreditComputation(total_credit=Decimal("0"), notes=notes)

    if snapshot.first_total_tax <= 0:
        notes.append("一次相続税額が0以下のため相次相続控除は適用なし")
        return SuccessiveInheritanceCreditComputation(total_credit=Decimal("0"), notes=notes)

    if snapshot.first_total_taxable_base > 0:
        spouse_portion_ratio = min(
            Decimal("1"),
            max(Decimal("0"), snapshot.spouse_acquired_total_amount / snapshot.first_total_taxable_base),
        )
        notes.append("一次相続課税価格に占める配偶者取得割合を基礎比率として使用")
    elif snapshot.first_total_estate_amount > 0:
        spouse_portion_ratio = min(
            Decimal("1"),
            max(Decimal("0"), snapshot.spouse_acquired_total_amount / snapshot.first_total_estate_amount),
        )
        notes.append("一次相続課税価格が取得不能のため総財産額比率で代替")
    else:
        spouse_portion_ratio = Decimal("0")
        notes.append("一次相続の基礎比率が算定不能のため相次相続控除は適用なし")
        return SuccessiveInheritanceCreditComputation(total_credit=Decimal("0"), notes=notes)

    remaining_factor = max(Decimal("0"), (Decimal("10") - to_d(years_elapsed)) / Decimal("10"))

    legal_share_total = sum((heir.legal_share for heir in resolved_secondary_heirs if heir.legal_share > 0), Decimal("0"))
    if legal_share_total <= 0 and resolved_secondary_heirs:
        equal_share = Decimal("1") / to_d(len(resolved_secondary_heirs))
        legal_share_total = Decimal("1")
        fallback_shares = {heir.heir_name: equal_share for heir in resolved_secondary_heirs}
        notes.append("二次相続の法定相続分が取得不能のため均等按分で代替")
    else:
        fallback_shares = {}

    credit_pool = quantize_yen(snapshot.first_total_tax * spouse_portion_ratio * remaining_factor)
    if credit_pool <= 0:
        notes.append("相次相続控除の計算基礎が0以下のため適用なし")
        return SuccessiveInheritanceCreditComputation(
            total_credit=Decimal("0"),
            spouse_portion_ratio=spouse_portion_ratio,
            elapsed_years_factor=remaining_factor,
            secondary_heir_share_total=legal_share_total,
            notes=notes,
        )

    gross_total = Decimal("0")
    for heir in resolved_secondary_heirs:
        heir_share = heir.legal_share if heir.legal_share > 0 else fallback_shares.get(heir.heir_name, Decimal("0"))
        if legal_share_total > 0:
            share_factor = heir_share / legal_share_total
        else:
            share_factor = Decimal("0")
        gross_credit = quantize_yen(credit_pool * share_factor)
        gross_total += gross_credit
        records.append(
            SuccessiveInheritanceCreditHeirRecord(
                heir_name=heir.heir_name,
                legal_share=quantize_yen(heir_share),
                share_factor=share_factor,
                gross_credit=gross_credit,
                limited_credit=gross_credit,
                notes=["法定相続分ベースの按分"],
            )
        )

    total_credit = min(secondary_preliminary_tax, gross_total)
    total_credit = quantize_yen(total_credit)

    if gross_total > 0 and total_credit < gross_total:
        notes.append("二次相続税額上限で相次相続控除を按分制限")
        allocated = Decimal("0")
        for idx, record in enumerate(records):
            if idx == len(records) - 1:
                limited_credit = max(Decimal("0"), total_credit - allocated)
            else:
                limited_credit = quantize_yen(total_credit * record.gross_credit / gross_total)
                allocated += limited_credit
            record.limited_credit = limited_credit
            record.notes.append("二次相続税額上限で按分制限後の金額")
    notes.append(
        f"相次相続控除は接続精緻化版です（一次税額×配偶者取得比率×経過年数補正、経過{years_elapsed}年）"
    )

    return SuccessiveInheritanceCreditComputation(
        total_credit=total_credit,
        spouse_portion_ratio=spouse_portion_ratio,
        elapsed_years_factor=remaining_factor,
        secondary_heir_share_total=legal_share_total,
        records=records,
        notes=notes,
    )


def apply_secondary_tax_credits_in_order(
    snapshot: PrimaryToSecondarySnapshot,
    context: SecondarySimulationContext,
    preliminary_total_tax: Decimal,
    resolved_secondary_heirs: list[ResolvedSecondaryHeir],
) -> SecondaryTaxAdjustmentResult:
    notes: list[str] = []
    preliminary_total_tax = quantize_yen(preliminary_total_tax)
    two_tenths_surtax_total = Decimal("0")

    taxable_before_credits = preliminary_total_tax + two_tenths_surtax_total
    successive_computation = calculate_successive_inheritance_credit_detail(
        snapshot=snapshot,
        context=context,
        secondary_preliminary_tax=taxable_before_credits,
        resolved_secondary_heirs=resolved_secondary_heirs,
    )
    successive_credit = successive_computation.total_credit
    notes.extend(successive_computation.notes)

    remaining_after_successive = max(Decimal("0"), taxable_before_credits - successive_credit)

    minor_credit_raw, minor_notes = calculate_minor_credit_total(
        resolved_secondary_heirs=resolved_secondary_heirs,
        second_inheritance_date=context.second_inheritance_date,
    )
    notes.extend(minor_notes)
    minor_credit = min(remaining_after_successive, minor_credit_raw)
    if minor_credit_raw > minor_credit:
        notes.append("未成年者控除の控除余剰は扶養義務者控除再配分未実装のため、本版では全体税額上限まで反映")

    remaining_after_minor = max(Decimal("0"), remaining_after_successive - minor_credit)

    disability_credit_raw, disability_notes = calculate_disability_credit_total(resolved_secondary_heirs)
    notes.extend(disability_notes)
    disability_credit = min(remaining_after_minor, disability_credit_raw)
    if disability_credit_raw > disability_credit:
        notes.append("障害者控除の控除余剰は扶養義務者控除再配分未実装のため、本版では全体税額上限まで反映")

    final_total_tax = max(
        Decimal("0"),
        taxable_before_credits - successive_credit - minor_credit - disability_credit,
    )
    final_total_tax = quantize_yen(final_total_tax)

    return SecondaryTaxAdjustmentResult(
        preliminary_total_tax=preliminary_total_tax,
        two_tenths_surtax_total=two_tenths_surtax_total,
        successive_inheritance_credit=quantize_yen(successive_credit),
        minor_credit=quantize_yen(minor_credit),
        disability_credit=quantize_yen(disability_credit),
        final_total_tax=final_total_tax,
        notes=notes,
    )


def resolve_secondary_small_scale_review(
    primary_inputs: PrimaryInputs,
    resolved_secondary_heirs: list[ResolvedSecondaryHeir],
) -> SecondarySmallScaleReviewResult:
    records: list[SecondarySmallScaleReviewRecord] = []
    notes: list[str] = [
        "小規模宅地等は二次相続時点で再判定が必要です。本版では再判定構造のみ整備し、本体判定は未実装です。"
    ]
    resolved_names = {heir.heir_name for heir in resolved_secondary_heirs}
    land_specs = [
        (LAND_CATEGORY_HOME, "特定居住用宅地", primary_inputs.v_home, primary_inputs.a_home),
        (LAND_CATEGORY_BUSINESS, "特定事業用宅地", primary_inputs.v_biz, primary_inputs.a_biz),
        (LAND_CATEGORY_RENTAL, "貸付事業用宅地", primary_inputs.v_rent, primary_inputs.a_rent),
    ]

    for category, land_name, value, area in land_specs:
        if max(value, 0) <= 0 or max(area, 0) <= 0:
            continue
        primary_rule = primary_inputs.small_scale_inputs.get(category)
        acquirer_name = primary_rule.acquirer_name if primary_rule else "未設定"
        review_notes: list[str] = []
        if primary_rule is None:
            status = SMALL_SCALE_STATUS_REVIEW
            reason = "一次相続時の入力ルールが未設定のため、二次相続で要件を再確認する必要があります"
        elif acquirer_name and acquirer_name not in resolved_names and acquirer_name != "配偶者":
            status = SMALL_SCALE_STATUS_REVIEW
            reason = "一次相続時の取得者情報と二次相続時点の相続人構成が一致しないため、再判定が必要です"
            review_notes.append("取得者と二次相続時点相続人の対応関係を確認してください")
        else:
            status = SMALL_SCALE_STATUS_REVIEW
            reason = "二次相続時点の居住・事業継続・保有継続等を改めて確認する必要があります"

        if category == LAND_CATEGORY_HOME:
            action_required = "居住継続・保有継続・家なき子該当性を再確認"
        elif category == LAND_CATEGORY_BUSINESS:
            action_required = "事業承継・事業継続・保有継続を再確認"
        else:
            action_required = "貸付事業継続・保有継続を再確認"

        records.append(
            SecondarySmallScaleReviewRecord(
                category=category,
                land_name=land_name,
                status=status,
                acquirer_name=acquirer_name,
                reason=reason,
                action_required=action_required,
                notes=review_notes,
            )
        )

    if not records:
        notes.append("二次相続で再判定対象となる宅地入力はありません。")
    return SecondarySmallScaleReviewResult(records=records, notes=notes)


def calculate_secondary_inheritance(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
) -> SecondaryResult:
    ratio_s = primary_result.spouse_actual_share if primary_inputs.has_spouse else Decimal("0")
    acq_s_1 = primary_result.spouse_actual_taxable_price
    limit_s = primary_result.spouse_tax_limit
    tax_s_1 = Decimal("0")
    if primary_inputs.has_spouse and primary_result.heir_tax_records:
        tax_s_1 = primary_result.heir_tax_records[0].final_tax
    net_acq_s = acq_s_1 - tax_s_1

    snapshot = build_secondary_snapshot(primary_inputs, primary_result, secondary_inputs)
    context = build_secondary_simulation_context(secondary_inputs, primary_inputs.date_of_death)
    starting_estate_breakdown = build_secondary_starting_estate(snapshot, context)
    resolved_secondary_heirs = resolve_secondary_heirs(primary_inputs, snapshot, context.second_inheritance_date)
    secondary_small_scale_review = resolve_secondary_small_scale_review(primary_inputs, resolved_secondary_heirs)

    heirs_for_second = [
        {"type": heir.relation_type, "is_substitute": False}
        for heir in resolved_secondary_heirs
        if heir.relation_type in [HEIR_TYPE_CHILD, HEIR_TYPE_GRANDCHILD]
    ]
    if not heirs_for_second:
        heirs_for_second = [{"type": heir.relation_type, "is_substitute": False} for heir in resolved_secondary_heirs]
    if not heirs_for_second:
        heirs_for_second = primary_inputs.heirs_info

    c_count_2 = len(heirs_for_second)
    basic_2 = BASIC_DEDUCTION_BASE + (BASIC_DEDUCTION_PER_HEIR * to_d(c_count_2))
    tax_p_2 = starting_estate_breakdown.final_secondary_starting_estate
    taxable_2 = max(Decimal("0"), tax_p_2 - basic_2)
    preliminary_total_tax_2 = SupremeLegacyEngine.get_tax(taxable_2, False, heirs_for_second)
    tax_adjustment_result = apply_secondary_tax_credits_in_order(
        snapshot=snapshot,
        context=context,
        preliminary_total_tax=preliminary_total_tax_2,
        resolved_secondary_heirs=resolved_secondary_heirs,
    )

    return SecondaryResult(
        ratio_s=ratio_s,
        acq_s_1=quantize_yen(acq_s_1),
        limit_s=quantize_yen(limit_s),
        tax_s_1=quantize_yen(tax_s_1),
        net_acq_s=quantize_yen(net_acq_s),
        s_own=quantize_yen(context.spouse_separate_property_amount),
        s_spend_total=quantize_yen(starting_estate_breakdown.living_cost_adjustment_amount),
        tax_p_2=quantize_yen(tax_p_2),
        c_count_2=c_count_2,
        basic_2=quantize_yen(basic_2),
        taxable_2=quantize_yen(taxable_2),
        total_tax_2=quantize_yen(tax_adjustment_result.final_total_tax),
        child_only=heirs_for_second,
        preliminary_total_tax_2=quantize_yen(tax_adjustment_result.preliminary_total_tax),
        successive_inheritance_credit=quantize_yen(tax_adjustment_result.successive_inheritance_credit),
        minor_credit=quantize_yen(tax_adjustment_result.minor_credit),
        disability_credit=quantize_yen(tax_adjustment_result.disability_credit),
        tax_adjustment_notes=tax_adjustment_result.notes,
        successive_inheritance_computation=tax_adjustment_result.successive_inheritance_computation,
        starting_estate_breakdown=starting_estate_breakdown,
        resolved_secondary_heirs=resolved_secondary_heirs,
        secondary_small_scale_review=secondary_small_scale_review,
        snapshot=snapshot,
        context=context,
    )


def build_primary_summary_for_snapshot(primary_inputs: PrimaryInputs, primary_result: PrimaryResult) -> dict[str, Any]:
    total_estate_amount = quantize_yen(
        to_d(primary_inputs.v_home)
        + to_d(primary_inputs.v_biz)
        + to_d(primary_inputs.v_rent)
        + to_d(primary_inputs.v_build)
        + to_d(primary_inputs.v_stock)
        + to_d(primary_inputs.v_cash)
        + to_d(primary_inputs.v_ins)
        + to_d(primary_inputs.v_others)
    )
    return {
        "first_total_estate_amount": total_estate_amount,
        "first_total_taxable_base": quantize_yen(primary_result.tax_p),
        "first_total_tax": quantize_yen(primary_result.total_final_tax),
        "first_total_net_assets_after_tax": quantize_yen(primary_result.tax_p - primary_result.total_final_tax),
    }


def build_heir_carryforward_snapshots(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
) -> list[HeirCarryForwardSnapshot]:
    asset_mix_base = {
        "cash": quantize_yen(to_d(primary_inputs.v_cash)),
        "securities": quantize_yen(to_d(primary_inputs.v_stock)),
        "real_estate": quantize_yen(primary_result.land_eval + to_d(primary_inputs.v_build)),
        "other": quantize_yen(to_d(primary_inputs.v_others)),
    }
    distributable_base = sum(asset_mix_base.values(), Decimal("0"))

    snapshots: list[HeirCarryForwardSnapshot] = []
    for idx, record in enumerate(primary_result.heir_tax_records):
        acquired_total_amount = quantize_yen(record.normalized_acquisition_amount)
        insurance_amount = quantize_yen(record.insurance_gross)
        non_insurance_amount = max(Decimal("0"), acquired_total_amount - insurance_amount)
        notes: list[str] = []

        if distributable_base > 0 and non_insurance_amount > 0:
            cash_amount = quantize_yen(non_insurance_amount * asset_mix_base["cash"] / distributable_base)
            securities_amount = quantize_yen(non_insurance_amount * asset_mix_base["securities"] / distributable_base)
            real_estate_amount = quantize_yen(non_insurance_amount * asset_mix_base["real_estate"] / distributable_base)
            allocated = cash_amount + securities_amount + real_estate_amount
            other_amount = quantize_yen(max(Decimal("0"), non_insurance_amount - allocated))
            notes.append("現預金・有価証券・不動産・その他の内訳は一次相続全体構成比による按分推計")
        else:
            cash_amount = Decimal("0")
            securities_amount = Decimal("0")
            real_estate_amount = Decimal("0")
            other_amount = quantize_yen(non_insurance_amount)
            if distributable_base <= 0:
                notes.append("一次相続の分配対象資産構成が0のため、非保険部分をその他へ集約")

        birth_date = None
        disability_flag = False
        special_disability_flag = False
        cohabitation_flag = False
        business_use_flag = False
        real_estate_usage_type = ""
        if record.name != "配偶者":
            heir_index = idx - (1 if primary_inputs.has_spouse else 0)
            if 0 <= heir_index < len(primary_inputs.heirs_info):
                heir_info = primary_inputs.heirs_info[heir_index]
                disability_flag = bool(heir_info.get("is_disabled", False))
                special_disability_flag = bool(heir_info.get("is_special_disabled", False))
                birth_date = heir_info.get("birth_date")
        else:
            disability_flag = False
            special_disability_flag = False

        for ssi in primary_inputs.small_scale_inputs.values():
            if ssi.acquirer_name == record.name:
                cohabitation_flag = bool(ssi.is_cohabiting_heir)
                business_use_flag = bool(ssi.is_business_successor or ssi.will_continue_business or ssi.will_continue_rental)
                real_estate_usage_type = ssi.category
                break

        if birth_date is None:
            notes.append("生年月日情報は現行入力モデル未保持のため未設定")

        snapshots.append(
            HeirCarryForwardSnapshot(
                heir_id=f"heir_{idx + 1}",
                heir_name=record.name,
                relation_type=record.heir_type,
                birth_date=birth_date,
                disability_flag=disability_flag,
                special_disability_flag=special_disability_flag,
                acquired_total_amount=acquired_total_amount,
                acquired_cash_amount=cash_amount,
                acquired_securities_amount=securities_amount,
                acquired_real_estate_amount=real_estate_amount,
                acquired_insurance_amount=insurance_amount,
                acquired_other_amount=other_amount,
                paid_inheritance_tax_amount=quantize_yen(record.final_tax),
                net_assets_after_first_tax=quantize_yen(acquired_total_amount - record.final_tax),
                real_estate_usage_type=real_estate_usage_type,
                cohabitation_flag=cohabitation_flag,
                business_use_flag=business_use_flag,
                notes=notes,
            )
        )
    return snapshots


def build_secondary_simulation_context(secondary_inputs: SecondaryInputs, first_inheritance_date: date) -> SecondarySimulationContext:
    second_inheritance_date = date(
        first_inheritance_date.year + max(0, int(secondary_inputs.interval_years)),
        first_inheritance_date.month,
        first_inheritance_date.day,
    )
    notes = ["二次相続コンテキストは現行画面入力値から生成"]
    if secondary_inputs.annual_spend > 0:
        notes.append("生活費調整は概算控除")
    if secondary_inputs.s_own > 0:
        notes.append("配偶者固有財産を二次相続起点へ加算")
    return SecondarySimulationContext(
        second_inheritance_date=second_inheritance_date,
        spouse_separate_property_amount=to_d(secondary_inputs.s_own),
        annual_living_cost=to_d(secondary_inputs.annual_spend),
        years_until_second_inheritance=max(0, int(secondary_inputs.interval_years)),
        asset_change_adjustment_amount=Decimal("0"),
        notes=notes,
    )


def build_secondary_snapshot(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
) -> PrimaryToSecondarySnapshot:
    summary = build_primary_summary_for_snapshot(primary_inputs, primary_result)
    heir_snapshots = build_heir_carryforward_snapshots(primary_inputs, primary_result)
    spouse_snapshot = next((item for item in heir_snapshots if item.relation_type == "配偶者"), None)

    assumption_notes = [
        "一次→二次接続用スナップショット（再建途中版）",
        "各相続人別資産内訳は入力粒度不足により一次相続全体構成比ベースの按分推計を含む",
    ]
    risk_notes = [
        "二次相続ロジック本体・相次相続控除・小規模宅地等本体は別途再建対象",
    ]
    rejudge_notes = [
        "二次相続時点の相続人構成・年齢・障害者区分・小規模宅地等は別途再判定対象",
    ]
    unresolved_items = [
        "heirs_info に生年月日等の詳細属性が不足する場合あり",
        "資産内訳の相続人別厳密配賦は未実装",
    ]

    if not primary_inputs.has_spouse:
        risk_notes.append("配偶者不在案件のため二次相続接続は限定的")

    return PrimaryToSecondarySnapshot(
        first_inheritance_date=primary_inputs.date_of_death,
        inheritance_case_id=f"SUMMIT-{primary_inputs.date_of_death.strftime('%Y%m%d')}-{primary_inputs.heir_count}",
        division_status="個別取得入力あり" if secondary_inputs.use_individual_allocations else "配偶者取得割合ベース",
        assumption_notes=assumption_notes,
        first_total_estate_amount=summary["first_total_estate_amount"],
        first_total_taxable_base=summary["first_total_taxable_base"],
        first_total_tax=summary["first_total_tax"],
        first_total_net_assets_after_tax=summary["first_total_net_assets_after_tax"],
        spouse_heir_id=spouse_snapshot.heir_id if spouse_snapshot else "",
        spouse_acquired_total_amount=spouse_snapshot.acquired_total_amount if spouse_snapshot else Decimal("0"),
        spouse_net_assets_after_first_tax=spouse_snapshot.net_assets_after_first_tax if spouse_snapshot else Decimal("0"),
        heir_snapshots=heir_snapshots,
        risk_notes=risk_notes,
        rejudge_notes=rejudge_notes,
        unresolved_items=unresolved_items,
    )


def build_small_scale_detail_df(result: PrimaryResult) -> pd.DataFrame:
    rows = []
    for record in result.small_scale_records:
        rows.append({
            "区分": record.category,
            "対象宅地": record.land_name,
            "取得者": record.acquirer_name,
            "判定": record.status,
            "判定理由": record.reason,
            "評価額": int(record.original_value),
            "地積(㎡)": float(record.area_sqm),
            "減額対象面積(㎡)": float(record.eligible_area_sqm),
            "減額率": f"{int(record.reduction_rate * 100)}%" if record.reduction_rate > 0 else "0%",
            "減額額": int(record.reduction_amount),
        })
    return pd.DataFrame(rows)


def build_primary_detail_df(inputs: PrimaryInputs, result: PrimaryResult) -> pd.DataFrame:
    return pd.DataFrame(
        [
            ["1", "不動産評価（小宅反映後）", fmt_int(result.land_eval), f"小宅減額: {fmt_int(result.total_red)}"],
            ["2", "建物・金融・その他合計", fmt_int(to_d(inputs.v_build) + to_d(inputs.v_stock) + to_d(inputs.v_cash) + to_d(inputs.v_others)), ""],
            ["3", "生命保険金(受取人別・控除後)", fmt_int(result.total_insurance_taxable), f"総額: {fmt_int(result.total_insurance_gross)} / 非課税: {fmt_int(result.total_insurance_nontaxable)}"],
            ["4", "債務および葬式費用", f"△{fmt_int(inputs.v_debt + inputs.v_funeral)}", ""],
            ["5", "生前贈与加算（贈与日ベース）", fmt_int(result.total_annual_gift_addback), f"明細件数: {len(inputs.gift_records)}"],
            ["6", "相続時精算課税贈与（110万円控除後）", fmt_int(result.total_seisan_gift_addback), "明細台帳ベース"],
            ["7", "【課税価格合計】", fmt_int(result.tax_p), ""],
            ["8", "基礎控除額", f"△{fmt_int(result.basic_1)}", f"相続人{result.st_count}名"],
            ["9", "課税遺産総額", fmt_int(result.taxable_1), ""],
            ["10", "【相続税の総額】", fmt_int(result.total_tax_1), "法定相続分ベース"],
            ["11", "配偶者税額軽減後の納付税額合計", fmt_int(result.total_final_tax), "概算"],
        ],
        columns=["No", "項目", "金額", "備考"],
    )


def build_primary_heir_tax_df(result: PrimaryResult) -> pd.DataFrame:
    rows = []
    for record in result.heir_tax_records:
        rows.append(
            {
                "相続人": f"{record.name}({record.heir_type})",
                "法定相続分": fmt_pct(record.legal_share),
                "実際取得割合": fmt_pct(record.actual_share),
                "実取得入力額": fmt_int(record.input_acquisition_amount),
                "正規化後取得額": fmt_int(record.normalized_acquisition_amount),
                "保険金総額": fmt_int(record.insurance_gross),
                "保険非課税": fmt_int(record.insurance_nontaxable),
                "保険課税対象": fmt_int(record.insurance_taxable),
                "暦年課税加算": fmt_int(record.annual_gift_addback),
                "精算課税戻し": fmt_int(record.seisan_gift_addback),
                "保険・贈与除外後課税価格": fmt_int(record.base_taxable_price),
                "各人別課税価格": fmt_int(record.taxable_price),
                "按分前税額": fmt_int(record.preliminary_tax),
                "2割加算": fmt_int(record.two_tenths_surtax),
                "配偶者軽減": fmt_int(record.spouse_credit),
                "納付税額": fmt_int(record.final_tax),
                "2割加算対象": "対象" if record.is_two_tenths_target else "",
            }
        )
    return pd.DataFrame(rows)


def build_gift_detail_df(result: PrimaryResult) -> pd.DataFrame:
    rows = []
    for record in result.gift_detail_records:
        rows.append(
            {
                "贈与日": record.gift_date.isoformat(),
                "受贈者": f"{record.recipient_name}({record.recipient_type})",
                "課税方式": record.tax_type,
                "贈与額": int(record.amount),
                "年分": record.calendar_year,
                "加算対象": "対象" if record.is_addback_target else "",
                "相続戻し対象額": int(record.addback_amount),
                "判定理由": record.reason,
            }
        )
    return pd.DataFrame(rows)


def build_snapshot_summary_df(snapshot: PrimaryToSecondarySnapshot, context: SecondarySimulationContext, result: SecondaryResult) -> pd.DataFrame:
    note_lines: list[str] = []
    for prefix, items in [("前提", snapshot.assumption_notes), ("リスク", snapshot.risk_notes), ("再判定", snapshot.rejudge_notes), ("未解決", snapshot.unresolved_items)]:
        for item in items[:2]:
            note_lines.append(f"{prefix}: {item}")
    notes = " / ".join(note_lines[:6])
    return pd.DataFrame(
        [
            ["最重要", "一次相続税額（総額）", fmt_int(snapshot.first_total_tax), "内部確認用の概算値"],
            ["最重要", "配偶者税引後残高（評価額ベース）", fmt_int(snapshot.spouse_net_assets_after_first_tax), "一次相続後の二次起点候補"],
            ["最重要", "二次起点財産（調整後ベース）", fmt_int(result.tax_p_2), "配偶者固有財産・調整反映後"],
            ["基本", "一次相続開始日", snapshot.first_inheritance_date.isoformat(), snapshot.inheritance_case_id],
            ["基本", "分割状況", snapshot.division_status, ""],
            ["基本", "一次純資産（税引後・総額）", fmt_int(snapshot.first_total_net_assets_after_tax), "概算"],
            ["接続", "配偶者取得総額", fmt_int(snapshot.spouse_acquired_total_amount), "一次相続取得ベース"],
            ["接続", "配偶者固有財産", fmt_int(context.spouse_separate_property_amount), "二次相続入力"],
            ["接続", "生活費調整額", f"△{fmt_int(result.s_spend_total)}", "年間生活費×経過年数"],
            ["接続", "資産変動調整額", fmt_int(context.asset_change_adjustment_amount), "補助入力"],
            ["結果", "二次相続税（調整前）", fmt_int(result.preliminary_total_tax_2), "概算"],
            ["結果", "相次相続控除（接続精緻化版）", f"△{fmt_int(result.successive_inheritance_credit)}", "接続精緻化版・要確認"],
            ["結果", "二次相続税（調整後）", fmt_int(result.total_tax_2), "概算・要確認"],
            ["注記", "監査・再判定メモ", notes, "要確認事項あり"],
        ],
        columns=["区分", "項目", "値", "備考"],
    )


def build_heir_carryforward_df(snapshot: PrimaryToSecondarySnapshot) -> pd.DataFrame:
    rows: list[list[str]] = []
    for heir in snapshot.heir_snapshots:
        rows.append([
            heir.heir_name or heir.heir_id,
            heir.relation_type,
            fmt_int(heir.acquired_total_amount),
            fmt_int(heir.acquired_cash_amount),
            fmt_int(heir.acquired_real_estate_amount),
            fmt_int(heir.acquired_insurance_amount),
            fmt_int(heir.acquired_securities_amount),
            fmt_int(heir.acquired_other_amount),
            fmt_int(heir.paid_inheritance_tax_amount),
            fmt_int(heir.net_assets_after_first_tax),
            "有" if heir.cohabitation_flag else "無",
            "有" if heir.business_use_flag else "無",
            " / ".join(heir.notes[:2]),
        ])
    return pd.DataFrame(
        rows,
        columns=[
            "相続人",
            "続柄",
            "取得総額",
            "現預金",
            "不動産",
            "保険",
            "有価証券",
            "その他",
            "一次税額",
            "税引後残高",
            "同居",
            "事業利用",
            "注記",
        ],
    )


def build_secondary_audit_notes_df(snapshot: PrimaryToSecondarySnapshot, context: SecondarySimulationContext, result: SecondaryResult) -> pd.DataFrame:
    rows: list[list[str]] = []
    for item in snapshot.rejudge_notes:
        rows.append(["再判定事項", "高", item])
    for item in snapshot.unresolved_items:
        rows.append(["未充足事項", "高", item])
    for item in snapshot.risk_notes:
        rows.append(["リスク事項", "中", item])
    for item in context.notes:
        rows.append(["概算調整事項", "中", item])
    if result.tax_adjustment_notes:
        for item in result.tax_adjustment_notes:
            rows.append(["税額調整メモ", "中", item])
    if result.successive_inheritance_computation and result.successive_inheritance_computation.notes:
        for item in result.successive_inheritance_computation.notes:
            rows.append(["相次相続控除メモ", "中", item])
    if result.starting_estate_breakdown and result.starting_estate_breakdown.notes:
        for item in result.starting_estate_breakdown.notes:
            rows.append(["起点財産メモ", "低", item])
    if result.secondary_small_scale_review:
        for item in result.secondary_small_scale_review.notes:
            rows.append(["小宅再判定メモ", "高", item])
        for record in result.secondary_small_scale_review.records:
            rows.append(["小宅再判定事項", "高", f"{record.land_name}: {record.action_required}"])
    if not rows:
        rows.append(["監査メモ", "低", "重大な追加注記はありません。"])
    return pd.DataFrame(rows, columns=["分類", "優先度", "内容"])


def build_successive_inheritance_credit_df(result: SecondaryResult) -> pd.DataFrame:
    computation = result.successive_inheritance_computation
    rows: list[list[str]] = []
    if computation is None or not computation.records:
        rows.append(["相次相続控除", "対象なし", "", "", "控除明細はありません。", ""])
    else:
        for record in computation.records:
            rows.append([
                record.heir_name,
                fmt_pct(record.share_factor),
                fmt_int(record.gross_credit),
                fmt_int(record.limited_credit),
                " / ".join(record.notes),
                "法定相続分ベース按分",
            ])
    return pd.DataFrame(rows, columns=["相続人", "按分比率", "按分前控除額", "反映控除額", "注記", "備考"])


def build_secondary_small_scale_review_df(result: SecondaryResult) -> pd.DataFrame:
    rows: list[list[str]] = []
    review = result.secondary_small_scale_review
    if review is None or not review.records:
        rows.append(["小宅再判定", "対象なし", "", "", "再判定対象となる宅地入力はありません。", ""])
    else:
        for record in review.records:
            rows.append([
                record.land_name,
                record.status,
                record.acquirer_name,
                record.action_required,
                record.reason,
                " / ".join(record.notes),
            ])
    return pd.DataFrame(rows, columns=["対象宅地", "状態", "一次取得者", "再判定アクション", "理由", "注記"])


def build_secondary_detail_df(result: SecondaryResult) -> pd.DataFrame:
    breakdown_note = ""
    if result.starting_estate_breakdown and result.starting_estate_breakdown.notes:
        breakdown_note = " / ".join(result.starting_estate_breakdown.notes[:2])
    rejudge_note = ""
    if result.resolved_secondary_heirs:
        rejudge_note = f"再判定相続人{len(result.resolved_secondary_heirs)}名"
    adjustment_note = " / ".join(result.tax_adjustment_notes[:2]) if result.tax_adjustment_notes else ""
    return pd.DataFrame(
        [
            ["1", "一次からの純承継分", fmt_int(result.net_acq_s), f"配偶者取得{int(result.ratio_s * 100)}%時"],
            ["2", "配偶者固有財産", fmt_int(result.s_own), ""],
            ["3", "生活費・支出等控除", f"△{fmt_int(result.s_spend_total)}", breakdown_note],
            ["4", "【二次相続 課税価格】", fmt_int(result.tax_p_2), rejudge_note],
            ["5", "二次基礎控除額", f"△{fmt_int(result.basic_2)}", f"相続人{result.c_count_2}名"],
            ["6", "二次相続税（調整前）", fmt_int(result.preliminary_total_tax_2), "概算"],
            ["7", "相次相続控除（接続精緻化版）", f"△{fmt_int(result.successive_inheritance_credit)}", adjustment_note],
            ["8", "未成年者控除", f"△{fmt_int(getattr(result, 'minor_credit', Decimal('0')))}", "二次相続時点で判定"],
            ["9", "障害者控除", f"△{fmt_int(getattr(result, 'disability_credit', Decimal('0')))}", "二次相続時点で判定"],
            ["10", "【二次相続税 総額】", fmt_int(result.total_tax_2), "概算・税額調整後"],
        ],
        columns=["No", "項目", "金額", "備考"],
    )


def build_simulation_df(primary_inputs: PrimaryInputs, primary_result: PrimaryResult, secondary_inputs: SecondaryInputs) -> pd.DataFrame:
    sim_results: list[dict[str, Any]] = []
    heirs_for_second = [h for h in primary_inputs.heirs_info if h["type"] in [HEIR_TYPE_CHILD, HEIR_TYPE_GRANDCHILD]]
    heirs_for_second = heirs_for_second if heirs_for_second else primary_inputs.heirs_info
    ratio_candidates = range(0, 101, 10) if primary_inputs.has_spouse else [0]

    for i in ratio_candidates:
        sim_secondary_inputs = SecondaryInputs(
            spouse_acquisition_pct=i,
            s_own=secondary_inputs.s_own,
            annual_spend=secondary_inputs.annual_spend,
            interval_years=secondary_inputs.interval_years,
            use_individual_allocations=secondary_inputs.use_individual_allocations,
            actual_acquisition_inputs=build_simulation_allocation_inputs(
                total_taxable_price=primary_result.tax_p,
                current_inputs=secondary_inputs.actual_acquisition_inputs,
                has_spouse=primary_inputs.has_spouse,
                heirs_info=primary_inputs.heirs_info,
                spouse_acquisition_pct=i,
            ) if secondary_inputs.use_individual_allocations else [],
        )
        sim_primary = calculate_primary_inheritance(primary_inputs, sim_secondary_inputs)
        spouse_tax = sim_primary.heir_tax_records[0].final_tax if primary_inputs.has_spouse and sim_primary.heir_tax_records else Decimal("0")
        spouse_net = sim_primary.spouse_actual_taxable_price - spouse_tax
        tp2 = max(Decimal("0"), spouse_net + to_d(secondary_inputs.s_own) - (to_d(secondary_inputs.annual_spend) * to_d(secondary_inputs.interval_years)))
        c_count_2 = len(heirs_for_second)
        basic_2 = BASIC_DEDUCTION_BASE + (BASIC_DEDUCTION_PER_HEIR * to_d(c_count_2))
        t2 = SupremeLegacyEngine.get_tax(max(Decimal("0"), tp2 - basic_2), False, heirs_for_second)
        sim_results.append(
            {
                "配分(%)": f"{i}%",
                "一次相続税額": int(sim_primary.total_final_tax),
                "二次相続税額": int(quantize_yen(t2)),
                "合計納税額": int(quantize_yen(sim_primary.total_final_tax + t2)),
            }
        )
    return pd.DataFrame(sim_results)


# =========================================================
# 7. Output Logic
# =========================================================
def ensure_pdf_font_registered() -> str:
    font_name = "HeiseiKakuGo-W5"
    try:
        registerFont(UnicodeCIDFont(font_name))
    except Exception:
        pass
    return font_name


def _pdf_safe(value: Any) -> str:
    if value is None:
        return "-"
    text = str(value)
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")


def _pdf_plain_money(value: Any) -> str:
    """金額を「#,###円」形式の文字列に変換する。"""
    try:
        amount = int(Decimal(str(value)).quantize(Decimal("1"), ROUND_HALF_UP))
    except Exception:
        try:
            amount = int(float(value))
        except Exception:
            return str(value)
    return f"{amount:,}円"


def _pdf_money_display(value: Any) -> str:
    """金額を億・万円単位の短縮形式で返す（PDF用）。"""
    try:
        amount = int(Decimal(str(value)).quantize(Decimal("1"), ROUND_HALF_UP))
    except Exception:
        try:
            amount = int(float(value))
        except Exception:
            return str(value)
    abs_amount = abs(amount)
    sign = "△" if amount < 0 else ""
    if abs_amount >= 100_000_000:
        return f"{sign}{abs_amount / 100_000_000:.2f}億円"
    if abs_amount >= 10_000:
        return f"{sign}{abs_amount / 10_000:.1f}万円"
    return f"{sign}{abs_amount:,}円"


def _pdf_build_highlight_box(
    label: str,
    value: str,
    caption: str,
    width_mm: float,
    label_style: Any,
    value_style: Any,
    caption_style: Any,
) -> Table:
    """PDF用KPIハイライトボックスを生成する。"""
    box_width = width_mm * mm
    inner_data = [
        [Paragraph(_pdf_safe(label), label_style)],
        [Paragraph(_pdf_safe(value), value_style)],
        [Paragraph(_pdf_safe(caption), caption_style)],
    ]
    inner_table = Table(inner_data, colWidths=[box_width - 6 * mm])
    inner_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#EDF5FF")),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#1E5AA8")),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]
        )
    )
    return inner_table


def _pdf_build_ratio_chart(
    df_sim: pd.DataFrame,
    current_ratio: int,
    recommended_ratio: int,
) -> Table:
    """配偶者取得割合別の税額比較を棒グラフ（Drawing）として生成する。"""
    work = _ensure_dataframe(
        df_sim,
        ["配分(%)", "一次相続税額", "二次相続税額", "合計納税額"],
    )
    if work.empty:
        placeholder_style = ParagraphStyle(
            "ChartPlaceholder",
            fontName="HeiseiKakuGo-W5",
            fontSize=9,
            leading=12,
        )
        return Table(
            [[Paragraph("（グラフデータがありません）", placeholder_style)]],
            colWidths=[170 * mm],
        )

    work["配分数値"] = work["配分(%)"].apply(_extract_ratio_int)
    work = work.sort_values("配分数値").reset_index(drop=True)
    n = len(work)

    chart_width = 170 * mm
    chart_height = 60 * mm
    drawing = Drawing(chart_width, chart_height)

    bar_chart = VerticalBarChart()
    bar_chart.x = 30
    bar_chart.y = 20
    bar_chart.width = chart_width - 50
    bar_chart.height = chart_height - 30
    bar_chart.data = [
        [int(v) for v in work["一次相続税額"].tolist()],
        [int(v) for v in work["二次相続税額"].tolist()],
        [int(v) for v in work["合計納税額"].tolist()],
    ]
    bar_chart.categoryAxis.categoryNames = [
        f"{int(v)}%" for v in work["配分数値"].tolist()
    ]
    bar_chart.bars[0].fillColor = colors.HexColor("#5B9BD5")
    bar_chart.bars[1].fillColor = colors.HexColor("#A5A5A5")
    bar_chart.bars[2].fillColor = colors.HexColor("#4472C4")
    bar_chart.groupSpacing = 5
    bar_chart.barSpacing = 2
    bar_chart.valueAxis.forceZero = 1
    bar_chart.categoryAxis.labels.fontSize = 7
    bar_chart.valueAxis.labels.fontSize = 7

    drawing.add(bar_chart)

    # 推奨案・現状想定の縦線マーカー
    bar_total_width = bar_chart.width
    bar_unit = bar_total_width / max(n, 1)
    for idx, row in work.iterrows():
        ratio_val = int(row["配分数値"])
        x_center = bar_chart.x + bar_unit * (idx + 0.5)
        if ratio_val == recommended_ratio:
            line = Line(
                x_center, bar_chart.y,
                x_center, bar_chart.y + bar_chart.height,
                strokeColor=colors.HexColor("#1E5AA8"),
                strokeWidth=1.5,
                strokeDashArray=[4, 2],
            )
            drawing.add(line)
            label = String(
                x_center + 2, bar_chart.y + bar_chart.height - 8,
                f"★推奨 {ratio_val}%",
                fontSize=6.5,
                fillColor=colors.HexColor("#1E5AA8"),
            )
            drawing.add(label)
        elif ratio_val == current_ratio and current_ratio != recommended_ratio:
            line = Line(
                x_center, bar_chart.y,
                x_center, bar_chart.y + bar_chart.height,
                strokeColor=colors.HexColor("#a61d24"),
                strokeWidth=1.2,
                strokeDashArray=[3, 2],
            )
            drawing.add(line)
            label = String(
                x_center + 2, bar_chart.y + bar_chart.height - 16,
                f"■現状 {ratio_val}%",
                fontSize=6.5,
                fillColor=colors.HexColor("#a61d24"),
            )
            drawing.add(label)

    wrapper = Table([[drawing]], colWidths=[chart_width])
    wrapper.setStyle(TableStyle([("TOPPADDING", (0, 0), (-1, -1), 0), ("BOTTOMPADDING", (0, 0), (-1, -1), 0)]))
    return wrapper


def _pdf_recommendation_comparison_df(
    df_sim: pd.DataFrame,
    current_ratio: int,
    recommended_ratio: int,
) -> pd.DataFrame:
    """PDF用の推奨案比較DataFrameを生成する。"""
    work = _ensure_dataframe(
        df_sim,
        ["配分(%)", "一次相続税額", "二次相続税額", "合計納税額"],
    )
    if work.empty:
        return pd.DataFrame(
            columns=["区分", "配偶者割合", "一次相続", "二次相続", "合計税額"]
        )
    work["配分数値"] = work["配分(%)"].apply(_extract_ratio_int)
    rows: list[dict[str, Any]] = []
    for _, row in work.sort_values("合計納税額").head(5).iterrows():
        ratio = _extract_ratio_int(row["配分(%)"])
        if ratio == recommended_ratio:
            label = "推奨案"
        elif ratio == current_ratio:
            label = "現状想定"
        else:
            label = "比較案"
        rows.append(
            {
                "区分": label,
                "配偶者割合": f"{ratio}%",
                "一次相続": _pdf_plain_money(row["一次相続税額"]),
                "二次相続": _pdf_plain_money(row["二次相続税額"]),
                "合計税額": _pdf_plain_money(row["合計納税額"]),
            }
        )
    return pd.DataFrame(rows)


def _find_ratio_row(
    df: pd.DataFrame,
    ratio: int,
) -> Optional[pd.Series]:
    """シミュレーションDataFrameから指定割合に最も近い行を返す。"""
    if df is None or df.empty:
        return None
    work = df.copy()
    if "配分数値" not in work.columns:
        work["配分数値"] = work["配分(%)"].apply(_extract_ratio_int)
    exact = work[work["配分数値"] == ratio]
    if not exact.empty:
        return exact.iloc[0]
    nearest = work.iloc[(work["配分数値"] - ratio).abs().argsort()[:1]]
    return nearest.iloc[0] if not nearest.empty else None


def _ppt_rgb(hex_str: str) -> RGBColor:
    """16進数カラー文字列をRGBColorオブジェクトに変換する。"""
    clean = hex_str.lstrip("#").upper()
    if len(clean) == 6:
        return RGBColor(
            int(clean[0:2], 16),
            int(clean[2:4], 16),
            int(clean[4:6], 16),
        )
    return RGBColor(0, 0, 0)


def _ppt_set_shape_fill(
    shape: Any,
    fill_hex: str,
    line_hex: Optional[str] = None,
) -> None:
    """PowerPointシェイプの塗りつぶし色と枠線色を設定する。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ppt_rgb(fill_hex)
    if line_hex is not None:
        shape.line.color.rgb = _ppt_rgb(line_hex)
    else:
        shape.line.color.rgb = _ppt_rgb(fill_hex)


def _trim_df_for_pdf(df: pd.DataFrame, max_rows: int = 12) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame({"項目": ["データなし"]})
    trimmed = df.copy()
    if len(trimmed) > max_rows:
        trimmed = trimmed.head(max_rows).copy()
        ellipsis_row = {col: "..." for col in trimmed.columns}
        trimmed.loc[len(trimmed)] = ellipsis_row
    return trimmed


def _pdf_text_units(value: Any) -> float:
    text = _pdf_safe(value).replace("<br/>", " ").replace("&nbsp;", " ")
    units = 0.0
    for ch in text:
        if ch in "\r\n\t":
            continue
        units += 1.9 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 1.0
    return max(units, 2.0)


def _pdf_estimate_col_widths(df: pd.DataFrame, max_width: float, min_col_width: float = 18 * mm) -> list[float]:
    columns = list(df.columns)
    if not columns:
        return [max_width]
    weights = []
    for col in columns:
        header_units = _pdf_text_units(col) * 1.15
        body_units = max((_pdf_text_units(v) for v in df[col].head(12).tolist()), default=4.0)
        weights.append(max(header_units, min(body_units, 34.0)))
    total_weight = sum(weights) or len(weights)
    widths = [max(min_col_width, max_width * (w / total_weight)) for w in weights]
    current_total = sum(widths)
    if current_total > max_width:
        scale = max_width / current_total
        widths = [max(min_col_width, w * scale) for w in widths]
        overflow = sum(widths) - max_width
        idx = len(widths) - 1
        while overflow > 0.1 and idx >= 0:
            reducible = max(0.0, widths[idx] - min_col_width)
            delta = min(reducible, overflow)
            widths[idx] -= delta
            overflow -= delta
            idx -= 1
    elif current_total < max_width:
        widths[-1] += max_width - current_total
    return widths


def _build_pdf_table(
    df: pd.DataFrame,
    body_style: ParagraphStyle,
    header_style: ParagraphStyle,
    col_widths: Optional[list[float]] = None,
    max_width: float = 170 * mm,
    body_font_size: Optional[float] = None,
    body_leading: Optional[float] = None,
    cell_padding: tuple[float, float, float, float] = (6, 6, 6, 6),
) -> Table:
    trimmed = _trim_df_for_pdf(df)
    columns = list(trimmed.columns)
    if not columns:
        trimmed = pd.DataFrame({"項目": ["データなし"]})
        columns = list(trimmed.columns)

    resolved_widths = list(col_widths) if col_widths else _pdf_estimate_col_widths(trimmed, max_width=max_width)
    column_count = len(columns)
    if len(resolved_widths) != column_count:
        if len(resolved_widths) < column_count:
            resolved_widths.extend([max_width / max(column_count, 1)] * (column_count - len(resolved_widths)))
        resolved_widths = resolved_widths[:column_count]

    total_width = sum(resolved_widths) or max_width
    if total_width > max_width:
        shrink_ratio = max_width / total_width
        resolved_widths = [max(14 * mm, width * shrink_ratio) for width in resolved_widths]
        current_total = sum(resolved_widths)
        if current_total > max_width:
            resolved_widths[-1] = max(14 * mm, resolved_widths[-1] - (current_total - max_width))
    elif total_width < max_width:
        resolved_widths[-1] += max_width - total_width

    effective_body = ParagraphStyle(
        f"{body_style.name}_PDFTable_{column_count}",
        parent=body_style,
        fontSize=body_font_size if body_font_size is not None else body_style.fontSize,
        leading=body_leading if body_leading is not None else max(body_style.leading, (body_font_size or body_style.fontSize) + 3),
        wordWrap="CJK",
        splitLongWords=True,
        allowWidows=1,
        allowOrphans=1,
    )
    effective_header = ParagraphStyle(
        f"{header_style.name}_PDFTable_{column_count}",
        parent=header_style,
        wordWrap="CJK",
        splitLongWords=True,
        alignment=1,
        leading=max(header_style.leading, (header_style.fontSize or 8.7) + 2.5),
    )

    data = [[Paragraph(_pdf_safe(col), effective_header) for col in columns]]
    for _, row in trimmed.iterrows():
        data.append([Paragraph(_pdf_safe(row[col]), effective_body) for col in columns])

    table = Table(data, repeatRows=1, colWidths=resolved_widths, hAlign='LEFT')
    top_pad, bottom_pad, left_pad, right_pad = cell_padding
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{COLOR_NAVY}")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, -1), effective_body.fontName),
                ("FONTSIZE", (0, 0), (-1, -1), effective_body.fontSize),
                ("LEADING", (0, 0), (-1, -1), effective_body.leading),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#B7C0D0")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
                ("TOPPADDING", (0, 0), (-1, -1), top_pad),
                ("BOTTOMPADDING", (0, 0), (-1, -1), bottom_pad),
                ("LEFTPADDING", (0, 0), (-1, -1), left_pad),
                ("RIGHTPADDING", (0, 0), (-1, -1), right_pad),
            ]
        )
    )
    return table


def _build_pdf_note_box(text_value: str, body_style: ParagraphStyle, box_width: float = 170 * mm) -> Table:
    html_text = text_value.replace("&", "&amp;").replace("<br/>", "[[BR]]")
    html_text = html_text.replace("<", "&lt;").replace(">", "&gt;").replace("[[BR]]", "<br/>")
    note_style = ParagraphStyle(
        f"{body_style.name}_Note",
        parent=body_style,
        wordWrap="CJK",
        splitLongWords=True,
        leading=max(body_style.leading, body_style.fontSize + 4),
        allowWidows=1,
        allowOrphans=1,
    )
    table = Table([[Paragraph(html_text, note_style)]], colWidths=[box_width], hAlign='LEFT')
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#FFF6DD")),
                ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(f"#{COLOR_GOLD}")),
                ("TOPPADDING", (0, 0), (-1, -1), 9),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 9),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ]
        )
    )
    return table


def create_pdf_report(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
    secondary_result: SecondaryResult,
    df_sim: pd.DataFrame,
    df_snapshot_summary: pd.DataFrame,
    df_carryforward: pd.DataFrame,
    df_audit_notes: pd.DataFrame,
    df_small_scale_review: pd.DataFrame,
    df_successive_credit: pd.DataFrame,
) -> bytes:
    font_name = ensure_pdf_font_registered()
    output = BytesIO()
    doc = SimpleDocTemplate(
        output,
        pagesize=A4,
        leftMargin=16 * mm,
        rightMargin=16 * mm,
        topMargin=14 * mm,
        bottomMargin=14 * mm,
        title="相続税シミュレーションご提案レポート",
        author="山根会計",
    )
    sample = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleJP", parent=sample["Title"], fontName=font_name, fontSize=20, leading=24, textColor=colors.HexColor(f"#{COLOR_NAVY}"), spaceAfter=6)
    heading_style = ParagraphStyle("HeadingJP", parent=sample["Heading2"], fontName=font_name, fontSize=14, leading=18, textColor=colors.HexColor(f"#{COLOR_NAVY}"), spaceBefore=2, spaceAfter=6)
    body_style = ParagraphStyle("BodyJP", parent=sample["BodyText"], fontName=font_name, fontSize=9.4, leading=13.2, textColor=colors.black, spaceAfter=4)
    small_style = ParagraphStyle("SmallJP", parent=body_style, fontSize=8.2, leading=11, textColor=colors.HexColor("#666666"))
    header_style = ParagraphStyle("HeaderJP", parent=body_style, fontName=font_name, fontSize=8.7, leading=11, textColor=colors.white)
    summary_label_style = ParagraphStyle("SummaryLabel", parent=body_style, fontName=font_name, fontSize=9.5, leading=11, textColor=colors.HexColor(f"#{COLOR_BLUE}"))
    summary_value_style = ParagraphStyle("SummaryValue", parent=body_style, fontName=font_name, fontSize=18, leading=22, textColor=colors.HexColor(f"#{COLOR_BLUE}"))
    summary_caption_style = ParagraphStyle("SummaryCaption", parent=small_style, fontName=font_name, fontSize=8.2, leading=10.4, textColor=colors.HexColor("#666666"))

    def section(title: str, summary_text: Optional[str] = None) -> list[Any]:
        elems: list[Any] = [Paragraph(title, heading_style)]
        if summary_text:
            elems.append(Paragraph(_pdf_safe(summary_text), body_style))
            elems.append(Spacer(1, 2.5 * mm))
        return elems

    total_assets_amount = (
        to_d(primary_inputs.v_home) + to_d(primary_inputs.v_biz) + to_d(primary_inputs.v_rent)
        + to_d(primary_inputs.v_build) + to_d(primary_inputs.v_stock) + to_d(primary_inputs.v_cash)
        + to_d(primary_inputs.v_ins) + to_d(primary_inputs.v_others) - to_d(primary_inputs.v_debt) - to_d(primary_inputs.v_funeral)
    )
    current_ratio = int(getattr(secondary_inputs, "spouse_acquisition_pct", 0) or 0)
    recommendation = _choose_recommendation_plan(df_sim, current_ratio)
    recommended_ratio = int(recommendation.get("recommended_ratio", current_ratio) or current_ratio)
    recommended_total_tax = _to_int_safe(recommendation.get("recommended_total_tax", 0))
    # recommended_primary_tax / recommended_secondary_tax は
    # recommendation dict 経由で参照するため、ここでは変数に保持しない

    work_sim = _prepare_simulation_dataframe(df_sim)
    current_row = _find_ratio_row(work_sim, current_ratio) if not work_sim.empty else None
    current_total_tax = _to_int_safe(current_row["合計納税額"]) if current_row is not None else recommended_total_tax
    current_primary_tax = _to_int_safe(current_row["一次相続税額"]) if current_row is not None else _to_int_safe(primary_result.total_final_tax)
    current_secondary_tax = _to_int_safe(current_row["二次相続税額"]) if current_row is not None else _to_int_safe(secondary_result.total_tax_2)
    tax_reduction = max(current_total_tax - recommended_total_tax, 0)
    equal_benchmark = _to_int_safe(recommendation.get("equal_total_tax", 0))
    equal_reduction = max(equal_benchmark - recommended_total_tax, 0)

    spouse_acquired_amount = secondary_result.snapshot.spouse_acquired_total_amount if secondary_result.snapshot else Decimal("0")
    second_inheritance_date = secondary_result.context.second_inheritance_date if secondary_result.context else date.today()
    secondary_starting_estate = secondary_result.starting_estate_breakdown.final_secondary_starting_estate if secondary_result.starting_estate_breakdown else Decimal("0")

    risk_lines: list[str] = []
    if df_audit_notes is not None and not df_audit_notes.empty:
        level_col = "区分" if "区分" in df_audit_notes.columns else ("分類" if "分類" in df_audit_notes.columns else None)
        item_col = "項目" if "項目" in df_audit_notes.columns else None
        desc_col = "内容" if "内容" in df_audit_notes.columns else ("コメント" if "コメント" in df_audit_notes.columns else item_col)
        for _, row in df_audit_notes.head(6).iterrows():
            parts = []
            if level_col:
                parts.append(str(row.get(level_col, "")).strip())
            if item_col:
                parts.append(str(row.get(item_col, "")).strip())
            if desc_col:
                parts.append(str(row.get(desc_col, "")).strip())
            merged = " / ".join([p for p in parts if p and p != "nan"])
            if merged:
                risk_lines.append(merged)
    if not risk_lines:
        risk_lines = [
            "自宅や預金の分け方を先に整理しておくと、ご家族の話し合いを進めやすくなります。",
            "納税資金の確保方法を確認しておくと、不動産売却を急がずに進めやすくなります。",
            "二次相続まで含めて確認することで、今回の判断による将来差額を把握できます。",
        ]

    assumptions_df = pd.DataFrame(
        [
            {"項目": "相続人構成", "内容": f"配偶者: {'あり' if primary_inputs.has_spouse else 'なし'} / 子等: {primary_inputs.heir_count}人"},
            {"項目": "現状の配偶者取得割合", "内容": f"{current_ratio}%"},
            {"項目": "推奨する配偶者取得割合", "内容": f"{recommended_ratio}%"},
            {"項目": "総財産額（概算）", "内容": _pdf_plain_money(total_assets_amount)},
            {"項目": "二次相続の想定時点", "内容": second_inheritance_date.isoformat()},
            {"項目": "配偶者固有財産", "内容": _pdf_plain_money(secondary_inputs.s_own)},
        ]
    )

    current_analysis_df = pd.DataFrame(
        [
            {"項目": "現状想定の総税額", "内容": _pdf_plain_money(current_total_tax), "意味": "一次相続と二次相続を合算した目安です"},
            {"項目": "今回の相続税額", "内容": _pdf_plain_money(current_primary_tax), "意味": "今回発生する税金の目安です"},
            {"項目": "将来の相続税額", "内容": _pdf_plain_money(current_secondary_tax), "意味": "配偶者様の将来の相続まで含めた目安です"},
            {"項目": "配偶者取得額", "内容": _pdf_plain_money(spouse_acquired_amount) if primary_inputs.has_spouse else "-", "意味": "生活資金と二次相続に影響する主要項目です"},
            {"項目": "二次相続の起点財産", "内容": _pdf_plain_money(secondary_starting_estate), "意味": "配偶者様の税引後残高や固有財産を加味しています"},
        ]
    )

    comparison_df = _pdf_recommendation_comparison_df(df_sim, current_ratio, recommended_ratio)

    compare_rows: list[dict[str, str]] = []
    if not work_sim.empty:
        work_sim["配分数値"] = work_sim["配分(%)"].apply(_extract_ratio_int)
        for _, row in work_sim.sort_values("合計納税額").head(5).iterrows():
            ratio = _extract_ratio_int(row["配分(%)"])
            label = "推奨案" if ratio == recommended_ratio else ("現状想定" if ratio == current_ratio else "比較案")
            compare_rows.append(
                {
                    "区分": label,
                    "配偶者割合": f"{ratio}%",
                    "一次相続": _pdf_plain_money(row["一次相続税額"]),
                    "二次相続": _pdf_plain_money(row["二次相続税額"]),
                    "合計税額": _pdf_plain_money(row["合計納税額"]),
                }
            )
    compare_options_df = pd.DataFrame(compare_rows) if compare_rows else comparison_df.copy()

    recommendation_df = pd.DataFrame(
        [
            {"項目": "推奨案", "内容": f"配偶者 {recommended_ratio}% 取得案"},
            {"項目": "推奨理由", "内容": recommendation.get("recommended_reason", "税額とご家族の進めやすさを踏まえた総合案です。")},
            {"項目": "想定総税額", "内容": _pdf_plain_money(recommended_total_tax)},
            {"項目": "現状想定との差額", "内容": _pdf_plain_money(tax_reduction)},
            {"項目": "均等分割に近い案との差額", "内容": _pdf_plain_money(equal_reduction)},
            {"項目": "納税資金への影響", "内容": "配偶者様の生活資金と納税資金の両立を検討しやすい水準です。"},
            {"項目": "ご家族への配分上の留意点", "内容": "自宅・預金・保険の分け方をあわせて整理すると、円満な承継につながりやすくなります。"},
        ]
    )

    family_merit_df = pd.DataFrame(
        [
            {"対象": "配偶者様", "メリット": "生活資金を確保しやすく、住まいの方針も整理しやすくなります。"},
            {"対象": "お子様", "メリット": "二次相続まで含めた負担差を先に把握でき、ご家族で話し合いやすくなります。"},
            {"対象": "ご家族全体", "メリット": "税金だけでなく、納税資金・分けやすさ・円満な承継まで含めて整理できます。"},
        ]
    )

    next_steps_df = pd.DataFrame(
        [
            {"STEP": "1", "内容": "不動産・預金・保険・遺言の資料を確認し、前提条件を固めます。"},
            {"STEP": "2", "内容": "土地評価や特例適用の可否を含め、正式シミュレーションへ進みます。"},
            {"STEP": "3", "内容": "ご家族向け説明資料と実行プランを整理し、具体的な対策案を決定します。"},
        ]
    )

    note_text = (
        "本レポートは、現時点でご提示いただいた資料・条件に基づく試算です。正式申告や実行時には、"
        "土地評価、各種特例の適用可否、遺産分割内容、相続人構成、追加資料の有無等により金額が変動する場合があります。"
    )

    story: list[Any] = []

    # 1. 総合結論サマリー
    story.append(Paragraph("相続税シミュレーション ご提案レポート", title_style))
    story.append(Paragraph(f"現時点では、配偶者 {recommended_ratio}% 取得案が最もご提案しやすい候補です。", heading_style))
    story.append(Paragraph("税額だけでなく、生活資金、納税資金、ご家族間の分けやすさまで含めて整理しています。", body_style))
    story.append(Spacer(1, 3 * mm))
    kpi_table = Table(
        [[
            _pdf_build_highlight_box("推奨案", f"配偶者 {recommended_ratio}%", "ご家族に説明しやすい基準案", 52, summary_label_style, summary_value_style, summary_caption_style),
            _pdf_build_highlight_box("想定総税額", _pdf_money_display(recommended_total_tax), "一次相続 + 二次相続の合計目安", 52, summary_label_style, summary_value_style, summary_caption_style),
            _pdf_build_highlight_box("現状との差額", _pdf_money_display(tax_reduction), f"現状 {current_ratio}% 想定との比較", 52, summary_label_style, summary_value_style, summary_caption_style),
        ]],
        colWidths=[56 * mm, 56 * mm, 56 * mm],
    )
    kpi_table.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP")]))
    story.append(kpi_table)
    story.append(Spacer(1, 4 * mm))
    summary_points = pd.DataFrame(
        [
            {"要点": "結論", "内容": f"現時点では配偶者 {recommended_ratio}% 取得案を推奨します。"},
            {"要点": "比較結果", "内容": f"現状想定より {tax_reduction:,}円、均等分割に近い案より {equal_reduction:,}円の改善余地があります。"},
            {"要点": "推奨理由", "内容": recommendation.get("recommended_reason", "税額とご家族の進めやすさを踏まえた総合案です。")},
            {"要点": "次の一歩", "内容": "正式シミュレーションでは、不動産評価や各種特例を確認し、精度を高めます。"},
        ]
    )
    story.append(_build_pdf_table(summary_points, body_style, header_style, [34 * mm, 136 * mm]))
    story.append(Spacer(1, 3 * mm))
    story.append(_build_pdf_note_box(note_text, body_style))
    story.append(PageBreak())

    # 2. 現状分析
    story.extend(section("1. 現状分析", "まずは現状の相続税の見通しと、ご家族の前提条件を整理します。"))
    story.append(_build_pdf_table(current_analysis_df, body_style, header_style, [42 * mm, 42 * mm, 86 * mm]))
    story.append(Spacer(1, 4 * mm))
    story.append(_build_pdf_table(assumptions_df, body_style, header_style, [50 * mm, 120 * mm]))
    story.append(PageBreak())

    # 3. 課題・留意点
    story.extend(section("2. 課題・留意点", "相続税額だけでなく、納税資金、家族間配分、将来への備えまで含めて確認します。"))
    issue_df = pd.DataFrame(
        [
            {"観点": "納税資金", "内容": "現金化しやすい資産と納税のタイミングを確認しておくと、実行時の負担が軽くなります。"},
            {"観点": "ご家族間配分", "内容": "自宅・預金・保険の分け方を整理すると、話し合いが進めやすくなります。"},
            {"観点": "将来への備え", "内容": "今回だけでなく二次相続まで含めて確認することで、将来差額を把握できます。"},
        ]
    )
    story.append(_build_pdf_table(issue_df, body_style, header_style, [35 * mm, 135 * mm]))
    story.append(Spacer(1, 4 * mm))
    risk_df = pd.DataFrame({"留意点": risk_lines[:5]})
    story.append(_build_pdf_table(risk_df, body_style, header_style, [170 * mm]))
    story.append(PageBreak())

    # 4. 対策案比較
    story.extend(section("3. 対策案比較", "配偶者取得割合を変えた場合の税額差を比較し、現状想定と推奨案の違いを確認します。"))
    story.append(_build_pdf_table(compare_options_df, body_style, header_style, [24 * mm, 28 * mm, 38 * mm, 38 * mm, 42 * mm]))
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph("配偶者取得割合別 税額比較グラフ", heading_style))
    story.append(_pdf_build_ratio_chart(df_sim, current_ratio, recommended_ratio))
    story.append(Spacer(1, 3 * mm))
    story.append(Paragraph("青の縦線が推奨案、赤の縦線が現状想定です。一次相続と二次相続の合計額だけでなく、その内訳も確認できます。", small_style))
    story.append(PageBreak())

    # 5. 最適案の提示
    story.extend(section("4. 最適案の提示", "ご家族にとって納得しやすく、税額面でも説明しやすい案を現時点の推奨案として整理します。"))
    story.append(_build_pdf_table(recommendation_df, body_style, header_style, [46 * mm, 124 * mm]))
    story.append(Spacer(1, 4 * mm))
    story.append(_build_pdf_table(comparison_df, body_style, header_style, [36 * mm, 28 * mm, 35 * mm, 35 * mm, 36 * mm]))
    story.append(PageBreak())

    # 6. 今後のアクション
    story.extend(section("5. 今後のアクション", "次回面談では、正式シミュレーションと実行プランの整理に進むことをおすすめします。"))
    story.append(_build_pdf_table(family_merit_df, body_style, header_style, [32 * mm, 138 * mm]))
    story.append(Spacer(1, 4 * mm))
    story.append(_build_pdf_table(next_steps_df, body_style, header_style, [18 * mm, 152 * mm]))
    story.append(Spacer(1, 4 * mm))
    story.append(_build_pdf_note_box(
        "次回面談では、不動産評価資料、保険契約内容、遺言書の有無、登記内容を確認し、より精度の高いご提案に進めます。",
        body_style,
    ))
    story.append(PageBreak())

    # 7. 前提条件・注意事項
    story.extend(section("6. 前提条件・注意事項", "現行法令等を前提に試算していますが、個別事情により最終結果は変動する場合があります。"))
    story.append(_build_pdf_table(_trim_df_for_pdf(df_snapshot_summary, 10), body_style, header_style))
    story.append(Spacer(1, 4 * mm))
    if df_carryforward is not None and not df_carryforward.empty:
        story.append(Paragraph("二次相続試算への引継状況", heading_style))
        story.append(_build_pdf_table(
            _trim_df_for_pdf(df_carryforward, 8),
            body_style,
            header_style,
            max_width=160 * mm,
            body_font_size=8.0,
            body_leading=10.4,
            cell_padding=(7.0, 7.0, 8.5, 8.5),
        ))
        story.append(Spacer(1, 4 * mm))
    if df_successive_credit is not None and not df_successive_credit.empty:
        story.append(Paragraph("相次相続控除等の反映状況", heading_style))
        story.append(_build_pdf_table(_trim_df_for_pdf(df_successive_credit, 8), body_style, header_style))
        story.append(Spacer(1, 4 * mm))
    if df_small_scale_review is not None and not df_small_scale_review.empty:
        story.append(Paragraph("特例適用等の確認事項", heading_style))
        story.append(_build_pdf_table(_trim_df_for_pdf(df_small_scale_review, 8), body_style, header_style))
        story.append(Spacer(1, 4 * mm))
    story.append(_build_pdf_note_box(
        "本レポートは現時点の試算結果です。正式申告・実行時には、土地評価、各種特例の適用、遺産分割内容、追加資料の有無等に応じて金額が変動する場合があります。最終提案時には個別事情を確認のうえご提案します。",
        body_style,
        box_width=166 * mm,
    ))

    doc.build(story)
    output.seek(0)
    return output.getvalue()


# =========================================================
# PPT Output Logic
# =========================================================
def _ppt_safe(value: Any) -> str:
    if value is None:
        return "―"
    text = str(value).strip()
    return text if text else "―"


def _ppt_money(value: Any) -> str:
    if value is None:
        return "―"
    if isinstance(value, Decimal):
        return f"{int(value):,}円"
    if isinstance(value, Number):
        return f"{int(value):,}円"
    try:
        return f"{int(value):,}円"
    except Exception:
        return _ppt_safe(value)


def _ppt_add_textbox(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 14, bold: bool = False, color: str = COLOR_NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.08
    font = p.font
    font.name = "Meiryo"
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = _ppt_rgb(color)
    return box


def _ppt_add_note(slide, text: str, left: float, top: float, width: float, height: float):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string('FFF4CC')
    shape.line.color.rgb = RGBColor.from_string(COLOR_GOLD.upper())
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.05
    p.font.name = "Meiryo"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor.from_string('7A5C00')
    return shape


def _ppt_add_table(slide, headers: list[str], rows: list[list[Any]], left: float, top: float, width: float, height: float, font_size: int = 10):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)).table

    def _ppt_text_units_local(value: Any) -> float:
        text_value = _ppt_safe(value)
        units = 0.0
        for ch in text_value:
            units += 1.9 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 1.0
        return max(units, 2.0)

    column_units = []
    for idx, header in enumerate(headers):
        body_max = max((_ppt_text_units_local(row[idx]) for row in rows if idx < len(row)), default=2.0)
        header_units = _ppt_text_units_local(header) * 1.15
        column_units.append(max(header_units, min(body_max, 26.0)))
    total_units = sum(column_units) or len(headers)
    for idx, unit in enumerate(column_units):
        table.columns[idx].width = Inches(width * (unit / total_units))

    base_row_height = max(0.34, min(0.64, height / max(len(rows) + 1, 1)))
    table.rows[0].height = Inches(max(base_row_height, 0.42))
    for ridx, row in enumerate(rows, start=1):
        line_factor = 1
        for val in row:
            line_factor = max(line_factor, _ppt_safe(val).count("\n") + 1)
        table.rows[ridx].height = Inches(min(max(base_row_height + (line_factor - 1) * 0.08, 0.34), 0.82))

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _ppt_rgb(COLOR_NAVY)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        font = p.font
        font.name = "Meiryo"
        font.size = Pt(font_size)
        font.bold = True
        font.color.rgb = _ppt_rgb('FFFFFF')

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = _ppt_safe(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            is_number = isinstance(val, Number) or (isinstance(val, str) and bool(re.fullmatch(r"[\d,.%-]+", val.strip())))
            p.alignment = PP_ALIGN.RIGHT if is_number else PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            font = p.font
            font.name = "Meiryo"
            font.size = Pt(font_size)
            font.bold = False
            font.color.rgb = _ppt_rgb(COLOR_TEXT)
    return table


def _ppt_apply_text_style(text_frame, font_size: int = 14, bold: bool = False, color: str = COLOR_DARK) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    if text_frame.vertical_anchor is None:
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Inches(0.04)
    text_frame.margin_right = Inches(0.04)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.08
        paragraph.font.name = "Meiryo"
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _ppt_rgb(color)
        if not paragraph.runs:
            paragraph.text = paragraph.text
        for run in paragraph.runs:
            run.font.name = "Meiryo"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = _ppt_rgb(color)


def _ppt_add_header_band(slide, title: str, subtitle: Optional[str] = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.75))
    _ppt_set_shape_fill(band, COLOR_DARK, COLOR_DARK)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(8.7), Inches(0.3))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    _ppt_apply_text_style(tf, 24, True, 'FFFFFF')
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(10.5), Inches(0.2))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        _ppt_apply_text_style(stf, 10, False, 'D8E4F4')


def _ppt_add_card(slide, title: str, body: str, left: float, top: float, width: float, height: float, fill: str = 'FFFFFF', line: str = COLOR_LIGHT_GRAY, title_color: str = COLOR_BLUE, body_color: str = COLOR_DARK) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _ppt_set_shape_fill(shape, fill, line)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    _ppt_apply_text_style(tf, 16, True, title_color)
    p2 = tf.add_paragraph()
    p2.text = body
    _ppt_apply_text_style(tf, 13, False, body_color)


def _ppt_add_kpi_card(slide, title: str, value: str, subtext: str, left: float, top: float, width: float, height: float, accent: str = COLOR_BLUE, fill: str = COLOR_LIGHT_BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _ppt_set_shape_fill(shape, fill, accent)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    _ppt_apply_text_style(tf, 12, True, accent)
    p2 = tf.add_paragraph()
    p2.text = value
    _ppt_apply_text_style(tf, 24, True, accent)
    p3 = tf.add_paragraph()
    p3.text = subtext
    _ppt_apply_text_style(tf, 10, False, COLOR_GRAY)


def _ppt_add_bullets(slide, title: str, bullets: list[str], left: float, top: float, width: float, height: float, title_color: str = COLOR_BLUE) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    _ppt_apply_text_style(tf, 16, True, title_color)
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f'• {bullet}'
        _ppt_apply_text_style(tf, 13, False, COLOR_DARK)


def _ppt_find_ratio_row(df_sim: pd.DataFrame, ratio: int) -> Optional[pd.Series]:
    work = _prepare_simulation_dataframe(df_sim)
    if work.empty:
        return None
    work['配分数値'] = work['配分(%)'].apply(_extract_ratio_int)
    work = work.sort_values('配分数値').reset_index(drop=True)
    exact = work[work['配分数値'] == ratio]
    if not exact.empty:
        return exact.iloc[0]
    nearest = work.iloc[(work['配分数値'] - ratio).abs().argsort()[:1]]
    return nearest.iloc[0] if not nearest.empty else None


def _ppt_compact_yen(value: Any) -> str:
    amount = _to_int_safe(value, 0)
    abs_amount = abs(amount)
    sign = '▲' if amount < 0 else ''
    if abs_amount >= 100000000:
        return f'{sign}{abs_amount / 100000000:.2f}億円'
    if abs_amount >= 10000:
        return f'{sign}{abs_amount / 10000:.1f}万円'
    return f'{sign}{abs_amount:,}円'


def _ppt_build_plan_text(primary_inputs: PrimaryInputs, ratio: int) -> tuple[str, str]:
    spouse_text = f'配偶者 {ratio}%' if primary_inputs.has_spouse else '配偶者なし'
    heir_types = [str(h.get('type', '相続人')) for h in primary_inputs.heirs_info]
    if not heir_types:
        return spouse_text, '残りはご家族で分ける前提です'
    unique_types = []
    for heir_type in heir_types:
        if heir_type not in unique_types:
            unique_types.append(heir_type)
    other_ratio = max(0, 100 - ratio)
    if len(unique_types) == 1:
        return spouse_text, f'{unique_types[0]} 合計 {other_ratio}% を分割'
    return spouse_text, f'その他のご家族 合計 {other_ratio}% を分割'


def _ppt_reason_lines(recommendation: dict[str, Any], current_ratio: int) -> list[str]:
    lines = [
        '一次相続と二次相続の合計負担が、全体として抑えやすい案です。',
        '配偶者の生活資金を確保しつつ、お子さま側への将来負担も偏りにくい水準です。',
    ]
    diff_vs_equal = _to_int_safe(recommendation.get('diff_vs_equal', 0), 0)
    if diff_vs_equal < 0:
        lines.append(f'均等分割に近い案と比べて、合計で {_ppt_compact_yen(abs(diff_vs_equal))} 程度の圧縮余地があります。')
    elif diff_vs_equal > 0:
        lines.append(f'均等分割に近い案より税額は {_ppt_compact_yen(diff_vs_equal)} 上がりますが、分けやすさと生活資金を重視した案です。')
    if abs(recommendation.get('recommended_ratio', current_ratio) - current_ratio) >= 20:
        lines.append('現在のお考えから差が大きい場合は、自宅・預金・納税資金の順に確認すると進めやすくなります。')
    return lines[:3]


def _ppt_risk_lines(df_audit_notes: Optional[pd.DataFrame], recommendation: dict[str, Any], current_ratio: int) -> list[str]:
    lines = []
    if current_ratio != recommendation.get('recommended_ratio', current_ratio):
        lines.append('現状の配分のままだと、税負担が重いまま固定される可能性があります。')
    if recommendation.get('recommended_secondary_tax', 0) > recommendation.get('recommended_primary_tax', 0):
        lines.append('配偶者に寄せすぎると、将来の二次相続で負担が膨らみやすくなります。')
    lines.append('不動産の分け方が曖昧なままだと、ご家族の話し合いが長引く原因になります。')
    if df_audit_notes is not None and not df_audit_notes.empty:
        for _, row in df_audit_notes.head(3).iterrows():
            content = str(row.iloc[-1]).strip()
            if content and content not in lines:
                lines.append(_sanitize_customer_text(content))
            if len(lines) >= 3:
                break
    return lines[:3]


def _ppt_add_tax_chart(slide, df_sim: pd.DataFrame, current_ratio: int, recommended_ratio: int, left: float, top: float, width: float, height: float) -> None:
    work = _prepare_simulation_dataframe(df_sim)
    if work.empty:
        _ppt_add_card(slide, '税額比較データ', '比較データが未作成のため、グラフ表示を省略しました。', left, top, width, 1.4)
        return
    work['配分数値'] = work['配分(%)'].apply(_extract_ratio_int)
    work = work.sort_values('配分数値').reset_index(drop=True)

    chart_data = CategoryChartData()
    chart_data.categories = [f"{int(v)}%" for v in work['配分数値'].tolist()]
    chart_data.add_series('一次相続', work['一次相続税額'].astype(float).tolist())
    chart_data.add_series('二次相続', work['二次相続税額'].astype(float).tolist())
    chart_data.add_series('合計', work['合計納税額'].astype(float).tolist())

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)

    series_colors = [COLOR_SKY, 'DADFE8', 'C9D3E6']
    for series, base_color in zip(chart.series, series_colors):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _ppt_rgb(base_color)
        series.format.line.color.rgb = _ppt_rgb(base_color)

    total_series = chart.series[2]
    for idx, point in enumerate(total_series.points):
        ratio_value = int(work.loc[idx, '配分数値'])
        point.format.fill.solid()
        if ratio_value == recommended_ratio:
            point.format.fill.fore_color.rgb = _ppt_rgb(COLOR_BLUE)
            point.format.line.color.rgb = _ppt_rgb(COLOR_BLUE)
        elif ratio_value == current_ratio and current_ratio != recommended_ratio:
            point.format.fill.fore_color.rgb = _ppt_rgb(COLOR_RED)
            point.format.line.color.rgb = _ppt_rgb(COLOR_RED)
        else:
            point.format.fill.fore_color.rgb = _ppt_rgb('BFC8D6')
            point.format.line.color.rgb = _ppt_rgb('BFC8D6')

    total_series.has_data_labels = True
    total_series.data_labels.number_format = '#,##0'
    total_series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    total_series.data_labels.font.size = Pt(9)

    note = f'★推奨: {recommended_ratio}%   ' + (f'■現状想定: {current_ratio}%' if current_ratio != recommended_ratio else '現状想定と推奨案は同一です')
    _ppt_add_textbox(slide, note, left, top + height + 0.1, width, 0.25, 10, False, COLOR_GRAY)


def create_ppt_report(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
    secondary_result: SecondaryResult,
    df_sim: pd.DataFrame,
    df_snapshot_summary: pd.DataFrame,
    df_carryforward: pd.DataFrame,
    df_audit_notes: pd.DataFrame,
    df_small_scale_review: pd.DataFrame,
    df_successive_credit: pd.DataFrame,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    current_ratio = _to_int_safe(getattr(secondary_inputs, 'spouse_acquisition_pct', 50), 50)
    recommendation = _choose_recommendation_plan(df_sim, current_ratio)
    recommended_ratio = _to_int_safe(recommendation.get('recommended_ratio', current_ratio), current_ratio)
    current_row = _ppt_find_ratio_row(df_sim, current_ratio)
    equal_row = _ppt_find_ratio_row(df_sim, 50)

    recommended_total_tax = recommendation.get('recommended_total_tax', 0)
    current_total_tax = _to_int_safe(current_row['合計納税額'], recommended_total_tax) if current_row is not None else recommended_total_tax
    current_primary_tax = _to_int_safe(current_row['一次相続税額'], 0) if current_row is not None else _to_int_safe(primary_result.total_final_tax, 0)
    current_secondary_tax = _to_int_safe(current_row['二次相続税額'], 0) if current_row is not None else _to_int_safe(secondary_result.total_tax_2, 0)
    savings_vs_current = max(0, current_total_tax - _to_int_safe(recommended_total_tax, 0))
    savings_vs_equal = max(0, _to_int_safe(equal_row['合計納税額'], current_total_tax) - _to_int_safe(recommended_total_tax, 0)) if equal_row is not None else max(0, -_to_int_safe(recommendation.get('diff_vs_equal', 0), 0))
    primary_secondary_total = _to_int_safe(recommendation.get('recommended_primary_tax', 0), 0) + _to_int_safe(recommendation.get('recommended_secondary_tax', 0), 0)
    spouse_line, family_line = _ppt_build_plan_text(primary_inputs, recommended_ratio)
    reason_lines = _ppt_reason_lines(recommendation, current_ratio)
    risk_lines = _ppt_risk_lines(df_audit_notes, recommendation, current_ratio)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '相続税の負担を抑え、ご家族が話し合いやすい分け方のご提案', '5分で全体像がつかめるよう、結論から先に整理しています。')
    _ppt_add_textbox(slide, '相続対策シミュレーション ご提案資料', 0.7, 1.2, 7.8, 0.7, 28, True, COLOR_DARK)
    _ppt_add_textbox(slide, '税金・生活資金・ご家族の納得感をまとめて見える化しました。', 0.72, 1.95, 7.4, 0.35, 14, False, COLOR_GRAY)
    _ppt_add_kpi_card(slide, '現時点のおすすめ配分', f'{recommended_ratio}%', '配偶者の取得割合', 0.8, 3.0, 2.3, 1.55)
    _ppt_add_kpi_card(slide, 'おすすめ案の総税額', _ppt_compact_yen(recommended_total_tax), '一次＋二次の合計', 3.35, 3.0, 2.7, 1.55)
    _ppt_add_kpi_card(slide, '見直し効果', _ppt_compact_yen(max(savings_vs_current, savings_vs_equal)), '現状想定・均等案との比較', 6.3, 3.0, 2.9, 1.55)
    _ppt_add_card(slide, 'この資料で分かること', '①どの分け方が有力か\n②どれくらい税額差があるか\n③次に確認すべきこと', 9.55, 2.85, 3.0, 2.1, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_textbox(slide, f'作成日: {date.today().isoformat()} / 山根会計', 0.8, 6.65, 3.5, 0.2, 10, False, COLOR_GRAY)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, 'この資料の目的', '相続税の金額だけでなく、家族で進めやすい形まで整理することが目的です。')
    _ppt_add_card(slide, 'このままだと起こりやすい問題', '税額が高いまま確定する\n自宅や預金の分け方で話し合いが長引く\n配偶者に寄せすぎると二次相続が重くなる', 0.8, 1.2, 5.7, 2.2, fill=COLOR_LIGHT_RED, line=COLOR_RED, title_color=COLOR_RED)
    _ppt_add_card(slide, '今回の資料で確認するポイント', '税金が抑えやすい配分\n配偶者の生活資金確保\nご家族が納得しやすい進め方', 6.85, 1.2, 5.7, 2.2, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_bullets(slide, '今やる理由', [
        '元気なうちに決めておくほど、ご家族の選択肢が広がります。',
        '納税資金や不動産の持ち方を早めに確認すると、慌てずに進められます。',
        '一度整理しておくと、次回面談で正式な提案に進みやすくなります。',
    ], 0.9, 4.05, 11.5, 2.1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '結論', 'まずは最も伝わる結論からご確認ください。')
    _ppt_add_textbox(slide, f'おすすめは 配偶者 {recommended_ratio}% 取得案 です', 0.8, 1.1, 8.0, 0.7, 28, True, COLOR_BLUE)
    _ppt_add_textbox(slide, '税負担・生活資金・将来のバランスを総合的に見て、この案が最も説明しやすい状態です。', 0.82, 1.8, 8.6, 0.35, 13, False, COLOR_GRAY)
    _ppt_add_kpi_card(slide, '想定総税額', _ppt_compact_yen(recommended_total_tax), '一次相続＋二次相続', 0.8, 2.5, 2.7, 1.55)
    _ppt_add_kpi_card(slide, '現状想定との差', _ppt_compact_yen(savings_vs_current), f'現状 {current_ratio}% 想定との比較', 3.7, 2.5, 2.8, 1.55, accent=COLOR_BLUE, fill='F5FAFF')
    _ppt_add_kpi_card(slide, '均等案との差', _ppt_compact_yen(savings_vs_equal), '50%付近との比較', 6.7, 2.5, 2.5, 1.55, accent=COLOR_BLUE, fill='F5FAFF')
    _ppt_add_card(slide, '推奨理由', '\n'.join(reason_lines[:3]), 9.5, 2.35, 3.0, 2.2, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, '次にやること', '不動産・預金・保険契約を確認し、正式シミュレーションへ進みます。', 0.8, 4.65, 11.7, 1.25, fill='FFFFFF', line=COLOR_LIGHT_GRAY, title_color=COLOR_GREEN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '最適分割案', '「誰がどの程度受け取ると、全体の納得感が高いか」を整理しています。')
    _ppt_add_card(slide, 'おすすめ配分', f'{spouse_line}\n{family_line}', 0.8, 1.2, 4.0, 2.0, fill=COLOR_LIGHT_BLUE, line=COLOR_BLUE)
    _ppt_add_card(slide, '現状想定との違い', f'現状想定: 配偶者 {current_ratio}%\nおすすめ案: 配偶者 {recommended_ratio}%\n差額: {_ppt_compact_yen(savings_vs_current)} の改善余地', 5.0, 1.2, 3.8, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, 'ご家族への伝え方', '「税額を下げるため」だけでなく\n「生活資金を守り、将来の負担も抑えるため」\nと説明すると納得されやすくなります。', 9.0, 1.2, 3.4, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    breakdown_rows = []
    if df_carryforward is not None and not df_carryforward.empty:
        use_cols = [c for c in ['相続人', '続柄', '取得総額', '税引後残高'] if c in df_carryforward.columns]
        for _, row in df_carryforward.head(5)[use_cols].iterrows():
            breakdown_rows.append([row.get(c, '―') for c in use_cols])
    _ppt_add_table(slide, ['相続人', '続柄', '取得総額', '税引後残高'] if breakdown_rows else ['項目', '内容'], breakdown_rows if breakdown_rows else [['配分情報', '取得明細は入力後に自動反映されます']], 0.9, 3.75, 11.7, 2.1, 11)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, 'なぜこの案が良いのか', 'ご家族に説明しやすい理由を、3つに絞って整理しています。')
    _ppt_add_card(slide, '① 一次相続', f'今回の相続税は {_ppt_compact_yen(recommendation.get("recommended_primary_tax", 0))} を目安に整理。\n配偶者の税負担を抑えやすい配分です。', 0.8, 1.3, 3.8, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, '② 二次相続', f'将来の相続税は {_ppt_compact_yen(recommendation.get("recommended_secondary_tax", 0))} が目安。\n寄せすぎを避けることで、次の相続も重くなりすぎません。', 4.8, 1.3, 3.8, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, '③ 実務面', '自宅・預金・納税資金を分けて考えやすく、家族会議で説明しやすい案です。', 8.8, 1.3, 3.7, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_bullets(slide, '営業・面談で伝えるポイント', [
        '税額最小だけを追うのではなく、生活資金と分けやすさまで含めて判断しています。',
        'お子さま側の将来負担も見ているため、長期で見た納得感があります。',
        recommendation.get('recommended_reason', '総合的に説明しやすい案です。'),
    ], 0.9, 3.85, 11.5, 2.0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '一次相続・二次相続の比較', '今回だけでなく、将来まで含めた合計額で判断することが大切です。')
    _ppt_add_kpi_card(slide, '一次相続の目安', _ppt_compact_yen(recommendation.get('recommended_primary_tax', 0)), '今回かかる税金', 0.9, 1.45, 2.7, 1.5)
    _ppt_add_kpi_card(slide, '二次相続の目安', _ppt_compact_yen(recommendation.get('recommended_secondary_tax', 0)), '将来かかる税金', 3.85, 1.45, 2.7, 1.5, accent=COLOR_DARK, fill='F4F6F9')
    _ppt_add_kpi_card(slide, '合計の目安', _ppt_compact_yen(primary_secondary_total), 'ご家族全体でみた税額', 6.8, 1.45, 2.7, 1.5)
    _ppt_add_card(slide, 'わかりやすい補足', '一次相続 = 今回の相続\n二次相続 = 将来、配偶者様の相続\n両方を合わせて判断することが大切です。', 9.75, 1.28, 2.7, 1.8, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    comparison_rows = [
        ['現状想定', f'{current_ratio}%', _ppt_compact_yen(current_primary_tax), _ppt_compact_yen(current_secondary_tax), _ppt_compact_yen(current_total_tax)],
        ['おすすめ案', f'{recommended_ratio}%', _ppt_compact_yen(recommendation.get('recommended_primary_tax', 0)), _ppt_compact_yen(recommendation.get('recommended_secondary_tax', 0)), _ppt_compact_yen(recommended_total_tax)],
    ]
    _ppt_add_table(slide, ['案', '配偶者割合', '一次相続', '二次相続', '合計'], comparison_rows, 0.9, 3.55, 11.6, 2.1, 11)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '税額比較グラフ', '配偶者の取得割合を変えると、税額がどう動くかを一目で確認できます。')
    _ppt_add_tax_chart(slide, df_sim, current_ratio, recommended_ratio, 0.8, 1.15, 11.8, 4.6)
    _ppt_add_card(slide, 'グラフの見方', '青の★が推奨案、赤は現状想定です。\n一番下がっている帯だけでなく、その前後も含めてバランスを見ています。', 0.9, 6.0, 11.4, 0.9, fill='FFFFFF', line=COLOR_LIGHT_GRAY, title_color=COLOR_BLUE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '放置した場合の主な注意点', '相続対策は「やるかどうか」より「いつ整理するか」で差が出ます。')
    _ppt_add_card(slide, '争族の注意点', risk_lines[0] if len(risk_lines) > 0 else '自宅や預金の分け方が曖昧だと、ご家族の話し合いが長引きやすくなります。', 0.8, 1.35, 3.75, 2.0, fill=COLOR_LIGHT_RED, line=COLOR_RED, title_color=COLOR_RED)
    _ppt_add_card(slide, '納税資金の注意点', risk_lines[1] if len(risk_lines) > 1 else '納税のための現金を確保できないと、不動産の売却を急ぐことがあります。', 4.8, 1.35, 3.75, 2.0, fill='FFF8F8', line='E9B8BB', title_color=COLOR_RED)
    _ppt_add_card(slide, '二次相続の注意点', risk_lines[2] if len(risk_lines) > 2 else '今回だけを見て決めると、将来の二次相続で税額が増えることがあります。', 8.8, 1.35, 3.75, 2.0, fill='FFF8F8', line='E9B8BB', title_color=COLOR_RED)
    _ppt_add_bullets(slide, '今のうちに確認しておきたい項目', [
        '自宅を誰が取得するか',
        '預金の分け方と納税資金',
        '保険契約・遺言書・不動産登記の内容',
    ], 0.9, 4.0, 11.3, 1.8, title_color=COLOR_RED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '今後の進め方', '面談後に何をすればよいかを、3ステップで整理しています。')
    _ppt_add_card(slide, 'STEP 1', '資料の確認\n不動産・預金・保険・遺言の有無を整理します。', 0.85, 1.8, 3.7, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY, title_color=COLOR_BLUE)
    _ppt_add_card(slide, 'STEP 2', '正式シミュレーション\n評価資料を反映し、より精度の高い数字に整えます。', 4.8, 1.8, 3.7, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY, title_color=COLOR_BLUE)
    _ppt_add_card(slide, 'STEP 3', '実行プランの決定\n分割案・遺言・納税資金対策まで具体化します。', 8.75, 1.8, 3.7, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY, title_color=COLOR_BLUE)
    _ppt_add_card(slide, '次回面談でできること', '① 正式版の税額比較\n② ご家族向け説明資料の作成\n③ 実行時の注意点整理', 0.9, 4.45, 11.5, 1.55, fill=COLOR_LIGHT_BLUE, line=COLOR_BLUE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_header_band(slide, '個別相談のご案内', '「この人に相談したい」と感じていただけるよう、次の一歩を明確にしています。')
    _ppt_add_textbox(slide, '次回は、正式シミュレーションと実行プランのご相談をおすすめします。', 0.8, 1.2, 10.8, 0.5, 22, True, COLOR_BLUE)
    _ppt_add_card(slide, '次回面談で確認すること', 'ご自宅の持ち方\n預金の分け方\n保険・遺言・登記の確認', 0.9, 2.0, 3.9, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, 'ご相談いただくメリット', '税金だけでなく、ご家族の納得感まで含めて整理できます。\n「今やるべきこと」が明確になります。', 5.0, 2.0, 3.9, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_card(slide, 'この資料の位置づけ', '現時点の資料・条件に基づくご提案です。\n正式版では評価資料を確認し、精度をさらに高めます。', 9.1, 2.0, 3.0, 2.0, fill='FFFFFF', line=COLOR_LIGHT_GRAY)
    _ppt_add_kpi_card(slide, '今回の見直し効果', _ppt_compact_yen(max(savings_vs_current, savings_vs_equal)), '相談の入口として十分な差額です', 0.95, 4.55, 3.2, 1.5)
    _ppt_add_card(slide, 'ご案内', '資料確認後、正式シミュレーションの面談日程をご相談ください。\n次回はご家族向けの説明資料までご用意できます。', 4.4, 4.45, 8.0, 1.6, fill=COLOR_LIGHT_BLUE, line=COLOR_BLUE, title_color=COLOR_BLUE)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def _to_int_safe(value: Any, default: int = 0) -> int:
    try:
        if pd.isna(value):
            return default
    except Exception:
        pass
    try:
        return int(Decimal(str(value)).quantize(Decimal("1"), ROUND_HALF_UP))
    except Exception:
        try:
            return int(float(value))
        except Exception:
            return default


def _extract_ratio_int(value: Any) -> int:
    text_value = str(value).replace("%", "").strip()
    if not text_value:
        return 0
    return _to_int_safe(text_value, 0)


def _yen_text(value: Any) -> str:
    amount = _to_int_safe(value, 0)
    sign = "△" if amount < 0 else ""
    return f"{sign}{abs(amount):,}円"


def _ensure_dataframe(df: Optional[pd.DataFrame], columns: list[str]) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame(columns=columns)
    return df.copy()


def _drop_duplicate_columns_keep_first(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return _ensure_dataframe(df, [])
    if not df.columns.duplicated().any():
        return df.copy()
    return df.loc[:, ~df.columns.duplicated()].copy()


def _upsert_column(df: pd.DataFrame, index: int, column_name: str, value: Any) -> pd.DataFrame:
    work = _ensure_dataframe(df, [])
    if column_name in work.columns:
        work[column_name] = value
        ordered_columns = [col for col in work.columns if col != column_name]
        safe_index = max(0, min(index, len(ordered_columns)))
        ordered_columns.insert(safe_index, column_name)
        return work.loc[:, ordered_columns]
    safe_index = max(0, min(index, len(work.columns)))
    work.insert(safe_index, column_name, value)
    return work


def _prepare_simulation_dataframe(df_sim: Optional[pd.DataFrame]) -> pd.DataFrame:
    work = _ensure_dataframe(df_sim, ['配分(%)', '一次相続税額', '二次相続税額', '合計納税額'])
    required_columns = ['配分(%)', '一次相続税額', '二次相続税額', '合計納税額']
    for column_name in required_columns:
        if column_name not in work.columns:
            work[column_name] = 0
    work['配分(%)'] = work['配分(%)'].apply(_extract_ratio_int)
    for col in ['一次相続税額', '二次相続税額', '合計納税額']:
        work[col] = pd.to_numeric(work[col], errors='coerce').fillna(0)
    work['配分表示'] = work['配分(%)'].apply(lambda v: f"{_extract_ratio_int(v)}%")
    return work


def _normalize_customer_sheet_name(name: str) -> str:
    replacements = {
        "内部確認": "ご提案用",
        "要確認": "確認事項",
        "carry_forward": "二次相続試算用の引継財産",
        "監査メモ": "確認事項一覧",
        "snapshot": "要点整理",
        "小宅判定": "小規模宅地等の特例 判定結果",
        "小宅再判定": "小規模宅地等の特例 再確認事項",
    }
    normalized = name
    for before, after in replacements.items():
        normalized = normalized.replace(before, after)
    return normalized


def _sanitize_customer_text(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (int, float, Decimal)):
        return value
    text_value = str(value)
    replacements = {
        "内部確認用": "ご提案資料用",
        "内部確認": "ご提案用",
        "要確認": "今後の確認事項",
        "snapshot": "要点整理",
        "carry forward": "二次相続試算用の引継財産",
        "carry_forward": "二次相続試算用の引継財産",
        "監査メモ": "確認事項一覧",
        "内部ロジック": "計算前提",
        "リスク": "確認事項",
        "不明": "未取得資料あり",
    }
    for before, after in replacements.items():
        text_value = text_value.replace(before, after)
    return text_value


def _customerize_dataframe(df: Optional[pd.DataFrame], column_map: dict[str, str], drop_columns: Optional[list[str]] = None) -> pd.DataFrame:
    work = _ensure_dataframe(df, list(column_map.keys()) if column_map else [])
    if column_map:
        rename_targets = {col: column_map.get(col, col) for col in work.columns}
        work = work.rename(columns=rename_targets)
        work = _drop_duplicate_columns_keep_first(work)
    if drop_columns:
        keep_cols = [col for col in work.columns if col not in drop_columns]
        work = work[keep_cols]
    map_method = getattr(work, "map", None)
    if callable(map_method):
        work = map_method(_sanitize_customer_text)
    else:
        work = work.apply(lambda col: col.map(_sanitize_customer_text))
    return work


def _choose_recommendation_plan(df_sim: pd.DataFrame, current_ratio: int) -> dict[str, Any]:
    work = _prepare_simulation_dataframe(df_sim)
    if work.empty:
        return {
            "recommended_ratio": current_ratio,
            "min_tax_ratio": current_ratio,
            "recommended_total_tax": 0,
            "recommended_primary_tax": 0,
            "recommended_secondary_tax": 0,
            "min_total_tax": 0,
            "equal_total_tax": 0,
            "diff_vs_min": 0,
            "diff_vs_equal": 0,
            "practical_ratio": current_ratio,
            "recommended_reason": "比較対象データがないため、入力いただいた配分条件を前提に整理しています。",
            "is_same_as_min": True,
            "is_same_as_practical": True,
        }

    work["配分数値"] = work["配分(%)"].apply(_extract_ratio_int)
    work = work.sort_values("配分数値").reset_index(drop=True)
    min_idx = int(work["合計納税額"].idxmin())
    min_row = work.loc[min_idx]
    min_total_tax = _to_int_safe(min_row["合計納税額"])
    min_ratio = _extract_ratio_int(min_row["配分(%)"])

    equal_candidates = work.iloc[(work["配分数値"] - 50).abs().argsort()[:1]]
    equal_row = equal_candidates.iloc[0]
    equal_total_tax = _to_int_safe(equal_row["合計納税額"])

    tolerance = max(int(min_total_tax * 0.03), 1000000)
    candidate_mask = (
        (work["合計納税額"] <= min_total_tax + tolerance)
        & (work["配分数値"] >= 30)
        & (work["配分数値"] <= 70)
    )
    candidates = work[candidate_mask].copy()
    if candidates.empty:
        candidates = work.copy()
    candidates["実務スコア"] = (
        (candidates["配分数値"] - 50).abs() * 1.3
        + ((candidates["合計納税額"] - min_total_tax) / max(min_total_tax, 1)) * 100
    )
    practical_row = candidates.sort_values(["実務スコア", "配分数値"]).iloc[0]
    practical_ratio = _extract_ratio_int(practical_row["配分(%)"])

    recommended_row = practical_row
    recommended_ratio = practical_ratio
    if abs(practical_ratio - current_ratio) <= 10 and int(practical_row["合計納税額"]) <= min_total_tax + tolerance:
        recommended_reason = "税額だけでなく、配偶者の生活資金確保、今後の分けやすさ、二次相続のバランスを踏まえて、現状の配分意向に近い範囲で整理した案です。"
    elif practical_ratio != min_ratio:
        recommended_reason = "合計税額が最も低い案と比べても差額が大きくなく、配偶者の生活資金、自宅の持ち方、二次相続まで含めた納得感を重視した総合案です。"
    else:
        recommended_reason = "税額面でも実務面でもバランスがよく、現時点では総合的に説明しやすい案です。"

    return {
        "recommended_ratio": recommended_ratio,
        "min_tax_ratio": min_ratio,
        "recommended_total_tax": _to_int_safe(recommended_row["合計納税額"]),
        "recommended_primary_tax": _to_int_safe(recommended_row["一次相続税額"]),
        "recommended_secondary_tax": _to_int_safe(recommended_row["二次相続税額"]),
        "min_total_tax": min_total_tax,
        "equal_total_tax": equal_total_tax,
        "diff_vs_min": _to_int_safe(recommended_row["合計納税額"]) - min_total_tax,
        "diff_vs_equal": _to_int_safe(recommended_row["合計納税額"]) - equal_total_tax,
        "practical_ratio": practical_ratio,
        "recommended_reason": recommended_reason,
        "is_same_as_min": recommended_ratio == min_ratio,
        "is_same_as_practical": recommended_ratio == practical_ratio,
    }


def _build_summary_sheet_df(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
    secondary_result: SecondaryResult,
    df_sim: pd.DataFrame,
    df_audit_notes: pd.DataFrame,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    recommendation = _choose_recommendation_plan(df_sim, secondary_inputs.spouse_acquisition_pct)
    diff_vs_min = recommendation["diff_vs_min"]
    diff_vs_equal = recommendation["diff_vs_equal"]

    comparison_comment = "均等分割に近い案と比べて"
    comparison_delta = diff_vs_equal
    if comparison_delta == 0:
        comparison_sentence = "均等分割に近い案との比較でも大きな差はありません。"
    elif comparison_delta < 0:
        comparison_sentence = f"{comparison_comment} 合計税額が {abs(comparison_delta):,}円低く、税負担の圧縮効果があります。"
    else:
        comparison_sentence = f"{comparison_comment} 合計税額が {abs(comparison_delta):,}円高い一方、分けやすさや生活資金の観点を織り込んだ案です。"

    if diff_vs_min == 0:
        min_comparison_sentence = "この案は合計税額が最も低い案と同水準です。"
    else:
        min_comparison_sentence = f"税額が最も低い案と比べると、合計税額は {abs(diff_vs_min):,}円高くなります。"

    priority_high = []
    priority_mid = []
    priority_low = []
    for _, row in _ensure_dataframe(df_audit_notes, ["分類", "優先度", "内容"]).iterrows():
        content = _sanitize_customer_text(row.get("内容", ""))
        priority = str(row.get("優先度", "")).strip()
        if not content:
            continue
        if priority == "高":
            priority_high.append(content)
        elif priority == "中":
            priority_mid.append(content)
        else:
            priority_low.append(content)

    if not priority_high:
        priority_high = ["自宅をどなたが取得するか、遺産分割の方向性を確認します。"]
    if not priority_mid:
        priority_mid = ["預金の配分方針と、配偶者の今後の生活資金見込みを確認します。"]
    if not priority_low:
        priority_low = ["保険受取人、登記、名義の整備状況を順次確認します。"]

    summary_df = pd.DataFrame(
        [
            ["現時点推奨案", f"配偶者取得割合 {recommendation['recommended_ratio']}% 案"],
            ["税額が最も低い案", f"配偶者取得割合 {recommendation['min_tax_ratio']}% 案"],
            ["実務上の推奨案", f"配偶者取得割合 {recommendation['practical_ratio']}% 案"],
            ["一次相続税額", _yen_text(recommendation["recommended_primary_tax"])],
            ["二次相続税額", _yen_text(recommendation["recommended_secondary_tax"])],
            ["合計税額", _yen_text(recommendation["recommended_total_tax"])],
            ["他案との差額", min_comparison_sentence],
            ["比較コメント", comparison_sentence],
            ["推奨理由", recommendation["recommended_reason"]],
            ["推奨アクション", "現時点では上記割合をたたき台として、配偶者の生活資金・自宅の帰属・納税資金を確認しながら分割方針を具体化することをおすすめします。"],
            ["今後の確認事項（優先度 高）", " / ".join(priority_high[:3])],
            ["今後の確認事項（優先度 中）", " / ".join(priority_mid[:3])],
            ["今後の確認事項（優先度 低）", " / ".join(priority_low[:3])],
        ],
        columns=["項目", "内容"],
    )

    title_text = (
        f"現時点では、配偶者取得割合 {recommendation['recommended_ratio']}% 案を総合的なご提案案として整理しています。"
        f" 一次相続税額は {_yen_text(recommendation['recommended_primary_tax'])}、"
        f"二次相続税額は {_yen_text(recommendation['recommended_secondary_tax'])}、"
        f"合計税額は {_yen_text(recommendation['recommended_total_tax'])} です。"
        f" {min_comparison_sentence} {comparison_sentence}"
    )
    title_df = pd.DataFrame([["結論サマリー", title_text]], columns=["区分", "内容"])
    return summary_df, title_df


def _build_comparison_sheet_df(df_sim: pd.DataFrame, secondary_inputs: SecondaryInputs) -> pd.DataFrame:
    work = _prepare_simulation_dataframe(df_sim)
    if work.empty:
        return pd.DataFrame(columns=["配偶者取得割合", "一次相続税額", "二次相続税額", "合計税額", "比較コメント", "位置づけ"])
    work["配偶者取得割合"] = work["配分表示"]
    work["配分数値"] = work["配分(%)"].apply(_extract_ratio_int)
    min_total = int(work["合計納税額"].min()) if not work.empty else 0
    equal_row = work.iloc[(work["配分数値"] - 50).abs().argsort()[:1]].iloc[0] if not work.empty else None

    rows: list[dict[str, Any]] = []
    for _, row in work.sort_values(by="配分数値").iterrows():
        total = _to_int_safe(row["合計納税額"])
        ratio = _extract_ratio_int(row["配分(%)"])
        if total == min_total:
            position = "税額が最も低い案"
        elif ratio == secondary_inputs.spouse_acquisition_pct:
            position = "現在入力中の案"
        elif equal_row is not None and ratio == _extract_ratio_int(equal_row["配分(%)"]):
            position = "均等分割に近い案"
        else:
            position = "比較案"

        diff = total - min_total
        if diff == 0:
            comment = "合計税額が最も低い水準です。"
        elif row["一次相続税額"] < row["二次相続税額"]:
            comment = f"最小税額案との差額は {abs(diff):,}円。一次は抑えやすい一方、二次で税額が増えやすい傾向です。"
        else:
            comment = f"最小税額案との差額は {abs(diff):,}円。一次の納税負担は増える一方、二次の圧縮効果が見込まれます。"

        rows.append(
            {
                "配偶者取得割合": row["配偶者取得割合"],
                "一次相続税額": _to_int_safe(row["一次相続税額"]),
                "二次相続税額": _to_int_safe(row["二次相続税額"]),
                "合計税額": total,
                "比較コメント": comment,
                "位置づけ": position,
            }
        )
    return pd.DataFrame(rows)


def _build_primary_overview_sheet_df(df1: pd.DataFrame, df_heirs: pd.DataFrame, df_small: pd.DataFrame, df_gifts: pd.DataFrame) -> pd.DataFrame:
    primary_overview = _customerize_dataframe(
        df1,
        {
            "項目": "項目",
            "金額": "金額",
            "備考": "内容",
            "内容": "内容",
        },
    )
    heir_sheet = _customerize_dataframe(
        df_heirs,
        {
            "相続人": "相続人",
            "続柄": "続柄",
            "課税価格": "各相続人ごとの課税対象額",
            "各人別課税価格": "各相続人ごとの課税対象額",
            "按分前税額": "配分調整前の税額",
            "最終税額": "最終税額",
            "配偶者軽減": "配偶者軽減額",
            "保険非課税": "生命保険の非課税額",
            "加算税額": "2割加算額",
            "備考": "内容",
            "注記": "内容",
        },
    )
    if not heir_sheet.empty:
        heir_sheet = _upsert_column(heir_sheet, 0, "区分", "相続人別整理")
    if not primary_overview.empty:
        primary_overview = _upsert_column(primary_overview, 0, "区分", "一次相続の概要")

    small_sheet = _customerize_dataframe(
        df_small,
        {
            "対象宅地": "対象宅地",
            "宅地": "対象宅地",
            "状態": "判定結果",
            "適用状況": "判定結果",
            "理由": "内容",
            "減額額": "減額見込額",
            "備考": "内容",
            "注記": "内容",
        },
    )
    if not small_sheet.empty:
        small_sheet = _upsert_column(small_sheet, 0, "区分", "小規模宅地等の特例 判定結果")

    gifts_sheet = _customerize_dataframe(
        df_gifts,
        {
            "贈与日": "贈与日",
            "受贈者": "受贈者",
            "課税方式": "贈与方式",
            "贈与額": "贈与額",
            "加算対象": "相続財産へ反映",
            "相続戻し対象額": "相続財産へ反映する金額",
            "判定理由": "内容",
        },
        drop_columns=["年分"],
    )
    if not gifts_sheet.empty:
        gifts_sheet = _upsert_column(gifts_sheet, 0, "区分", "贈与・保険等の確認事項")

    sections = [df for df in [primary_overview, heir_sheet, small_sheet, gifts_sheet] if df is not None and not df.empty]
    if not sections:
        return pd.DataFrame(columns=["区分", "項目", "内容"])
    return pd.concat(sections, ignore_index=True, sort=False)


def _build_secondary_overview_sheet_df(df2: pd.DataFrame, df_carryforward: pd.DataFrame, df_successive_credit: pd.DataFrame) -> pd.DataFrame:
    secondary_overview = _customerize_dataframe(
        df2,
        {
            "No": "No",
            "項目": "項目",
            "金額": "金額",
            "備考": "内容",
        },
    )
    if not secondary_overview.empty:
        secondary_overview = _upsert_column(secondary_overview, 0, "区分", "二次相続の概要")

    carryforward_sheet = _customerize_dataframe(
        df_carryforward,
        {
            "相続人": "相続人",
            "続柄": "続柄",
            "取得総額": "一次相続での取得総額",
            "現預金": "現預金",
            "不動産": "不動産",
            "保険": "生命保険",
            "有価証券": "有価証券",
            "その他": "その他",
            "一次税額": "一次相続税額",
            "税引後残高": "税引後の残額",
            "同居": "自宅同居の有無",
            "事業利用": "事業利用の有無",
            "注記": "内容",
        },
    )
    if not carryforward_sheet.empty:
        carryforward_sheet = _upsert_column(carryforward_sheet, 0, "区分", "二次相続試算用の引継財産")

    credit_sheet = _customerize_dataframe(
        df_successive_credit,
        {
            "相続人": "相続人",
            "按分比率": "按分比率",
            "按分前控除額": "按分前の控除額",
            "反映控除額": "反映控除額",
            "注記": "内容",
            "備考": "計算前提",
        },
    )
    if not credit_sheet.empty:
        credit_sheet = _upsert_column(credit_sheet, 0, "区分", "相次相続控除の整理")

    sections = [df for df in [secondary_overview, carryforward_sheet, credit_sheet] if df is not None and not df.empty]
    if not sections:
        return pd.DataFrame(columns=["区分", "項目", "内容"])
    return pd.concat(sections, ignore_index=True, sort=False)


def _build_confirmation_sheet_df(df_audit_notes: pd.DataFrame, df_small_scale_review: pd.DataFrame) -> pd.DataFrame:
    notes_sheet = _customerize_dataframe(
        df_audit_notes,
        {
            "分類": "区分",
            "優先度": "優先度",
            "内容": "今後の確認事項",
        },
    )
    if not notes_sheet.empty:
        notes_sheet["区分"] = notes_sheet["区分"].replace(
            {
                "再判定事項": "今後の確認事項",
                "未充足事項": "今後の確認事項",
                "リスク事項": "今後の確認事項",
                "概算調整事項": "今後の確認事項",
                "税額調整メモ": "今後の確認事項",
                "小宅再判定事項": "今後の確認事項",
                "小宅再判定メモ": "今後の確認事項",
            }
        )

    review_sheet = _customerize_dataframe(
        df_small_scale_review,
        {
            "対象宅地": "対象宅地",
            "状態": "判定結果",
            "一次取得者": "一次相続での取得者",
            "再判定アクション": "次回までに確認したい事項",
            "理由": "内容",
            "注記": "補足",
        },
    )
    if not review_sheet.empty:
        review_sheet = _upsert_column(review_sheet, 0, "区分", "小規模宅地等の特例に関する確認事項")
        review_sheet = _upsert_column(review_sheet, 1, "優先度", "高")

    if notes_sheet.empty and review_sheet.empty:
        notes_sheet = pd.DataFrame(
            [["今後の確認事項", "中", "現時点では追加の重大論点はありません。分割方針と資料確認を進める想定です。"]],
            columns=["区分", "優先度", "今後の確認事項"],
        )
    sections = [df for df in [notes_sheet, review_sheet] if df is not None and not df.empty]
    return pd.concat(sections, ignore_index=True, sort=False)


def _build_assumptions_sheet_df(primary_inputs: PrimaryInputs, secondary_inputs: SecondaryInputs, primary_result: PrimaryResult, secondary_result: SecondaryResult) -> pd.DataFrame:
    heir_summary = []
    if primary_inputs.has_spouse:
        heir_summary.append("配偶者")
    for idx, heir in enumerate(primary_inputs.heirs_info, start=1):
        heir_summary.append(f"相続人{idx}（{heir.get('type', '')}）")

    assumption_rows = [
        ["相続人構成の前提", "、".join(heir_summary) if heir_summary else "入力情報に基づき作成"],
        ["財産評価の前提", f"一次相続の純資産総額は {_yen_text(primary_result.pure_as)} を基準に整理しています。"],
        ["不動産評価の前提", "土地・建物は入力いただいた評価額を前提としており、正式評価や現地確認により変動する場合があります。"],
        ["特例適用の前提", "配偶者の税額軽減、生命保険非課税、小規模宅地等の特例、相次相続控除は入力内容に基づく試算です。"],
        ["二次相続の前提", f"配偶者固有財産 {_yen_text(secondary_inputs.s_own)}、年間生活費 {_yen_text(secondary_inputs.annual_spend)}、経過年数 {secondary_inputs.interval_years}年 を前提にしています。"],
        ["未確定資料がある場合", "遺言書、保険受取人、登記、借入残高、贈与履歴などの確認により結果が変わることがあります。"],
        ["本試算の位置づけ", OUTPUT_RISK_NOTICE],
    ]
    if secondary_result.tax_adjustment_notes:
        assumption_rows.append(["補足前提", " / ".join([_sanitize_customer_text(x) for x in secondary_result.tax_adjustment_notes[:3]])])
    return pd.DataFrame(assumption_rows, columns=["項目", "内容"])


def _build_next_steps_sheet_df(recommendation_df: pd.DataFrame) -> pd.DataFrame:
    return pd.DataFrame(
        [
            ["1", "自宅をどなたが取得するかを確認し、分割方針を具体化します。"],
            ["2", "預金・生命保険・有価証券の配分方針を整理し、納税資金を確認します。"],
            ["3", "二次相続をどの程度重視するかをご家族で共有します。"],
            ["4", "遺言書の有無、保険受取人、登記名義を確認します。"],
            ["5", "確認後、必要に応じて配分割合を再試算し、最終提案案を固めます。"],
        ],
        columns=["No", "次回までにご確認いただきたい事項"],
    )


def _write_dataframe_to_sheet(
    ws,
    df: pd.DataFrame,
    start_row: int,
    title: str,
    description: str,
    styles: dict[str, Any],
) -> int:
    max_col = max(2, len(df.columns) if not df.empty else 2)

    ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=max_col)
    title_cell = ws.cell(row=start_row, column=1, value=title)
    title_cell.font = styles["section_title_font"]
    title_cell.fill = styles["section_fill"]
    title_cell.alignment = styles["left_align"]

    ws.merge_cells(start_row=start_row + 1, start_column=1, end_row=start_row + 1, end_column=max_col)
    desc_cell = ws.cell(row=start_row + 1, column=1, value=description)
    desc_cell.alignment = styles["left_align"]
    desc_cell.fill = styles["soft_fill"]
    desc_cell.border = styles["thin_border"]

    header_row = start_row + 3
    if df.empty:
        ws.cell(row=header_row, column=1, value="項目")
        ws.cell(row=header_row, column=2, value="内容")
        ws.cell(row=header_row + 1, column=1, value="ご案内")
        ws.cell(row=header_row + 1, column=2, value="現時点では表示対象データがありません。")
        data_columns = ["項目", "内容"]
        data_rows = [["ご案内", "現時点では表示対象データがありません。"]]
        actual_cols = 2
    else:
        data_columns = list(df.columns)
        data_rows = df.values.tolist()
        actual_cols = len(data_columns)
        for col_idx, col_name in enumerate(data_columns, start=1):
            ws.cell(row=header_row, column=col_idx, value=col_name)

    for col_idx in range(1, actual_cols + 1):
        cell = ws.cell(row=header_row, column=col_idx)
        cell.fill = styles["header_fill"]
        cell.font = styles["header_font"]
        cell.alignment = styles["center_align"]
        cell.border = styles["thin_border"]

    for row_offset, row_values in enumerate(data_rows, start=1):
        row_idx = header_row + row_offset
        for col_idx, value in enumerate(row_values, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.border = styles["thin_border"]
            if isinstance(value, Number):
                cell.number_format = '#,##0'
                cell.alignment = styles["right_align"]
            else:
                cell.alignment = styles["left_align"]

    return header_row + len(data_rows) + 2


def _apply_sheet_layout(ws, title: str, subtitle: str, orientation: str, fit_width: int = 1) -> None:
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A5"
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.orientation = orientation
    ws.page_setup.fitToPage = True
    ws.page_setup.fitToWidth = fit_width
    ws.page_setup.fitToHeight = 0
    ws.page_margins = PageMargins(left=0.35, right=0.35, top=0.5, bottom=0.5, header=0.2, footer=0.2)
    ws.oddHeader.center.text = f"&\"Meiryo,Bold\"&14 {title}"
    ws.oddFooter.right.text = "&P / &N"
    ws["A1"] = EXCEL_TITLE
    ws["A1"].font = Font(name="Meiryo", size=16, bold=True, color="1F2C4D")
    ws["A2"] = subtitle
    ws["A2"].font = Font(name="Meiryo", size=10, color="556070")
    ws["A3"] = OUTPUT_RISK_NOTICE
    ws["A3"].font = Font(name="Meiryo", size=9, color="556070")
    ws.row_dimensions[1].height = 24
    ws.row_dimensions[2].height = 20
    ws.row_dimensions[3].height = 30


def _excel_char_width_units(value: Any) -> float:
    if value is None:
        return 0.0
    text_value = str(value).replace("\r", "\n")
    widest_line = 0.0
    for line in text_value.split("\n"):
        line_units = 0.0
        for ch in line:
            if ch == "\t":
                line_units += 4.0
            elif unicodedata.east_asian_width(ch) in ("W", "F", "A"):
                line_units += 1.9
            else:
                line_units += 1.0
        widest_line = max(widest_line, line_units)
    return widest_line


def _excel_is_numeric_header(header_value: Any, numeric_keywords: list[str]) -> bool:
    if header_value is None:
        return False
    header_text = str(header_value)
    return any(keyword in header_text for keyword in numeric_keywords)


def _excel_estimate_wrapped_lines(value: Any, column_width: float, is_numeric: bool) -> int:
    if value is None:
        return 1
    text_value = str(value).replace("\r", "\n")
    if not text_value:
        return 1
    if is_numeric:
        return max(1, text_value.count("\n") + 1)
    usable_width = max(column_width - 2.2, 6.0)
    total_lines = 0
    for line in text_value.split("\n"):
        units = _excel_char_width_units(line)
        total_lines += max(1, int(math.ceil(units / usable_width)))
    return max(total_lines, 1)


def _autosize_and_format_sheet(ws, numeric_keywords: Optional[list[str]] = None) -> None:
    numeric_keywords = numeric_keywords or ["税額", "金額", "額", "残高", "財産", "控除", "減額", "差額", "取得総額", "取得額", "割合", "%"]
    column_meta: dict[int, dict[str, Any]] = {}
    header_candidates: dict[int, str] = {}
    merged_cells = {coord for merged in ws.merged_cells.ranges for coord in merged.cells}

    for row in ws.iter_rows():
        if not row:
            continue
        for cell in row:
            if cell.coordinate in merged_cells and cell.coordinate != cell.merged_cells.start_cell.coordinate:
                continue
            value = cell.value
            base_font = cell.font or Font()
            cell.font = Font(
                name="Meiryo",
                size=base_font.sz or 10,
                bold=base_font.bold,
                italic=base_font.italic,
                color=(base_font.color.rgb if base_font.color and base_font.color.type == "rgb" else None),
            )
            if value is None:
                continue

            text_value = str(value)
            if 1 <= cell.row <= 10 and cell.column not in header_candidates and text_value.strip():
                header_candidates[cell.column] = text_value

            meta = column_meta.setdefault(cell.column, {"max_width": 0.0, "is_numeric": False, "max_numeric_width": 0.0, "max_text_width": 0.0})
            header_text = header_candidates.get(cell.column, "")
            width_units = _excel_char_width_units(value)
            is_numeric_column = isinstance(value, Number) or _excel_is_numeric_header(header_text, numeric_keywords)
            if is_numeric_column:
                meta["is_numeric"] = True
                meta["max_numeric_width"] = max(meta["max_numeric_width"], width_units)
                meta["max_width"] = max(meta["max_width"], min(max(width_units + 2.0, 10.5), 18.5))
                if isinstance(value, Number):
                    cell.number_format = '#,##0'
                cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=False)
            else:
                meta["max_text_width"] = max(meta["max_text_width"], width_units)
                meta["max_width"] = max(meta["max_width"], min(max(width_units + 2.8, 13.0), 42.0))
                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

            if cell.row <= 3:
                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

    resolved_widths: dict[int, float] = {}
    for col_idx, meta in column_meta.items():
        width = meta["max_width"]
        if meta["is_numeric"] and meta["max_text_width"] <= 4.0:
            width = min(max(width, 11.0), 18.0)
        else:
            width = min(max(width, 14.0), 42.0)
        ws.column_dimensions[get_column_letter(col_idx)].width = width
        resolved_widths[col_idx] = width

    for row in ws.iter_rows():
        if not row:
            continue
        row_idx = row[0].row
        max_required_lines = 1
        has_wrapped_text = False
        for cell in row:
            if cell.coordinate in merged_cells and cell.coordinate != cell.merged_cells.start_cell.coordinate:
                continue
            value = cell.value
            if value is None:
                continue
            meta = column_meta.get(cell.column, {"is_numeric": False})
            col_width = resolved_widths.get(cell.column, ws.column_dimensions[get_column_letter(cell.column)].width or 12.0)
            line_count = _excel_estimate_wrapped_lines(value, col_width, bool(meta.get("is_numeric")))
            max_required_lines = max(max_required_lines, line_count)
            has_wrapped_text = has_wrapped_text or line_count > 1
            if meta.get("is_numeric"):
                cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=False)
            else:
                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
            if row_idx <= 3:
                cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

        if row_idx <= 3:
            current_height = ws.row_dimensions[row_idx].height or 0
            default_height = 18.0 if row_idx == 1 else (20.0 if row_idx == 2 else 30.0)
            ws.row_dimensions[row_idx].height = max(current_height, default_height)
        else:
            base_height = 18.0
            if has_wrapped_text:
                base_height += (max_required_lines - 1) * 12.8
            ws.row_dimensions[row_idx].height = max(ws.row_dimensions[row_idx].height or 0, min(base_height, 96.0))


def _insert_comparison_chart(ws, df_comp: pd.DataFrame, recommended_ratio: int) -> None:
    if df_comp is None or df_comp.empty:
        return

    chart_source_col_start = 10  # J
    chart_source_headers = ["配偶者取得割合", "一次相続税額", "二次相続税額", "合計税額"]
    safe_df = df_comp.copy()

    for required_col in chart_source_headers:
        if required_col not in safe_df.columns:
            return

    safe_df["配偶者取得割合"] = safe_df["配偶者取得割合"].apply(_extract_ratio_int)
    safe_df["一次相続税額"] = safe_df["一次相続税額"].apply(lambda v: _to_int_safe(v, 0))
    safe_df["二次相続税額"] = safe_df["二次相続税額"].apply(lambda v: _to_int_safe(v, 0))
    safe_df["合計税額"] = safe_df["合計税額"].apply(lambda v: _to_int_safe(v, 0))
    safe_df = safe_df.dropna(subset=["配偶者取得割合"]).sort_values("配偶者取得割合").reset_index(drop=True)

    if safe_df.empty:
        return

    chart_start_row = 5
    data_start_row = chart_start_row + 18

    for offset, header in enumerate(chart_source_headers):
        target_cell = ws.cell(row=data_start_row, column=chart_source_col_start + offset, value=header)
        target_cell.font = Font(name="Meiryo", size=9, bold=True)

    for row_offset, (_, row) in enumerate(safe_df.iterrows(), start=1):
        ws.cell(row=data_start_row + row_offset, column=chart_source_col_start + 0, value=int(row["配偶者取得割合"]))
        ws.cell(row=data_start_row + row_offset, column=chart_source_col_start + 1, value=int(row["一次相続税額"]))
        ws.cell(row=data_start_row + row_offset, column=chart_source_col_start + 2, value=int(row["二次相続税額"]))
        ws.cell(row=data_start_row + row_offset, column=chart_source_col_start + 3, value=int(row["合計税額"]))

    data_end_row = data_start_row + len(safe_df)
    if data_end_row <= data_start_row:
        return

    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    chart.overlap = 0
    chart.grouping = "clustered"
    chart.title = "配偶者取得割合ごとの税額比較"
    chart.y_axis.title = "税額（円）"
    chart.x_axis.title = "配偶者取得割合"
    chart.height = 8.5
    chart.width = 15.5
    chart.legend.position = "b"
    chart.varyColors = False
    chart.gapWidth = 80

    cats = Reference(ws, min_col=chart_source_col_start, min_row=data_start_row + 1, max_row=data_end_row)
    for series_col in range(chart_source_col_start + 1, chart_source_col_start + 4):
        values = Reference(ws, min_col=series_col, min_row=data_start_row, max_row=data_end_row)
        chart.add_data(values, titles_from_data=True)
    chart.set_categories(cats)

    series_palette = ["5B9BD5", "A5A5A5", "4472C4"]
    for idx, series in enumerate(chart.series):
        fill_color = series_palette[idx % len(series_palette)]
        series.graphicalProperties.solidFill = fill_color
        series.graphicalProperties.line.solidFill = fill_color

    ws.add_chart(chart, "A5")

    for col in range(chart_source_col_start, chart_source_col_start + len(chart_source_headers)):
        ws.column_dimensions[get_column_letter(col)].hidden = True
    for row in range(data_start_row, data_end_row + 1):
        ws.row_dimensions[row].hidden = True


def create_excel_file(
    primary_inputs: PrimaryInputs,
    primary_result: PrimaryResult,
    secondary_inputs: SecondaryInputs,
    secondary_result: SecondaryResult,
    df1: pd.DataFrame,
    df_heirs: pd.DataFrame,
    df_small: pd.DataFrame,
    df_gifts: pd.DataFrame,
    df2: pd.DataFrame,
    df_sim: pd.DataFrame,
    df_snapshot_summary: pd.DataFrame,
    df_carryforward: pd.DataFrame,
    df_audit_notes: pd.DataFrame,
    df_small_scale_review: pd.DataFrame,
    df_successive_credit: pd.DataFrame,
) -> bytes:
    summary_df, summary_text_df = _build_summary_sheet_df(
        primary_inputs,
        primary_result,
        secondary_inputs,
        secondary_result,
        df_sim,
        df_audit_notes,
    )
    recommendation = _choose_recommendation_plan(df_sim, secondary_inputs.spouse_acquisition_pct)
    comparison_df = _build_comparison_sheet_df(df_sim, secondary_inputs)
    primary_overview_df = _build_primary_overview_sheet_df(df1, df_heirs, df_small, df_gifts)
    secondary_overview_df = _build_secondary_overview_sheet_df(df2, df_carryforward, df_successive_credit)
    confirmation_df = _build_confirmation_sheet_df(df_audit_notes, df_small_scale_review)
    assumptions_df = _build_assumptions_sheet_df(primary_inputs, secondary_inputs, primary_result, secondary_result)
    next_steps_df = _build_next_steps_sheet_df(summary_df)

    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        pd.DataFrame().to_excel(writer, sheet_name="結論サマリー", index=False)
        pd.DataFrame().to_excel(writer, sheet_name="税額比較一覧", index=False)
        primary_overview_df.to_excel(writer, sheet_name="一次相続の整理", index=False, startrow=3)
        secondary_overview_df.to_excel(writer, sheet_name="二次相続の整理", index=False, startrow=3)
        confirmation_df.to_excel(writer, sheet_name="今後の確認事項", index=False, startrow=3)
        assumptions_df.to_excel(writer, sheet_name="前提条件", index=False, startrow=3)
        next_steps_df.to_excel(writer, sheet_name="今後の進め方", index=False, startrow=3)

    output.seek(0)
    wb = load_workbook(output)

    styles = {
        "header_fill": PatternFill(start_color="1F2C4D", end_color="1F2C4D", fill_type="solid"),
        "header_font": Font(name="Meiryo", color="FFFFFF", bold=True),
        "section_fill": PatternFill(start_color="DCE6F2", end_color="DCE6F2", fill_type="solid"),
        "soft_fill": PatternFill(start_color="F4F6F8", end_color="F4F6F8", fill_type="solid"),
        "highlight_fill": PatternFill(start_color="E5F0EA", end_color="E5F0EA", fill_type="solid"),
        "accent_fill": PatternFill(start_color="FFF4DE", end_color="FFF4DE", fill_type="solid"),
        "section_title_font": Font(name="Meiryo", size=12, bold=True, color="1F2C4D"),
        "thin_border": Border(left=Side(style="thin", color="C9CED6"), right=Side(style="thin", color="C9CED6"), top=Side(style="thin", color="C9CED6"), bottom=Side(style="thin", color="C9CED6")),
        "left_align": Alignment(horizontal="left", vertical="center", wrap_text=True),
        "center_align": Alignment(horizontal="center", vertical="center", wrap_text=True),
        "right_align": Alignment(horizontal="right", vertical="center", wrap_text=False),
    }

    summary_ws = wb["結論サマリー"]
    _apply_sheet_layout(summary_ws, "結論サマリー", "まず最初に、税額と実務面を合わせた現時点のご提案を整理しています。", "landscape", 1)
    current_row = 5
    current_row = _write_dataframe_to_sheet(
        summary_ws,
        summary_text_df,
        current_row,
        "総括コメント",
        "ご家族での話し合いの起点になるよう、結論を文章でも読みやすく整理しています。",
        styles,
    )
    current_row = _write_dataframe_to_sheet(
        summary_ws,
        summary_df,
        current_row + 1,
        "結論サマリー",
        "税額だけでなく、配偶者の生活資金・自宅の持ち方・二次相続まで含めて整理した現時点の推奨内容です。",
        styles,
    )
    _insert_comparison_chart(summary_ws, comparison_df[["配偶者取得割合", "一次相続税額", "二次相続税額", "合計税額"]], recommendation["recommended_ratio"])
    summary_ws["A5"].fill = styles["highlight_fill"]
    summary_ws["A5"].font = styles["section_title_font"]
    for row in summary_ws.iter_rows(min_row=5, max_row=summary_ws.max_row):
        for cell in row:
            cell.border = styles["thin_border"]

    comparison_ws = wb["税額比較一覧"]
    _apply_sheet_layout(comparison_ws, "税額比較一覧", "配偶者取得割合ごとの一次・二次・合計税額を比較しています。", "landscape", 1)
    _write_dataframe_to_sheet(
        comparison_ws,
        comparison_df,
        5,
        "税額比較一覧",
        "税額が最も低い案、実務上の推奨案、現在の想定案の違いを一覧で確認できます。",
        styles,
    )

    primary_ws = wb["一次相続の整理"]
    _apply_sheet_layout(primary_ws, "一次相続の整理", "一次相続の計算根拠、相続人ごとの税額、小規模宅地等の特例や贈与整理をまとめています。", "portrait", 1)

    secondary_ws = wb["二次相続の整理"]
    _apply_sheet_layout(secondary_ws, "二次相続の整理", "一次相続から二次相続へどう影響するかを整理しています。", "portrait", 1)

    confirmation_ws = wb["今後の確認事項"]
    _apply_sheet_layout(confirmation_ws, "今後の確認事項", "不安を煽らない形で、次回までに確認したい事項を優先度別に整理しています。", "portrait", 1)

    assumptions_ws = wb["前提条件"]
    _apply_sheet_layout(assumptions_ws, "前提条件", "本試算の前提条件と、正式申告までに変動し得るポイントを整然とまとめています。", "portrait", 1)

    next_steps_ws = wb["今後の進め方"]
    _apply_sheet_layout(next_steps_ws, "今後の進め方", "次回面談までにご確認いただきたい事項と、今後の進め方を整理しています。", "portrait", 1)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    cell.value = _sanitize_customer_text(cell.value)
                if cell.row >= 4:
                    cell.border = styles["thin_border"]
        _autosize_and_format_sheet(ws)
        ws.print_options.gridLines = False
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.print_title_rows = "$1:$4"

    comparison_ws.freeze_panes = "A5"
    comparison_ws.auto_filter.ref = f"A8:{get_column_letter(max(1, comparison_ws.max_column))}{comparison_ws.max_row}"
    primary_ws.auto_filter.ref = f"A4:{get_column_letter(max(1, primary_ws.max_column))}{primary_ws.max_row}"
    secondary_ws.auto_filter.ref = f"A4:{get_column_letter(max(1, secondary_ws.max_column))}{secondary_ws.max_row}"
    confirmation_ws.auto_filter.ref = f"A4:{get_column_letter(max(1, confirmation_ws.max_column))}{confirmation_ws.max_row}"
    assumptions_ws.auto_filter.ref = f"A4:{get_column_letter(max(1, assumptions_ws.max_column))}{assumptions_ws.max_row}"
    next_steps_ws.auto_filter.ref = f"A4:{get_column_letter(max(1, next_steps_ws.max_column))}{next_steps_ws.max_row}"

    final_output = BytesIO()
    wb.save(final_output)
    return final_output.getvalue()


def build_simulation_figure(df_sim: pd.DataFrame) -> go.Figure:
    plot_df = _prepare_simulation_dataframe(df_sim)
    if plot_df.empty:
        fig = go.Figure()
        fig.update_layout(
            title="配分シミュレーション結果",
            xaxis_title="配偶者取得割合",
            yaxis_title="相続税額（円）",
            barmode="group",
        )
        return fig
    plot_df["配分数値"] = plot_df["配分(%)"].apply(_extract_ratio_int)
    plot_df["配分(%)"] = plot_df["配分表示"]
    plot_df = plot_df.sort_values("配分数値").reset_index(drop=True)

    x_labels = plot_df["配分(%)"]
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=x_labels,
            y=plot_df["一次相続税額"],
            name="一次相続税（概算）",
            marker_color=f"#{COLOR_NAVY}",
            hovertemplate="配偶者取得割合: %{x}<br>一次相続税（概算）: %{y:,}円<extra></extra>",
        )
    )
    fig.add_trace(
        go.Bar(
            x=x_labels,
            y=plot_df["二次相続税額"],
            name="二次相続税（概算）",
            marker_color=f"#{COLOR_GOLD}",
            hovertemplate="配偶者取得割合: %{x}<br>二次相続税（概算）: %{y:,}円<extra></extra>",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=x_labels,
            y=plot_df["合計納税額"],
            name="合計納税額（概算）",
            line=dict(color=f"#{COLOR_RED}", width=4),
            hovertemplate="配偶者取得割合: %{x}<br>合計納税額（概算）: %{y:,}円<extra></extra>",
        )
    )
    fig.update_layout(
        barmode="stack",
        title="配偶者取得割合別 税額比較シミュレーション",
        xaxis_title="横軸：配偶者取得割合（%）",
        yaxis_title="縦軸：概算税額（円）",
        legend_title="表示項目",
        hovermode="x unified",
    )
    fig.update_xaxes(type="category", categoryorder="array", categoryarray=list(x_labels))
    fig.update_yaxes(tickformat=",", rangemode="tozero")
    return fig


def render_audit_evidence() -> None:
    st.markdown(
        f"""
        <div style="background-color: #f9f9f9; border: 2px solid #{COLOR_GOLD}; padding: 20px; border-radius: 5px;">
            <p style="color: #{COLOR_NAVY}; font-weight: bold; margin-bottom: 10px;">🛡️ 山根会計 監査証跡エビデンス (v31.16)</p>
            <p style="font-size: 0.9em; line-height: 1.6;">担当: 川東 / 本版はご提案用の試算レポートです。<br>
            各人別課税価格・各人別税額・配偶者税額軽減の土台に加え、生命保険金非課税の受取人別管理、2割加算、小規模宅地等の要件判定付き概算ロジックを反映しています。<br>
            二次相続・相次相続控除・小規模宅地等の個別論点は、実務利用前に別途確認が必要です。顧客提出前提の完成資料ではありません。提出前に別途レビューを実施してください。</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


# =========================================================
# 8. UI Rendering
# =========================================================
def render_sidebar() -> None:
    st.sidebar.markdown(f"###  🏢  {APP_TITLE}")
    st.sidebar.info(APP_LOGIN_USER_LABEL)


def render_risk_notice(message: str, level: str = "warning") -> None:
    if level == "error":
        st.error(message)
    elif level == "info":
        st.info(message)
    else:
        st.warning(message)


def render_tab_basic() -> tuple[int, bool, list[dict[str, Any]]]:
    add_print_button("1. 基本構成")
    st.subheader("相続関係の設定（内部確認用）")
    c1, c2 = st.columns(2)
    heir_count = c1.number_input("相続人の人数（配偶者除く）", min_value=1, max_value=10, value=2, key="in_child")
    has_spouse = c2.checkbox("配偶者は健在", value=True, key="in_spouse")
    heirs_info: list[dict[str, Any]] = []
    for i in range(heir_count):
        st.write(f"##### 相続人 {i + 1}")
        top_col1, top_col2 = st.columns(2)
        h_type = top_col1.selectbox(f"相続人 {i + 1} の続柄", HEIR_TYPE_OPTIONS, key=f"rel_{i}")
        is_substitute = False
        if h_type == HEIR_TYPE_GRANDCHILD:
            is_substitute = top_col2.checkbox(f"相続人 {i + 1} は代襲相続人", value=False, key=f"substitute_{i}")
        detail_col1, detail_col2, detail_col3 = st.columns(3)
        has_birth_date = detail_col1.checkbox(f"相続人 {i + 1} の生年月日を入力する", value=False, key=f"birth_enabled_{i}")
        birth_date = None
        if has_birth_date:
            birth_date = detail_col1.date_input(
                f"相続人 {i + 1} の生年月日",
                value=date(2000, 1, 1),
                min_value=date(1900, 1, 1),
                max_value=date.today(),
                key=f"birth_date_{i}",
            )
        is_disabled = detail_col2.checkbox(f"相続人 {i + 1} は障害者", value=False, key=f"disabled_{i}")
        is_special_disabled = False
        if is_disabled:
            is_special_disabled = detail_col3.checkbox(f"相続人 {i + 1} は特別障害者", value=False, key=f"special_disabled_{i}")
        heirs_info.append({
            "type": h_type,
            "is_substitute": is_substitute,
            "birth_date": birth_date,
            "is_disabled": is_disabled,
            "is_special_disabled": is_special_disabled,
        })
    return heir_count, has_spouse, heirs_info


def render_small_scale_input_section(category: str, title: str, has_spouse: bool, heirs_info: list[dict[str, Any]]) -> SmallScaleInput:
    st.write(f"##### {title}：小規模宅地等の特例判定")
    st.caption(SMALL_SCALE_RISK_NOTICE)
    options = build_heir_labels(has_spouse, heirs_info)
    option_labels = [label for label, _ in options] if options else ["未設定"]
    apply_special_rule = st.checkbox(f"{title}で小宅特例を検討する", value=(category == LAND_CATEGORY_HOME), key=f"apply_small_{category}")
    acquirer_name = st.selectbox(f"{title}の取得者", option_labels, key=f"acquirer_{category}")

    if category == LAND_CATEGORY_HOME:
        is_spouse_acquirer = st.checkbox("配偶者が取得", value=(acquirer_name == "配偶者"), key=f"home_spouse_{category}")
        is_cohabiting_heir = st.checkbox("同居親族が取得", value=False, key=f"home_cohab_{category}")
        is_no_house_heir = st.checkbox("家なき子要件に該当", value=False, key=f"home_nohouse_{category}")
        will_hold_until_filing = st.checkbox("申告期限まで保有する", value=True, key=f"home_hold_{category}")
        will_reside_until_filing = st.checkbox("申告期限まで居住継続する", value=True, key=f"home_live_{category}")
        return SmallScaleInput(category=category, acquirer_name=acquirer_name, apply_special_rule=apply_special_rule, is_spouse_acquirer=is_spouse_acquirer, is_cohabiting_heir=is_cohabiting_heir, is_no_house_heir=is_no_house_heir, will_hold_until_filing=will_hold_until_filing, will_reside_until_filing=will_reside_until_filing)

    if category == LAND_CATEGORY_BUSINESS:
        is_business_successor = st.checkbox("取得者は事業承継者", value=False, key=f"biz_successor_{category}")
        will_continue_business = st.checkbox("申告期限まで事業継続", value=False, key=f"biz_continue_{category}")
        will_hold_until_filing = st.checkbox("申告期限まで保有継続", value=False, key=f"biz_hold_{category}")
        return SmallScaleInput(category=category, acquirer_name=acquirer_name, apply_special_rule=apply_special_rule, is_business_successor=is_business_successor, will_continue_business=will_continue_business, will_hold_until_filing=will_hold_until_filing)

    will_continue_rental = st.checkbox("申告期限まで貸付継続", value=False, key=f"rent_continue_{category}")
    will_hold_until_filing = st.checkbox("申告期限まで保有継続", value=False, key=f"rent_hold_{category}")
    return SmallScaleInput(category=category, acquirer_name=acquirer_name, apply_special_rule=apply_special_rule, will_continue_rental=will_continue_rental, will_hold_until_filing=will_hold_until_filing)


def render_tab_primary_inputs(heir_count: int, has_spouse: bool, heirs_info: list[dict[str, Any]]) -> PrimaryInputs:
    add_print_button("2. 一次財産詳細")
    st.subheader("一次相続：財産・贈与入力（内部確認用）")
    render_risk_notice(INSURANCE_GIFT_RISK_NOTICE, level="info")
    col_a, col_b, col_c = st.columns(3)
    with col_a:
        st.write("#### 🏗️ 不動産")
        v_home = st.number_input("特定居住用：評価額", min_value=0, value=32781936, key="v_home")
        a_home = st.number_input("特定居住用：面積(㎡)", min_value=0, value=330, key="a_home")
        home_rule = render_small_scale_input_section(LAND_CATEGORY_HOME, "特定居住用", has_spouse, heirs_info)
        st.divider()
        v_biz = st.number_input("特定事業用：評価額", min_value=0, value=0, key="v_biz")
        a_biz = st.number_input("特定事業用：面積(㎡)", min_value=0, value=400, key="a_biz")
        biz_rule = render_small_scale_input_section(LAND_CATEGORY_BUSINESS, "特定事業用", has_spouse, heirs_info)
        st.divider()
        v_rent = st.number_input("貸付事業用：評価額", min_value=0, value=0, key="v_rent")
        a_rent = st.number_input("貸付事業用：面積(㎡)", min_value=0, value=200, key="a_rent")
        rent_rule = render_small_scale_input_section(LAND_CATEGORY_RENTAL, "貸付事業用", has_spouse, heirs_info)
        v_build = st.number_input("建物評価", min_value=0, value=1700044, key="v_build")
    with col_b:
        st.write("#### 💵 金融財産")
        v_stock = st.number_input("有価証券", min_value=0, value=45132788, key="v_stock")
        v_cash = st.number_input("現預金", min_value=0, value=45573502, key="v_cash")
        v_ins = st.number_input("生命保険金（総額）", min_value=0, value=3651514, key="v_ins")
        recipient_options = build_recipient_options(has_spouse, heirs_info)
        recipient_labels = [f"{label}（{heir_type}）" for label, heir_type, _, _ in recipient_options]
        recipient_map = {f"{label}（{heir_type}）": (label, heir_type, is_statutory_heir, is_two_tenths_target) for label, heir_type, is_statutory_heir, is_two_tenths_target in recipient_options}
        insurance_entry_count = st.number_input("生命保険金の受取人明細数", min_value=0, max_value=MAX_INSURANCE_RECIPIENT_ROWS, value=min(1, len(recipient_labels)), key="insurance_entry_count")
        insurance_entries: list[InsuranceRecipientInput] = []
        for idx in range(int(insurance_entry_count)):
            recipient_key = st.selectbox(f"保険金受取人 {idx + 1}", recipient_labels, key=f"insurance_recipient_{idx}")
            amount = st.number_input(f"保険金受取額 {idx + 1}", min_value=0, value=v_ins if idx == 0 else 0, key=f"insurance_amount_{idx}")
            label, heir_type, is_statutory_heir, is_two_tenths_target = recipient_map[recipient_key]
            insurance_entries.append(InsuranceRecipientInput(recipient_name=label, recipient_type=heir_type, amount=amount, is_statutory_heir=is_statutory_heir, is_two_tenths_target=is_two_tenths_target))
        if insurance_entries:
            entered_total = sum(entry.amount for entry in insurance_entries)
            st.caption(f"※受取人明細合計: {entered_total:,}円 / 総額入力: {v_ins:,}円。差異がある場合は内部で総額に合わせて按分補正します。")
        v_others = st.number_input("その他", min_value=0, value=1662687, key="v_others")
        st.write("#### 📉 債務・葬式")
        v_debt = st.number_input("債務", min_value=0, value=322179, key="v_debt")
        v_funeral = st.number_input("葬式費用", min_value=0, value=41401, key="v_funeral")
    with col_c:
        st.write("#### 🎁 贈与台帳（明細管理）")
        date_of_death = st.date_input("相続開始日", value=date.today(), key="date_of_death")
        gift_recipient_options = build_gift_recipient_options(has_spouse, heirs_info)
        gift_recipient_labels = [f"{label}（{heir_type}）" for label, heir_type in gift_recipient_options]
        gift_recipient_map = {f"{label}（{heir_type}）": (label, heir_type) for label, heir_type in gift_recipient_options}
        gift_record_count = st.number_input("贈与明細件数", min_value=0, max_value=MAX_GIFT_RECORD_ROWS, value=0, key="gift_record_count")
        gift_records: list[GiftRecord] = []
        for idx in range(int(gift_record_count)):
            st.markdown(f"**贈与明細 {idx + 1}**")
            gift_date_value = st.date_input(f"贈与日 {idx + 1}", value=date.today(), key=f"gift_date_{idx}")
            recipient_key = st.selectbox(f"受贈者 {idx + 1}", gift_recipient_labels, key=f"gift_recipient_{idx}")
            gift_amount = st.number_input(f"贈与額 {idx + 1}", min_value=0, value=0, key=f"gift_amount_{idx}")
            gift_tax_type = st.selectbox(f"課税方式 {idx + 1}", GIFT_TYPE_OPTIONS, key=f"gift_tax_type_{idx}")
            recipient_name, recipient_type = gift_recipient_map[recipient_key]
            gift_records.append(
                GiftRecord(
                    gift_date=gift_date_value,
                    recipient_name=recipient_name,
                    recipient_type=recipient_type,
                    amount=gift_amount,
                    tax_type=gift_tax_type,
                )
            )
        if gift_records:
            annual_total = sum(record.amount for record in gift_records if record.tax_type == GIFT_TYPE_ANNUAL)
            seisan_total = sum(record.amount for record in gift_records if record.tax_type == GIFT_TYPE_SEISAN)
            st.caption(f"暦年課税入力合計: {annual_total:,}円 / 相続時精算課税入力合計: {seisan_total:,}円")
        else:
            st.caption("贈与明細がない場合は0件のままで構いません。")

    return PrimaryInputs(
        heir_count=heir_count,
        has_spouse=has_spouse,
        heirs_info=heirs_info,
        date_of_death=date_of_death,
        v_home=v_home,
        a_home=a_home,
        v_biz=v_biz,
        a_biz=a_biz,
        v_rent=v_rent,
        a_rent=a_rent,
        small_scale_inputs={
            LAND_CATEGORY_HOME: home_rule,
            LAND_CATEGORY_BUSINESS: biz_rule,
            LAND_CATEGORY_RENTAL: rent_rule,
        },
        v_build=v_build,
        v_stock=v_stock,
        v_cash=v_cash,
        v_ins=v_ins,
        insurance_entries=insurance_entries,
        v_others=v_others,
        v_debt=v_debt,
        v_funeral=v_funeral,
        gift_records=gift_records,
    )


def render_tab_primary_detail(df1: pd.DataFrame, df_heirs: pd.DataFrame, df_small: pd.DataFrame, df_gifts: pd.DataFrame) -> None:
    add_print_button("3. 一次相続明細")
    st.subheader("一次相続：計算明細（概算・内部確認用）")
    st.table(df1)
    st.divider()
    st.subheader("小規模宅地等の特例 判定結果（概算・要確認）")
    st.dataframe(df_small, width="stretch")
    st.divider()
    st.subheader("贈与加算・相続時精算課税 明細（概算・要確認）")
    if df_gifts.empty:
        st.caption("贈与明細はありません。")
    else:
        st.dataframe(df_gifts, width="stretch")
    st.divider()
    st.subheader("各人別課税価格・各人別税額（概算・要個別確認）")
    st.dataframe(df_heirs, width="stretch")


def render_tab_secondary_detail(df2: pd.DataFrame, df_snapshot_summary: pd.DataFrame, df_carryforward: pd.DataFrame, df_audit_notes: pd.DataFrame, df_small_scale_review: pd.DataFrame, df_successive_credit: pd.DataFrame) -> None:
    add_print_button("4. 二次相続明細")
    st.subheader("二次相続：計算明細予測（概算参考・要確認）")
    render_risk_notice(SECONDARY_RISK_NOTICE)
    st.table(df2)
    st.divider()
    st.subheader("一次→二次 接続サマリー（内部確認用）")
    st.dataframe(df_snapshot_summary, width="stretch")
    st.divider()
    st.subheader("各人別 carry forward 一覧（内部確認用）")
    st.dataframe(df_carryforward, width="stretch")
    st.divider()
    st.subheader("監査メモ・再判定事項（内部確認用）")
    st.dataframe(df_audit_notes, width="stretch")
    st.divider()
    st.subheader("相次相続控除 明細（内部確認用）")
    st.dataframe(df_successive_credit, width="stretch")
    st.divider()
    st.subheader("小規模宅地等 再判定レビュー（内部確認用）")
    st.dataframe(df_small_scale_review, width="stretch")


def estimate_total_taxable_price_reference(primary_inputs: PrimaryInputs) -> Decimal:
    total_red, land_eval, _ = calculate_small_scale_reduction(primary_inputs)
    st_count = primary_inputs.heir_count + (1 if primary_inputs.has_spouse else 0)
    ins_ded = calculate_life_insurance_deduction(primary_inputs.v_ins, st_count)
    pure_as = land_eval + to_d(primary_inputs.v_build) + to_d(primary_inputs.v_stock) + to_d(primary_inputs.v_cash) + to_d(primary_inputs.v_ins) + to_d(primary_inputs.v_others)
    labels = build_heir_labels(primary_inputs.has_spouse, primary_inputs.heirs_info)
    annual_addbacks, seisan_addbacks, _ = calculate_gift_addbacks(primary_inputs.gift_records, primary_inputs.date_of_death, labels)
    tax_p = pure_as - ins_ded - to_d(primary_inputs.v_debt) - to_d(primary_inputs.v_funeral) + sum(annual_addbacks, Decimal("0")) + sum(seisan_addbacks, Decimal("0"))
    return quantize_yen(max(Decimal("0"), tax_p))


def build_default_acquisition_input_amounts(total_taxable_price: Decimal, has_spouse: bool, heirs_info: list[dict[str, Any]]) -> list[int]:
    fallback_shares = allocate_actual_shares(has_spouse, heirs_info, 50 if has_spouse else 0)
    defaults = normalize_amounts_to_total(total_taxable_price, [], fallback_shares)
    return [int(amount) for amount in defaults]


def build_simulation_allocation_inputs(
    total_taxable_price: Decimal,
    current_inputs: list[int],
    has_spouse: bool,
    heirs_info: list[dict[str, Any]],
    spouse_acquisition_pct: int,
) -> list[int]:
    fallback_shares = allocate_actual_shares(has_spouse, heirs_info, spouse_acquisition_pct)
    if not has_spouse:
        normalized = normalize_amounts_to_total(total_taxable_price, [to_d(max(0, amount)) for amount in current_inputs], fallback_shares)
        return [int(amount) for amount in normalized]

    non_spouse_count = len(heirs_info)
    spouse_amount = quantize_yen(total_taxable_price * to_d(spouse_acquisition_pct) / PERCENT_DENOMINATOR)
    remaining_amount = max(Decimal("0"), total_taxable_price - spouse_amount)

    desired_non_spouse = [to_d(max(0, amount)) for amount in current_inputs[1: 1 + non_spouse_count]]
    fallback_non_spouse_shares = fallback_shares[1:] if len(fallback_shares) > 1 else [Decimal("0")] * non_spouse_count
    normalized_non_spouse = normalize_amounts_to_total(remaining_amount, desired_non_spouse, fallback_non_spouse_shares)

    combined = [spouse_amount] + normalized_non_spouse
    return [int(amount) for amount in combined]


def render_tab_secondary_parameters(has_spouse: bool, heirs_info: list[dict[str, Any]], estimated_tax_p: Decimal) -> SecondaryInputs:
    add_print_button("5. 二次推移予測")
    st.subheader("二次推移パラメータ設定（参考試算用）")
    render_risk_notice(SECONDARY_RISK_NOTICE)
    cp1, cp2 = st.columns(2)
    spouse_acquisition_pct = 0
    if has_spouse:
        spouse_acquisition_pct = cp1.slider("一次相続における配偶者の取得割合(%)", min_value=0, max_value=100, value=50, key="in_spouse_ratio")
    else:
        cp1.caption("配偶者がいないため、配偶者取得割合は0%で固定です。")

    use_individual_allocations = st.checkbox("一次相続の実取得額を相続人ごとに入力する（推奨）", value=True, key="use_individual_allocations")
    labels = build_heir_labels(has_spouse, heirs_info)
    default_amounts = build_default_acquisition_input_amounts(estimated_tax_p, has_spouse, heirs_info)
    actual_acquisition_inputs: list[int] = []

    if use_individual_allocations:
        st.caption(f"参考：一次相続の課税価格合計 {fmt_int(estimated_tax_p)}円 を基準に、各人の実取得額（概算）を入力してください。合計が一致しない場合は内部で比率按分します。")
        for idx, (label, heir_type) in enumerate(labels):
            key = f"actual_acq_{idx}"
            default_val = default_amounts[idx] if idx < len(default_amounts) else 0
            amount = st.number_input(
                f"{label}（{heir_type}）の実取得額",
                min_value=0,
                value=default_val,
                key=key,
            )
            actual_acquisition_inputs.append(int(amount))
        entered_total = sum(actual_acquisition_inputs)
        discrepancy = int(estimated_tax_p) - entered_total
        st.caption(f"入力合計: {entered_total:,}円 / 参考課税価格合計: {int(estimated_tax_p):,}円 / 差額: {discrepancy:,}円")
    else:
        actual_acquisition_inputs = default_amounts

    s_own = cp1.number_input("配偶者の固有財産", min_value=0, value=50000000, key="in_s_own")
    interval_years = cp1.slider("二次までの想定期間(年)", min_value=0, max_value=20, value=10, key="in_interval")
    annual_spend = cp2.number_input("年間生活費・支出(減価)", min_value=0, value=5000000, key="in_s_spend")
    return SecondaryInputs(
        spouse_acquisition_pct=spouse_acquisition_pct,
        s_own=s_own,
        annual_spend=annual_spend,
        interval_years=interval_years,
        use_individual_allocations=use_individual_allocations,
        actual_acquisition_inputs=actual_acquisition_inputs,
    )


def render_tab_analysis(primary_inputs: PrimaryInputs, primary_result: PrimaryResult, secondary_inputs: SecondaryInputs, secondary_result: SecondaryResult, df_sim: pd.DataFrame, iryu_df: pd.DataFrame, df1: pd.DataFrame, df_heirs: pd.DataFrame, df_small: pd.DataFrame, df_gifts: pd.DataFrame, df2: pd.DataFrame, df_snapshot_summary: pd.DataFrame, df_carryforward: pd.DataFrame, df_audit_notes: pd.DataFrame, df_small_scale_review: pd.DataFrame, df_successive_credit: pd.DataFrame) -> None:
    add_print_button("6. 精密分析結果")
    st.subheader("配偶者取得割合別の税額推移分析（内部確認用・概算・横軸=配偶者取得割合 / 縦軸=税額）")
    render_risk_notice("配偶者取得割合別の比較は内部検討用の参考表示です。差額の背景にある個別論点確認前に断定利用しないでください。", level="info")
    st.plotly_chart(build_simulation_figure(df_sim), width="stretch")
    st.dataframe(
        df_sim.style.format({"一次相続税額": "{:,}", "二次相続税額": "{:,}", "合計納税額": "{:,}"}),
        width="stretch",
    )

    st.divider()
    st.subheader("⚠️ 遺留分侵害額の参考表示")
    st.table(iryu_df)

    st.divider()
    render_audit_evidence()

    st.divider()
    st.subheader("📥 ご提案資料の出力")
    render_risk_notice(OUTPUT_RISK_NOTICE)
    col_excel, col_pdf, col_ppt = st.columns(3)
    try:
        excel_data = create_excel_file(primary_inputs, primary_result, secondary_inputs, secondary_result, df1, df_heirs, df_small, df_gifts, df2, df_sim, df_snapshot_summary, df_carryforward, df_audit_notes, df_small_scale_review, df_successive_credit)
        with col_excel:
            st.download_button(
                label="📊 Excelファイルをダウンロード（ご提案資料）",
                data=excel_data,
                file_name=EXCEL_FILE_NAME,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
    except Exception as exc:
        with col_excel:
            st.error(f"Excel出力エラー: {exc}")

    try:
        pdf_data = create_pdf_report(primary_inputs, primary_result, secondary_inputs, secondary_result, df_sim, df_snapshot_summary, df_carryforward, df_audit_notes, df_small_scale_review, df_successive_credit)
        with col_pdf:
            st.download_button(
                label="📄 PDFファイルをダウンロード（ご提案資料）",
                data=pdf_data,
                file_name=PDF_FILE_NAME,
                mime="application/pdf",
            )
    except Exception as exc:
        with col_pdf:
            st.error(f"PDF出力エラー: {exc}")

    try:
        ppt_data = create_ppt_report(primary_inputs, primary_result, secondary_inputs, secondary_result, df_sim, df_snapshot_summary, df_carryforward, df_audit_notes, df_small_scale_review, df_successive_credit)
        with col_ppt:
            st.download_button(
                label="🖥️ PPTファイルをダウンロード（ご提案資料）",
                data=ppt_data,
                file_name=PPT_FILE_NAME,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
    except Exception as exc:
        with col_ppt:
            st.error(f"PPT出力エラー: {exc}")


# =========================================================
# 9. main()
# =========================================================
def main() -> None:
    inject_print_css()
    if not authenticate_user():
        return

    render_sidebar()
    render_risk_notice(GLOBAL_RISK_NOTICE)
    tabs = st.tabs(TAB_LABELS)

    with tabs[0]:
        heir_count, has_spouse, heirs_info = render_tab_basic()

    with tabs[1]:
        primary_inputs = render_tab_primary_inputs(heir_count, has_spouse, heirs_info)

    estimated_tax_p = estimate_total_taxable_price_reference(primary_inputs)

    with tabs[4]:
        secondary_inputs = render_tab_secondary_parameters(has_spouse, heirs_info, estimated_tax_p)

    primary_result = calculate_primary_inheritance(primary_inputs, secondary_inputs)
    secondary_result = calculate_secondary_inheritance(primary_inputs, primary_result, secondary_inputs)
    df1 = build_primary_detail_df(primary_inputs, primary_result)
    df_heirs = build_primary_heir_tax_df(primary_result)
    df_small = build_small_scale_detail_df(primary_result)
    df_gifts = build_gift_detail_df(primary_result)
    df2 = build_secondary_detail_df(secondary_result)
    df_snapshot_summary = build_snapshot_summary_df(secondary_result.snapshot, secondary_result.context, secondary_result)
    df_carryforward = build_heir_carryforward_df(secondary_result.snapshot)
    df_audit_notes = build_secondary_audit_notes_df(secondary_result.snapshot, secondary_result.context, secondary_result)
    df_small_scale_review = build_secondary_small_scale_review_df(secondary_result)
    df_successive_credit = build_successive_inheritance_credit_df(secondary_result)
    df_sim = build_simulation_df(primary_inputs, primary_result, secondary_inputs)
    iryu_df = build_iryubun_reference(primary_inputs, primary_result)

    with tabs[2]:
        render_tab_primary_detail(df1, df_heirs, df_small, df_gifts)

    with tabs[3]:
        render_tab_secondary_detail(df2, df_snapshot_summary, df_carryforward, df_audit_notes, df_small_scale_review, df_successive_credit)

    with tabs[5]:
        render_tab_analysis(primary_inputs, primary_result, secondary_inputs, secondary_result, df_sim, iryu_df, df1, df_heirs, df_small, df_gifts, df2, df_snapshot_summary, df_carryforward, df_audit_notes, df_small_scale_review, df_successive_credit)


if __name__ == "__main__":
    main()
